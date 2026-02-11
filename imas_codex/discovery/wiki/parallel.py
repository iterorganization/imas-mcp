"""
Parallel wiki discovery engine with async workers.

Architecture:
- Two async workers: Scorer, Ingester (scan is bulk upfront)
- Graph + claimed_at timestamp for coordination (same pattern as paths discovery)
- Status transitions:
  - scanned → scored (content-aware scoring in single pass)
  - scored → ingested (embed for score >= 0.5)
- Workers claim pages by setting claimed_at, release by clearing it
- Orphan recovery: pages with claimed_at > 5 min old are reclaimed

Resilience:
- Supervised workers with automatic restart on crash (via base.supervision)
- Exponential backoff on infrastructure errors (Neo4j, network)
- Graceful degradation when services are temporarily unavailable

Workflow (after bulk discovery):
1. SCORE: Fetch content + LLM scoring in single pass → sets score, updates to 'scored'
2. INGEST: Chunk and embed high-value pages (score >= 0.5) → updates to 'ingested'
"""

from __future__ import annotations

import asyncio
import functools
import logging
import random
import subprocess
import time
import urllib.parse
from dataclasses import dataclass, field
from typing import TYPE_CHECKING, Any

from neo4j.exceptions import TransientError

from imas_codex.discovery.base.progress import WorkerStats
from imas_codex.discovery.base.supervision import (
    SupervisedWorkerGroup,
    supervised_worker,
)
from imas_codex.graph import GraphClient
from imas_codex.graph.models import WikiArtifactStatus, WikiPageStatus

if TYPE_CHECKING:
    from collections.abc import Callable

logger = logging.getLogger(__name__)

# Claim timeout - pages claimed longer than this are reclaimed
CLAIM_TIMEOUT_SECONDS = 300  # 5 minutes


# =============================================================================
# Shared Graph Helpers — single source of truth for wiki node creation
# =============================================================================


def _bulk_create_wiki_pages(
    gc: GraphClient,
    facility: str,
    batch_data: list[dict],
    *,
    batch_size: int = 500,
    on_progress: Callable | None = None,
) -> int:
    """Create WikiPage nodes with FACILITY_ID relationship in batches.

    Every WikiPage gets a [:FACILITY_ID]->(:Facility) relationship at creation
    time, not deferred to ingestion.

    Args:
        gc: Open GraphClient
        facility: Facility ID
        batch_data: List of dicts with id, title, url keys
        batch_size: Nodes per batch
        on_progress: Optional progress callback

    Returns:
        Number of pages created/updated
    """
    created = 0
    total = len(batch_data)
    for i in range(0, total, batch_size):
        batch = batch_data[i : i + batch_size]
        result = gc.query(
            """
            UNWIND $pages AS page
            MERGE (wp:WikiPage {id: page.id})
            ON CREATE SET wp.title = page.title,
                          wp.url = page.url,
                          wp.facility_id = $facility,
                          wp.status = $scanned,
                          wp.link_depth = 1,
                          wp.discovered_at = datetime(),
                          wp.bulk_discovered = true
            ON MATCH SET wp.bulk_discovered = true
            WITH wp
            MATCH (f:Facility {id: $facility})
            MERGE (wp)-[:FACILITY_ID]->(f)
            RETURN count(wp) AS count
            """,
            pages=batch,
            facility=facility,
            scanned=WikiPageStatus.scanned.value,
        )
        if result:
            created += result[0]["count"]

        if on_progress:
            on_progress(f"creating pages ({i + len(batch)}/{total})", None)

    return created


def _bulk_create_wiki_artifacts(
    gc: GraphClient,
    facility: str,
    batch_data: list[dict],
    *,
    batch_size: int = 500,
    on_progress: Callable | None = None,
) -> int:
    """Create WikiArtifact nodes with FACILITY_ID and HAS_ARTIFACT relationships.

    Args:
        gc: Open GraphClient
        facility: Facility ID
        batch_data: List of dicts with id, filename, url, artifact_type,
                    and optionally size_bytes, mime_type, linked_pages
        batch_size: Nodes per batch
        on_progress: Optional progress callback

    Returns:
        Number of artifacts created/updated
    """
    created = 0
    total = len(batch_data)
    for i in range(0, total, batch_size):
        batch = batch_data[i : i + batch_size]
        result = gc.query(
            """
            UNWIND $artifacts AS a
            MERGE (wa:WikiArtifact {id: a.id})
            ON CREATE SET wa.facility_id = $facility,
                          wa.filename = a.filename,
                          wa.url = a.url,
                          wa.artifact_type = a.artifact_type,
                          wa.size_bytes = a.size_bytes,
                          wa.mime_type = a.mime_type,
                          wa.status = $discovered,
                          wa.discovered_at = datetime(),
                          wa.bulk_discovered = true
            ON MATCH SET wa.bulk_discovered = true
            WITH wa
            MATCH (f:Facility {id: $facility})
            MERGE (wa)-[:FACILITY_ID]->(f)
            RETURN count(wa) AS count
            """,
            artifacts=batch,
            facility=facility,
            discovered=WikiArtifactStatus.discovered.value,
        )
        if result:
            created += result[0]["count"]

        if on_progress:
            on_progress(f"created {i + len(batch)}/{total} artifacts", None)

    # Create HAS_ARTIFACT relationships from linked pages
    # This is done in a separate pass to handle pages that may not exist yet
    page_links = []
    for a in batch_data:
        artifact_id = a["id"]
        for page_name in a.get("linked_pages", []):
            page_id = f"{facility}:{page_name}"
            page_links.append({"artifact_id": artifact_id, "page_id": page_id})

    if page_links:
        # Link artifacts to existing pages (pages discovered via bulk discovery)
        gc.query(
            """
            UNWIND $links AS link
            MATCH (wa:WikiArtifact {id: link.artifact_id})
            MATCH (wp:WikiPage {id: link.page_id})
            MERGE (wp)-[:HAS_ARTIFACT]->(wa)
            """,
            links=page_links,
        )

    return created


# Retry configuration for Neo4j transient errors (deadlocks)
MAX_RETRY_ATTEMPTS = 5
RETRY_BASE_DELAY = 0.1  # seconds
RETRY_MAX_DELAY = 2.0  # seconds


def retry_on_deadlock(
    max_attempts: int = MAX_RETRY_ATTEMPTS,
    base_delay: float = RETRY_BASE_DELAY,
    max_delay: float = RETRY_MAX_DELAY,
):
    """Decorator to retry functions on Neo4j transient errors (e.g., deadlocks).

    Uses exponential backoff with jitter to reduce contention.

    Args:
        max_attempts: Maximum number of retry attempts
        base_delay: Initial delay in seconds
        max_delay: Maximum delay in seconds
    """

    def decorator(func):
        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            last_exception = None
            for attempt in range(max_attempts):
                try:
                    return func(*args, **kwargs)
                except TransientError as e:
                    last_exception = e
                    if attempt < max_attempts - 1:
                        # Exponential backoff with jitter
                        delay = min(base_delay * (2**attempt), max_delay)
                        jitter = random.uniform(0, delay * 0.5)
                        sleep_time = delay + jitter
                        logger.debug(
                            "%s: transient error (attempt %d/%d), "
                            "retrying in %.2fs: %s",
                            func.__name__,
                            attempt + 1,
                            max_attempts,
                            sleep_time,
                            e,
                        )
                        time.sleep(sleep_time)
                    else:
                        logger.warning(
                            "%s: transient error after %d attempts: %s",
                            func.__name__,
                            max_attempts,
                            e,
                        )
            raise last_exception  # type: ignore[misc]

        return wrapper

    return decorator


# =============================================================================
# Wiki Discovery State
# =============================================================================


@dataclass
class WikiDiscoveryState:
    """Shared state for parallel wiki discovery."""

    facility: str
    site_type: str  # mediawiki, confluence, twiki
    base_url: str
    portal_page: str
    ssh_host: str | None = None

    # Authentication for HTTP-based access
    auth_type: str | None = None  # tequila, session, keycloak, basic, or None
    credential_service: str | None = None  # Keyring service for credentials

    # Shared async wiki client for native async HTTP (avoids re-auth per page)
    # Initialized lazily on first use, shared across all workers
    _async_wiki_client: Any = field(default=None, repr=False)
    _async_wiki_client_lock: Any = field(default=None, repr=False)

    # Shared Keycloak async client (for keycloak auth via oauth2-proxy)
    _keycloak_client: Any = field(default=None, repr=False)
    _keycloak_client_lock: Any = field(default=None, repr=False)

    # Shared HTTP Basic auth async client (for basic auth wikis like JET)
    _basic_auth_client: Any = field(default=None, repr=False)
    _basic_auth_client_lock: Any = field(default=None, repr=False)

    # Shared HTTP fetch semaphore — bounds total concurrent wiki HTTP requests
    # across ALL workers (score + ingest) to prevent httpx PoolTimeout.
    # Default 30 keeps us well under httpx pool max_connections=50.
    _http_fetch_semaphore: Any = field(default=None, repr=False)

    # SSH fetch semaphore — bounds concurrent SSH subprocess calls.
    # SSH multiplexes over ControlMaster, but too many concurrent ssh
    # processes can race during master establishment or create load.
    # Default 4 keeps footprint low on remote hosts.
    _ssh_fetch_semaphore: Any = field(default=None, repr=False)

    # Limits
    cost_limit: float = 10.0
    page_limit: int | None = None
    max_depth: int | None = None
    focus: str | None = None

    # Worker stats (simplified: score + ingest only)
    scan_stats: WorkerStats = field(default_factory=WorkerStats)
    score_stats: WorkerStats = field(default_factory=WorkerStats)
    ingest_stats: WorkerStats = field(default_factory=WorkerStats)
    artifact_stats: WorkerStats = field(default_factory=WorkerStats)
    artifact_score_stats: WorkerStats = field(default_factory=WorkerStats)
    image_stats: WorkerStats = field(default_factory=WorkerStats)

    # Control
    stop_requested: bool = False
    scan_idle_count: int = 0
    score_idle_count: int = 0
    ingest_idle_count: int = 0
    artifact_idle_count: int = 0
    artifact_score_idle_count: int = 0
    image_idle_count: int = 0

    # SSH retry tracking
    ssh_retry_count: int = 0
    max_ssh_retries: int = 5
    ssh_error_message: str | None = None

    @property
    def http_fetch_semaphore(self) -> asyncio.Semaphore:
        """Shared semaphore bounding total concurrent wiki HTTP fetches.

        Lazily initialized because asyncio.Semaphore needs a running event loop.
        Shared across ALL score and ingest workers so total concurrency stays
        within the httpx connection pool capacity.
        """
        if self._http_fetch_semaphore is None:
            self._http_fetch_semaphore = asyncio.Semaphore(30)
        return self._http_fetch_semaphore

    @property
    def ssh_fetch_semaphore(self) -> asyncio.Semaphore:
        """Shared semaphore bounding concurrent SSH subprocess calls.

        SSH ControlMaster multiplexes over a single TCP connection, but
        too many concurrent ssh processes can race during establishment
        or create unnecessary load on remote hosts. Limit to 4.
        """
        if self._ssh_fetch_semaphore is None:
            self._ssh_fetch_semaphore = asyncio.Semaphore(4)
        return self._ssh_fetch_semaphore

    @property
    def effective_fetch_semaphore(self) -> asyncio.Semaphore:
        """Return SSH semaphore for SSH-based sites, HTTP semaphore otherwise."""
        if self.ssh_host and self.auth_type in (None, "none"):
            return self.ssh_fetch_semaphore
        return self.http_fetch_semaphore

    @property
    def total_cost(self) -> float:
        return self.score_stats.cost + self.ingest_stats.cost + self.image_stats.cost

    @property
    def budget_exhausted(self) -> bool:
        return self.total_cost >= self.cost_limit

    @property
    def page_limit_reached(self) -> bool:
        if self.page_limit is None:
            return False
        return self.score_stats.processed >= self.page_limit

    def should_stop(self) -> bool:
        """Check if ALL workers should terminate.

        Used by the main loop to determine when discovery is complete.
        """
        if self.stop_requested:
            return True
        # Stop if all workers idle for 3+ iterations AND no pending work
        all_idle = (
            self.scan_idle_count >= 3
            and self.score_idle_count >= 3
            and self.ingest_idle_count >= 3
            and self.artifact_idle_count >= 3
            and self.artifact_score_idle_count >= 3
            and self.image_idle_count >= 3
        )
        if all_idle:
            if (
                has_pending_work(self.facility)
                or has_pending_artifact_work(self.facility)
                or has_pending_image_work(self.facility)
            ):
                # Reset idle counts to force workers to re-poll
                self.scan_idle_count = 0
                self.score_idle_count = 0
                self.ingest_idle_count = 0
                self.artifact_idle_count = 0
                self.artifact_score_idle_count = 0
                self.image_idle_count = 0
                return False
            return True
        return False

    def should_stop_scanning(self) -> bool:
        """Check if scan workers should stop.

        Scan workers continue even when budget is exhausted. They only stop
        when explicitly requested or when no pending work remains.
        """
        if self.stop_requested:
            return True
        # Only stop scanning when idle with no work
        if self.scan_idle_count >= 3 and not has_pending_scan_work(self.facility):
            return True
        return False

    def should_stop_scoring(self) -> bool:
        """Check if score workers should stop.

        Score workers stop when budget exhausted or page limit reached.
        """
        if self.stop_requested:
            return True
        if self.budget_exhausted:
            return True
        if self.page_limit_reached:
            return True
        return False

    def should_stop_ingesting(self) -> bool:
        """Check if ingest workers should stop.

        Ingest workers continue AFTER budget is exhausted to drain the
        ingest queue. This ensures all scorable content gets ingested
        before termination. They only stop when:
        1. Explicitly requested
        2. Idle for 3+ iterations with no pending ingest work AND
           score workers are also idle (no more scoring happening)
        """
        if self.stop_requested:
            return True
        # Continue even when budget exhausted - drain the ingest queue
        # BUT don't exit early if scoring is still running - pages may arrive soon
        if self.ingest_idle_count >= 3:
            # Only stop if scoring is also idle AND no pending ingest work
            scoring_done = self.score_idle_count >= 3 or self.budget_exhausted
            if scoring_done and not has_pending_ingest_work(self.facility):
                return True
        return False

    def should_stop_artifact_worker(self) -> bool:
        """Check if artifact ingest workers should stop.

        Artifact ingest workers continue until no pending scored artifacts remain.
        They stop when explicitly requested or idle for 3+ iterations.
        """
        if self.stop_requested:
            return True
        if self.artifact_idle_count >= 3 and not has_pending_artifact_ingest_work(
            self.facility
        ):
            return True
        return False

    def should_stop_artifact_scoring(self) -> bool:
        """Check if artifact score workers should stop.

        Artifact score workers stop when budget exhausted or no more discovered artifacts.
        """
        if self.stop_requested:
            return True
        if self.budget_exhausted:
            return True
        if self.artifact_score_idle_count >= 3 and not has_pending_artifact_score_work(
            self.facility
        ):
            return True
        return False

    def should_stop_image_scoring(self) -> bool:
        """Check if image score workers should stop.

        Image score workers stop when budget exhausted or no pending images.
        They continue after page ingestion completes since page ingestion
        creates the Image nodes.
        """
        if self.stop_requested:
            return True
        if self.budget_exhausted:
            return True
        # Only stop if ingestion is also done AND no pending images
        ingestion_done = self.ingest_idle_count >= 3
        if (
            self.image_idle_count >= 3
            and ingestion_done
            and not has_pending_image_work(self.facility)
        ):
            return True
        return False

    async def get_async_wiki_client(self):
        """Get shared AsyncMediaWikiClient for Tequila auth.

        Lazily initializes an async client that persists session cookies
        across all page fetches. Uses native async HTTP for better performance.
        """
        if self._async_wiki_client is None:
            from imas_codex.discovery.wiki.mediawiki import AsyncMediaWikiClient

            if self._async_wiki_client_lock is None:
                self._async_wiki_client_lock = asyncio.Lock()

            async with self._async_wiki_client_lock:
                if self._async_wiki_client is None:
                    self._async_wiki_client = AsyncMediaWikiClient(
                        base_url=self.base_url,
                        credential_service=self.credential_service
                        or f"{self.facility}-wiki",
                        verify_ssl=False,
                    )
                    # Pre-authenticate to warm up session
                    try:
                        await self._async_wiki_client.authenticate()
                        logger.info(
                            "Initialized shared AsyncMediaWikiClient with Tequila auth"
                        )
                    except Exception as e:
                        logger.warning(
                            "Failed to pre-authenticate AsyncMediaWikiClient: %s", e
                        )

        return self._async_wiki_client

    async def close_async_wiki_client(self):
        """Close the shared async wiki client."""
        if self._async_wiki_client is not None:
            try:
                await self._async_wiki_client.close()
            except Exception:
                pass
            self._async_wiki_client = None

    async def get_keycloak_client(self) -> Any:
        """Get shared AsyncKeycloakSession for Keycloak OIDC auth.

        Lazily initializes a Keycloak OIDC session that persists cookies
        across all page fetches. One login authenticates across all wiki
        sites on the same domain.

        Returns:
            Authenticated httpx.AsyncClient or None if auth fails.
        """
        if self._keycloak_client is None:
            from imas_codex.discovery.wiki.keycloak import AsyncKeycloakSession

            if self._keycloak_client_lock is None:
                self._keycloak_client_lock = asyncio.Lock()

            async with self._keycloak_client_lock:
                if self._keycloak_client is None:
                    aks = AsyncKeycloakSession(self.credential_service or self.facility)
                    try:
                        self._keycloak_client = await aks.login(f"{self.base_url}/")
                        logger.info(
                            "Initialized shared Keycloak session for %s",
                            self.credential_service,
                        )
                    except Exception as e:
                        logger.warning("Keycloak auth failed: %s", e)

        return self._keycloak_client

    async def close_keycloak_client(self):
        """Close the shared Keycloak client."""
        if self._keycloak_client is not None:
            try:
                await self._keycloak_client.aclose()
            except Exception:
                pass
            self._keycloak_client = None

    async def get_basic_auth_client(self) -> Any:
        """Get shared httpx.AsyncClient with HTTP Basic auth.

        For wikis that use standard HTTP Basic authentication (RFC 7617),
        e.g. JET wiki sites at wiki.jetdata.eu.

        Returns:
            httpx.AsyncClient with BasicAuth or None if credentials unavailable.
        """
        if self._basic_auth_client is None:
            import httpx

            from imas_codex.discovery.wiki.auth import CredentialManager

            if self._basic_auth_client_lock is None:
                self._basic_auth_client_lock = asyncio.Lock()

            async with self._basic_auth_client_lock:
                if self._basic_auth_client is None:
                    cred_mgr = CredentialManager()
                    service = self.credential_service or self.facility
                    creds = cred_mgr.get_credentials(service, prompt_if_missing=False)
                    if not creds:
                        logger.warning(
                            "No credentials for %s - cannot use HTTP Basic auth",
                            service,
                        )
                        return None
                    username, password = creds
                    self._basic_auth_client = httpx.AsyncClient(
                        auth=httpx.BasicAuth(username, password),
                        timeout=httpx.Timeout(60.0),
                        follow_redirects=True,
                        verify=False,
                        headers={
                            "User-Agent": "imas-codex/1.0 (IMAS Data Mapping Tool)",
                            "Accept": "text/html,application/xhtml+xml,*/*;q=0.8",
                        },
                        limits=httpx.Limits(
                            max_connections=20,
                            max_keepalive_connections=10,
                        ),
                    )
                    logger.info(
                        "Initialized shared HTTP Basic auth client for %s",
                        service,
                    )

        return self._basic_auth_client

    async def close_basic_auth_client(self):
        """Close the shared HTTP Basic auth client."""
        if self._basic_auth_client is not None:
            try:
                await self._basic_auth_client.aclose()
            except Exception:
                pass
            self._basic_auth_client = None


# Artifact types we can extract text from
SUPPORTED_ARTIFACT_TYPES = {"pdf", "docx", "doc", "pptx", "ppt", "xlsx", "xls", "ipynb"}


def has_pending_artifact_work(facility: str) -> bool:
    """Check if there are pending artifacts (scoring or ingestion).

    Artifacts are pending when:
    - status = 'discovered' AND artifact_type in SUPPORTED_ARTIFACT_TYPES (needs scoring)
    - status = 'scored' AND score >= 0.5 (needs ingestion)
    - claimed_at is null or expired
    """
    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility})
            WHERE (
                (wa.status = $discovered AND wa.artifact_type IN $types)
                OR (wa.status = $scored AND wa.score >= 0.5 AND wa.artifact_type IN $types)
              )
              AND (wa.claimed_at IS NULL OR wa.claimed_at < datetime() - duration($cutoff))
            RETURN count(wa) AS pending
            """,
            facility=facility,
            discovered=WikiArtifactStatus.discovered.value,
            scored=WikiArtifactStatus.scored.value,
            types=list(SUPPORTED_ARTIFACT_TYPES),
            cutoff=cutoff,
        )
        return result[0]["pending"] > 0 if result else False


def has_pending_artifact_score_work(facility: str) -> bool:
    """Check if there are discovered artifacts awaiting scoring.

    Returns True if there are artifacts with:
    - status = 'discovered'
    - artifact_type in SUPPORTED_ARTIFACT_TYPES
    - claimed_at is null or expired
    """
    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility})
            WHERE wa.status = $discovered
              AND wa.artifact_type IN $types
              AND (wa.claimed_at IS NULL OR wa.claimed_at < datetime() - duration($cutoff))
            RETURN count(wa) AS pending
            """,
            facility=facility,
            discovered=WikiArtifactStatus.discovered.value,
            types=list(SUPPORTED_ARTIFACT_TYPES),
            cutoff=cutoff,
        )
        return result[0]["pending"] > 0 if result else False


def has_pending_artifact_ingest_work(facility: str) -> bool:
    """Check if there are scored artifacts awaiting ingestion.

    Returns True if there are artifacts with:
    - status = 'scored'
    - score >= 0.5
    - artifact_type in SUPPORTED_ARTIFACT_TYPES
    - claimed_at is null or expired
    """
    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility})
            WHERE wa.status = $scored
              AND wa.score >= 0.5
              AND wa.artifact_type IN $types
              AND (wa.claimed_at IS NULL OR wa.claimed_at < datetime() - duration($cutoff))
            RETURN count(wa) AS pending
            """,
            facility=facility,
            scored=WikiArtifactStatus.scored.value,
            types=list(SUPPORTED_ARTIFACT_TYPES),
            cutoff=cutoff,
        )
        return result[0]["pending"] > 0 if result else False


def has_pending_work(facility: str) -> bool:
    """Check if there's pending wiki work in the graph.

    Work exists if there are:
    - scanned pages awaiting scoring (claimed_at is null)
    - scored pages with score >= 0.5 awaiting ingest (claimed_at is null)
    """
    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE (wp.status = $scanned AND (wp.claimed_at IS NULL OR wp.claimed_at < datetime() - duration($cutoff)))
               OR (wp.status = $scored AND wp.score >= 0.5 AND (wp.claimed_at IS NULL OR wp.claimed_at < datetime() - duration($cutoff)))
            RETURN count(wp) AS pending
            """,
            facility=facility,
            scanned=WikiPageStatus.scanned.value,
            scored=WikiPageStatus.scored.value,
            cutoff=cutoff,
        )
        return result[0]["pending"] > 0 if result else False


def has_pending_scan_work(facility: str) -> bool:
    """Check if there's pending scoring work in the graph.

    With bulk discovery, there's no scan phase. This checks for:
    - scanned pages awaiting content-aware scoring
    """
    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE wp.status = $scanned
              AND (wp.claimed_at IS NULL OR wp.claimed_at < datetime() - duration($cutoff))
            RETURN count(wp) AS pending
            """,
            facility=facility,
            scanned=WikiPageStatus.scanned.value,
            cutoff=cutoff,
        )
        return result[0]["pending"] > 0 if result else False


def has_pending_ingest_work(facility: str) -> bool:
    """Check if there's pending ingest work in the graph.

    Returns True if there are scored pages with score >= 0.5 awaiting ingestion.
    """
    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE wp.status = $scored
              AND wp.score >= 0.5
              AND (wp.claimed_at IS NULL OR wp.claimed_at < datetime() - duration($cutoff))
            RETURN count(wp) AS pending
            """,
            facility=facility,
            scored=WikiPageStatus.scored.value,
            cutoff=cutoff,
        )
        return result[0]["pending"] > 0 if result else False


# =============================================================================
# Startup Reset
# =============================================================================


def reset_transient_pages(facility: str, *, silent: bool = False) -> dict[str, int]:
    """Reset claimed_at on all wiki pages and artifacts on CLI startup.

    Since only one CLI process runs per facility at a time, any items with
    claimed_at set are orphans from a previous crashed/killed process.
    Simply clear claimed_at to make them reclaimable.
    """
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE wp.claimed_at IS NOT NULL
            SET wp.claimed_at = null
            RETURN count(wp) AS reset_count
            """,
            facility=facility,
        )
        page_reset = result[0]["reset_count"] if result else 0

        # Also clear artifact claims
        result = gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility})
            WHERE wa.claimed_at IS NOT NULL
            SET wa.claimed_at = null
            RETURN count(wa) AS reset_count
            """,
            facility=facility,
        )
        artifact_reset = result[0]["reset_count"] if result else 0

    total_reset = page_reset + artifact_reset
    if not silent and total_reset > 0:
        logger.info(
            "Reset %d orphaned claims on startup (%d pages, %d artifacts)",
            total_reset,
            page_reset,
            artifact_reset,
        )

    return {"orphan_reset": page_reset, "artifact_reset": artifact_reset}


# =============================================================================
# Graph-based Work Claiming (uses claimed_at for coordination)
# =============================================================================


@retry_on_deadlock()
def claim_pages_for_scoring(facility: str, limit: int = 50) -> list[dict[str, Any]]:
    """Claim scanned pages for content-aware scoring.

    Workflow: scanned + unclaimed → set claimed_at
    Score worker fetches content and scores in single pass.
    After scoring: update status to 'scored' and set score field.

    Uses claim token pattern to handle race conditions between workers:
    1. Generate unique claim token
    2. Atomically SET token on unclaimed pages
    3. Read back only pages with OUR token (pages we actually won)
    """
    import uuid

    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"  # ISO 8601 duration
    claim_token = str(uuid.uuid4())

    with GraphClient() as gc:
        # Step 1: Attempt to claim pages with our unique token
        # Using random() in ORDER BY reduces collision probability further
        gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE wp.status = $scanned
              AND (wp.claimed_at IS NULL
                   OR wp.claimed_at < datetime() - duration($cutoff))
            WITH wp
            ORDER BY rand()
            LIMIT $limit
            SET wp.claimed_at = datetime(), wp.claim_token = $token
            """,
            facility=facility,
            scanned=WikiPageStatus.scanned.value,
            cutoff=cutoff,
            limit=limit,
            token=claim_token,
        )

        # Step 2: Read back only pages WE successfully claimed
        # If another worker raced us, they have a different token
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility, claim_token: $token})
            RETURN wp.id AS id, wp.title AS title, wp.url AS url
            """,
            facility=facility,
            token=claim_token,
        )
        claimed = list(result)

        logger.debug(
            "claim_pages_for_scoring: requested %d, won %d (token=%s)",
            limit,
            len(claimed),
            claim_token[:8],
        )
        return claimed


@retry_on_deadlock()
def claim_pages_for_ingesting(
    facility: str, min_score: float = 0.5, limit: int = 10
) -> list[dict[str, Any]]:
    """Claim scored pages for ingestion (chunking and embedding).

    Workflow: scored + score >= min_score + unclaimed → set claimed_at
    After ingest: update status to 'ingested'.

    Uses claim token pattern to handle race conditions between workers.
    """
    import uuid

    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    claim_token = str(uuid.uuid4())

    with GraphClient() as gc:
        # Step 1: Attempt to claim pages with our unique token
        gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE wp.status = $scored
              AND wp.score >= $min_score
              AND (wp.claimed_at IS NULL
                   OR wp.claimed_at < datetime() - duration($cutoff))
            WITH wp
            ORDER BY rand()
            LIMIT $limit
            SET wp.claimed_at = datetime(), wp.claim_token = $token
            """,
            facility=facility,
            scored=WikiPageStatus.scored.value,
            min_score=min_score,
            cutoff=cutoff,
            limit=limit,
            token=claim_token,
        )

        # Step 2: Read back only pages WE successfully claimed
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility, claim_token: $token})
            RETURN wp.id AS id, wp.title AS title, wp.url AS url,
                   wp.score AS score, wp.description AS description,
                   wp.physics_domain AS physics_domain,
                   wp.preview_text AS preview,
                   wp.in_degree AS in_degree, wp.out_degree AS out_degree
            """,
            facility=facility,
            token=claim_token,
        )
        claimed = list(result)

        logger.debug(
            "claim_pages_for_ingesting: requested %d, won %d (token=%s)",
            limit,
            len(claimed),
            claim_token[:8],
        )
        return claimed


# =============================================================================
# Mark Work Complete
# =============================================================================


def mark_pages_scored(
    facility: str,
    results: list[dict[str, Any]],
    skip_threshold: float = 0.5,
) -> int:
    """Mark pages as scored or skipped based on score threshold.

    Pages with score >= skip_threshold get status='scored' (proceed to ingest).
    Pages with score < skip_threshold get status='skipped' (filtered out).

    Uses batched UNWIND for O(1) graph operations instead of O(n) individual queries.
    """
    if not results:
        return 0

    # Prepare batch data with all scoring fields, split by threshold
    scored_batch: list[dict[str, Any]] = []
    skipped_batch: list[dict[str, Any]] = []
    # Safety: collect page IDs with empty preview_text to release for reprocessing
    released_ids: list[str] = []

    for r in results:
        page_id = r.get("id")
        if not page_id:
            continue

        # Pages without preview_text should not be marked as scored —
        # release them back to scanned so they can be retried
        if not r.get("preview_text"):
            released_ids.append(page_id)
            continue

        item = {
            "id": page_id,
            "score": r.get("score", 0.0),
            "purpose": r.get("purpose", r.get("page_purpose", "other")),
            "description": (r.get("description", "") or "")[:150],
            "reasoning": r.get("reasoning", ""),
            "keywords": r.get("keywords", []),
            "physics_domain": r.get("physics_domain"),
            "preview_text": r.get("preview_text", ""),
            "score_data_documentation": r.get("score_data_documentation", 0.0),
            "score_physics_content": r.get("score_physics_content", 0.0),
            "score_code_documentation": r.get("score_code_documentation", 0.0),
            "score_data_access": r.get("score_data_access", 0.0),
            "score_calibration": r.get("score_calibration", 0.0),
            "score_imas_relevance": r.get("score_imas_relevance", 0.0),
            "should_ingest": r.get("should_ingest", True),
            "skip_reason": r.get("skip_reason"),
            "is_physics": r.get("is_physics", False),
            "score_cost": r.get("score_cost", 0.0),
        }

        if item["score"] >= skip_threshold:
            scored_batch.append(item)
        else:
            skipped_batch.append(item)

    if not scored_batch and not skipped_batch and not released_ids:
        return 0

    with GraphClient() as gc:
        # Cypher template shared by scored and skipped batches
        _SET_SCORING_FIELDS = """
                UNWIND $batch AS item
                MATCH (wp:WikiPage {id: item.id})
                SET wp.status = $status,
                    wp.score = item.score,
                    wp.purpose = item.purpose,
                    wp.description = item.description,
                    wp.reasoning = item.reasoning,
                    wp.keywords = item.keywords,
                    wp.physics_domain = item.physics_domain,
                    wp.preview_text = item.preview_text,
                    wp.score_data_documentation = item.score_data_documentation,
                    wp.score_physics_content = item.score_physics_content,
                    wp.score_code_documentation = item.score_code_documentation,
                    wp.score_data_access = item.score_data_access,
                    wp.score_calibration = item.score_calibration,
                    wp.score_imas_relevance = item.score_imas_relevance,
                    wp.should_ingest = item.should_ingest,
                    wp.skip_reason = item.skip_reason,
                    wp.is_physics_content = item.is_physics,
                    wp.score_cost = item.score_cost,
                    wp.scored_at = datetime(),
                    wp.preview_fetched_at = datetime(),
                    wp.claimed_at = null
        """

        if scored_batch:
            gc.query(
                _SET_SCORING_FIELDS,
                batch=scored_batch,
                status=WikiPageStatus.scored.value,
            )

        if skipped_batch:
            gc.query(
                _SET_SCORING_FIELDS,
                batch=skipped_batch,
                status=WikiPageStatus.skipped.value,
            )
            logger.info(
                "mark_pages_scored: %d pages skipped (score < %.1f)",
                len(skipped_batch),
                skip_threshold,
            )

        # Release pages with empty preview_text back to scanned for reprocessing
        if released_ids:
            logger.info(
                "mark_pages_scored: releasing %d pages with empty preview_text "
                "back to scanned for retry",
                len(released_ids),
            )
            gc.query(
                """
                UNWIND $ids AS id
                MATCH (wp:WikiPage {id: id})
                SET wp.claimed_at = null
                """,
                ids=released_ids,
            )

    return len(scored_batch) + len(skipped_batch)


def mark_pages_ingested(
    facility: str,
    results: list[dict[str, Any]],
) -> int:
    """Mark pages as ingested with chunk data.

    Uses batched UNWIND for O(1) graph operations instead of O(n) individual queries.
    """
    if not results:
        return 0

    # Prepare batch data
    batch_data = [
        {"id": r.get("id"), "chunks": r.get("chunk_count", 0)}
        for r in results
        if r.get("id")
    ]

    if not batch_data:
        return 0

    with GraphClient() as gc:
        gc.query(
            """
            UNWIND $batch AS item
            MATCH (wp:WikiPage {id: item.id})
            SET wp.status = $ingested,
                wp.chunk_count = item.chunks,
                wp.ingested_at = datetime(),
                wp.claimed_at = null
            """,
            batch=batch_data,
            ingested=WikiPageStatus.ingested.value,
        )

    return len(batch_data)


def mark_page_failed(page_id: str, error: str, fallback_status: str) -> None:
    """Mark a page as failed with error message.

    If Neo4j is unavailable, logs a warning and silently fails.
    The page will be reclaimed after CLAIM_TIMEOUT_SECONDS.
    """
    try:
        with GraphClient() as gc:
            gc.query(
                """
                MATCH (wp:WikiPage {id: $id})
                SET wp.status = $status,
                    wp.error = $error,
                    wp.failed_at = datetime(),
                    wp.claimed_at = null
                """,
                id=page_id,
                status=fallback_status,
                error=error,
            )
    except Exception as e:
        logger.warning(
            "Could not mark page %s as failed (Neo4j unavailable): %s", page_id, e
        )


def _release_claimed_pages(page_ids: list[str]) -> None:
    """Release claimed pages back for reprocessing.

    Clears claimed_at without changing status, allowing pages to be
    reclaimed by the next worker iteration.

    If Neo4j is unavailable, logs a warning and silently fails.
    Pages will be reclaimed after CLAIM_TIMEOUT_SECONDS anyway.
    """
    if not page_ids:
        return

    try:
        with GraphClient() as gc:
            gc.query(
                """
                UNWIND $ids AS id
                MATCH (wp:WikiPage {id: id})
                SET wp.claimed_at = null
                """,
                ids=page_ids,
            )
    except Exception as e:
        logger.warning(
            "Could not release %d claimed pages (Neo4j unavailable): %s",
            len(page_ids),
            e,
        )


def release_orphaned_claims(facility: str) -> dict[str, int]:
    """Release all orphaned claims for a facility.

    Orphaned claims occur when a worker crashes without releasing its claimed
    pages. This function finds all pages with claimed_at older than the timeout
    and clears the claim.

    This is automatically called by the claim functions (they skip old claims),
    but can also be called explicitly to recover from stuck state.

    Args:
        facility: Facility ID

    Returns:
        Dict with counts: {"released": N, "pages": [...]}
    """
    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    try:
        with GraphClient() as gc:
            # Release stale page claims
            result = gc.query(
                """
                MATCH (wp:WikiPage {facility_id: $facility})
                WHERE wp.claimed_at IS NOT NULL
                  AND wp.claimed_at < datetime() - duration($cutoff)
                SET wp.claimed_at = null
                RETURN wp.id AS id, wp.status AS status
                """,
                facility=facility,
                cutoff=cutoff,
            )
            pages = list(result)

            # Release stale artifact claims
            result = gc.query(
                """
                MATCH (wa:WikiArtifact {facility_id: $facility})
                WHERE wa.claimed_at IS NOT NULL
                  AND wa.claimed_at < datetime() - duration($cutoff)
                SET wa.claimed_at = null
                RETURN wa.id AS id, wa.status AS status
                """,
                facility=facility,
                cutoff=cutoff,
            )
            artifacts = list(result)

            total = len(pages) + len(artifacts)
            if total > 0:
                logger.info(
                    "Released %d orphaned claims (%d pages, %d artifacts) for %s",
                    total,
                    len(pages),
                    len(artifacts),
                    facility,
                )

            return {
                "released_pages": len(pages),
                "released_artifacts": len(artifacts),
                "page_ids": [p["id"] for p in pages],
                "artifact_ids": [a["id"] for a in artifacts],
            }
    except Exception as e:
        logger.warning("Could not release orphaned claims: %s", e)
        return {"released_pages": 0, "released_artifacts": 0, "error": str(e)}


# =============================================================================
# Artifact Claim/Mark Functions
# =============================================================================


@retry_on_deadlock()
def claim_artifacts_for_scoring(
    facility: str,
    limit: int = 20,
) -> list[dict[str, Any]]:
    """Claim discovered artifacts for scoring.

    Claims artifacts with status='discovered' and supported artifact_type.
    Supported types: pdf, docx, pptx, xlsx, ipynb.

    Workflow: discovered + supported_type + unclaimed → set claimed_at
    Score worker extracts text preview and scores with LLM.
    After scoring: update status to 'scored' and set score/description/physics_domain.

    Uses claim token pattern to handle race conditions between workers.
    """
    import uuid

    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    claim_token = str(uuid.uuid4())

    with GraphClient() as gc:
        # Step 1: Attempt to claim artifacts with our unique token
        gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility})
            WHERE wa.status = $discovered
              AND wa.artifact_type IN $types
              AND (wa.claimed_at IS NULL
                   OR wa.claimed_at < datetime() - duration($cutoff))
            WITH wa
            ORDER BY rand()
            LIMIT $limit
            SET wa.claimed_at = datetime(), wa.claim_token = $token
            """,
            facility=facility,
            discovered=WikiArtifactStatus.discovered.value,
            types=list(SUPPORTED_ARTIFACT_TYPES),
            cutoff=cutoff,
            limit=limit,
            token=claim_token,
        )

        # Step 2: Read back only artifacts WE successfully claimed
        result = gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility, claim_token: $token})
            RETURN wa.id AS id, wa.url AS url, wa.filename AS filename,
                   wa.artifact_type AS artifact_type, wa.size_bytes AS size_bytes
            """,
            facility=facility,
            token=claim_token,
        )
        claimed = list(result)

        logger.debug(
            "claim_artifacts_for_scoring: requested %d, won %d (token=%s)",
            limit,
            len(claimed),
            claim_token[:8],
        )
        return claimed


@retry_on_deadlock()
def claim_artifacts_for_ingesting(
    facility: str,
    min_score: float = 0.5,
    limit: int = 5,
) -> list[dict[str, Any]]:
    """Claim scored artifacts for ingestion.

    Claims artifacts with status='scored' and score >= min_score.
    Supported types: pdf, docx, pptx, xlsx, ipynb.

    Workflow: scored + score >= threshold + unclaimed → set claimed_at
    After ingest: update status to 'ingested'.

    Uses claim token pattern to handle race conditions between workers.
    """
    import uuid

    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    claim_token = str(uuid.uuid4())

    with GraphClient() as gc:
        # Step 1: Attempt to claim artifacts with our unique token
        gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility})
            WHERE wa.status = $scored
              AND wa.score >= $min_score
              AND wa.artifact_type IN $types
              AND (wa.claimed_at IS NULL
                   OR wa.claimed_at < datetime() - duration($cutoff))
            WITH wa
            ORDER BY wa.score DESC, rand()
            LIMIT $limit
            SET wa.claimed_at = datetime(), wa.claim_token = $token
            """,
            facility=facility,
            scored=WikiArtifactStatus.scored.value,
            min_score=min_score,
            types=list(SUPPORTED_ARTIFACT_TYPES),
            cutoff=cutoff,
            limit=limit,
            token=claim_token,
        )

        # Step 2: Read back only artifacts WE successfully claimed
        result = gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility, claim_token: $token})
            RETURN wa.id AS id, wa.url AS url, wa.filename AS filename,
                   wa.artifact_type AS artifact_type, wa.score AS score,
                   wa.description AS description, wa.physics_domain AS physics_domain
            """,
            facility=facility,
            token=claim_token,
        )
        claimed = list(result)

        logger.debug(
            "claim_artifacts_for_ingesting: requested %d, won %d (token=%s)",
            limit,
            len(claimed),
            claim_token[:8],
        )
        return claimed


def mark_artifacts_scored(
    facility: str,
    results: list[dict[str, Any]],
) -> int:
    """Mark artifacts as scored with scoring data.

    Uses batched UNWIND for O(1) graph operations.
    All artifacts move to 'scored' status. The ingester filters by score >= threshold.
    """
    if not results:
        return 0

    batch_data = []
    for r in results:
        artifact_id = r.get("id")
        if not artifact_id:
            continue

        batch_data.append(
            {
                "id": artifact_id,
                "score": r.get("score", 0.0),
                "artifact_purpose": r.get("artifact_purpose", "other"),
                "description": (r.get("description", "") or "")[:150],
                "reasoning": r.get("reasoning", ""),
                "keywords": r.get("keywords", []),
                "physics_domain": r.get("physics_domain"),
                "preview_text": r.get("preview_text", "")[:500],
                "score_data_documentation": r.get("score_data_documentation", 0.0),
                "score_physics_content": r.get("score_physics_content", 0.0),
                "score_code_documentation": r.get("score_code_documentation", 0.0),
                "score_data_access": r.get("score_data_access", 0.0),
                "score_calibration": r.get("score_calibration", 0.0),
                "score_imas_relevance": r.get("score_imas_relevance", 0.0),
                "should_ingest": r.get("should_ingest", False),
                "skip_reason": r.get("skip_reason"),
                "score_cost": r.get("score_cost", 0.0),
            }
        )

    if not batch_data:
        return 0

    with GraphClient() as gc:
        gc.query(
            """
            UNWIND $batch AS item
            MATCH (wa:WikiArtifact {id: item.id})
            SET wa.status = $status,
                wa.score = item.score,
                wa.artifact_purpose = item.artifact_purpose,
                wa.description = item.description,
                wa.score_reasoning = item.reasoning,
                wa.keywords = item.keywords,
                wa.physics_domain = item.physics_domain,
                wa.preview_text = item.preview_text,
                wa.score_data_documentation = item.score_data_documentation,
                wa.score_physics_content = item.score_physics_content,
                wa.score_code_documentation = item.score_code_documentation,
                wa.score_data_access = item.score_data_access,
                wa.score_calibration = item.score_calibration,
                wa.score_imas_relevance = item.score_imas_relevance,
                wa.should_ingest = item.should_ingest,
                wa.skip_reason = item.skip_reason,
                wa.score_cost = item.score_cost,
                wa.scored_at = datetime(),
                wa.claimed_at = null
            """,
            batch=batch_data,
            status=WikiArtifactStatus.scored.value,
        )

    return len(batch_data)


def mark_artifacts_ingested(
    facility: str,
    results: list[dict[str, Any]],
) -> int:
    """Mark artifacts as ingested with chunk data.

    Uses batched UNWIND for O(1) graph operations.
    """
    if not results:
        return 0

    batch_data = [
        {"id": r.get("id"), "chunks": r.get("chunk_count", 0)}
        for r in results
        if r.get("id")
    ]

    if not batch_data:
        return 0

    with GraphClient() as gc:
        gc.query(
            """
            UNWIND $batch AS item
            MATCH (wa:WikiArtifact {id: item.id})
            SET wa.status = $ingested,
                wa.chunk_count = item.chunks,
                wa.ingested_at = datetime(),
                wa.claimed_at = null
            """,
            batch=batch_data,
            ingested=WikiArtifactStatus.ingested.value,
        )

    return len(batch_data)


def mark_artifact_failed(artifact_id: str, error: str) -> None:
    """Mark an artifact as failed with error message."""
    try:
        with GraphClient() as gc:
            gc.query(
                """
                MATCH (wa:WikiArtifact {id: $id})
                SET wa.status = $failed,
                    wa.error = $error,
                    wa.failed_at = datetime(),
                    wa.claimed_at = null
                """,
                id=artifact_id,
                failed=WikiArtifactStatus.failed.value,
                error=error,
            )
    except Exception as e:
        logger.warning(
            "Could not mark artifact %s as failed (Neo4j unavailable): %s",
            artifact_id,
            e,
        )


def mark_artifact_deferred(artifact_id: str, reason: str) -> None:
    """Mark an artifact as deferred (unsupported type or too large)."""
    try:
        with GraphClient() as gc:
            gc.query(
                """
                MATCH (wa:WikiArtifact {id: $id})
                SET wa.status = $deferred,
                    wa.defer_reason = $reason,
                    wa.claimed_at = null
                """,
                id=artifact_id,
                deferred=WikiArtifactStatus.deferred.value,
                reason=reason,
            )
    except Exception as e:
        logger.warning(
            "Could not mark artifact %s as deferred (Neo4j unavailable): %s",
            artifact_id,
            e,
        )


# =============================================================================
# Image Claim/Mark Functions
# =============================================================================


def has_pending_image_work(facility: str) -> bool:
    """Check if there are images pending scoring (status=ingested, no caption)."""
    try:
        with GraphClient() as gc:
            result = gc.query(
                """
                MATCH (img:Image {facility_id: $facility})
                WHERE img.status = 'ingested'
                  AND img.caption IS NULL
                RETURN count(img) > 0 AS has_work
                """,
                facility=facility,
            )
            rows = list(result)
            return rows[0]["has_work"] if rows else False
    except Exception:
        return False


@retry_on_deadlock()
def claim_images_for_scoring(
    facility: str,
    limit: int = 10,
) -> list[dict[str, Any]]:
    """Claim ingested images for VLM scoring.

    Images with status='ingested' and no caption are ready for VLM processing.
    Uses same claim token pattern as page/artifact claiming.
    """
    import uuid

    cutoff = f"PT{CLAIM_TIMEOUT_SECONDS}S"
    claim_token = str(uuid.uuid4())

    with GraphClient() as gc:
        gc.query(
            """
            MATCH (img:Image {facility_id: $facility})
            WHERE img.status = 'ingested'
              AND img.caption IS NULL
              AND (img.claimed_at IS NULL
                   OR img.claimed_at < datetime() - duration($cutoff))
            WITH img
            ORDER BY rand()
            LIMIT $limit
            SET img.claimed_at = datetime(), img.claim_token = $token
            """,
            facility=facility,
            cutoff=cutoff,
            limit=limit,
            token=claim_token,
        )

        result = gc.query(
            """
            MATCH (img:Image {facility_id: $facility, claim_token: $token})
            RETURN img.id AS id,
                   img.source_url AS source_url,
                   img.image_data AS image_data,
                   img.image_format AS image_format,
                   img.page_title AS page_title,
                   img.section AS section,
                   img.surrounding_text AS surrounding_text,
                   img.alt_text AS alt_text
            """,
            facility=facility,
            token=claim_token,
        )
        claimed = list(result)

        logger.debug(
            "claim_images_for_scoring: requested %d, won %d (token=%s)",
            limit,
            len(claimed),
            claim_token[:8],
        )
        return claimed


@retry_on_deadlock()
def mark_images_scored(
    facility: str,
    results: list[dict[str, Any]],
) -> int:
    """Mark images as captioned with VLM scoring results.

    Updates image status to 'captioned' and persists caption + scoring fields.
    Uses batched UNWIND for efficient graph updates.
    """
    if not results:
        return 0

    with GraphClient() as gc:
        gc.query(
            """
            UNWIND $batch AS item
            MATCH (img:Image {id: item.id})
            SET img.status = 'captioned',
                img.caption = item.caption,
                img.ocr_text = item.ocr_text,
                img.purpose = item.purpose,
                img.description = item.description,
                img.score = item.score,
                img.score_data_documentation = item.score_data_documentation,
                img.score_physics_content = item.score_physics_content,
                img.score_code_documentation = item.score_code_documentation,
                img.score_data_access = item.score_data_access,
                img.score_calibration = item.score_calibration,
                img.score_imas_relevance = item.score_imas_relevance,
                img.reasoning = item.reasoning,
                img.keywords = item.keywords,
                img.physics_domain = item.physics_domain,
                img.should_ingest = item.should_ingest,
                img.skip_reason = item.skip_reason,
                img.score_cost = item.score_cost,
                img.scored_at = datetime(),
                img.captioned_at = datetime(),
                img.claimed_at = null
            """,
            batch=results,
        )

    logger.info(
        "mark_images_scored: updated %d images to captioned for %s",
        len(results),
        facility,
    )
    return len(results)


def _release_claimed_images(image_ids: list[str]) -> None:
    """Release claimed images back to pool (e.g., on VLM failure)."""
    if not image_ids:
        return
    try:
        with GraphClient() as gc:
            gc.query(
                """
                UNWIND $ids AS id
                MATCH (img:Image {id: id})
                SET img.claimed_at = null
                """,
                ids=image_ids,
            )
    except Exception as e:
        logger.warning("Could not release %d claimed images: %s", len(image_ids), e)


# =============================================================================
# Bulk Page Discovery (unified)
# =============================================================================


def _persist_discovered_pages(
    facility: str,
    pages: list,
    on_progress: Callable | None = None,
) -> int:
    """Persist discovered pages to graph as 'scanned' status.

    Converts DiscoveredPage instances to batch_data and calls
    _bulk_create_wiki_pages for efficient UNWIND insertion.

    Args:
        facility: Facility ID
        pages: List of DiscoveredPage instances (name, url attrs)
        on_progress: Optional progress callback

    Returns:
        Number of pages created/updated
    """
    from imas_codex.discovery.wiki.scraper import canonical_page_id

    if not pages:
        return 0

    batch_data = [
        {
            "id": canonical_page_id(page.name, facility),
            "title": page.name,
            "url": page.url,
        }
        for page in pages
    ]

    with GraphClient() as gc:
        created = _bulk_create_wiki_pages(
            gc, facility, batch_data, on_progress=on_progress
        )

    logger.info("Created/updated %d pages in graph (scanned status)", created)
    if on_progress:
        on_progress(f"created {created} pages", None)

    return created


def bulk_discover_pages(
    facility: str,
    site_type: str,
    base_url: str,
    *,
    ssh_host: str | None = None,
    auth_type: str | None = None,
    credential_service: str | None = None,
    access_method: str = "direct",
    webs: list[str] | None = None,
    data_path: str | None = None,
    web_name: str = "Main",
    exclude_patterns: list[str] | None = None,
    exclude_prefixes: list[str] | None = None,
    on_progress: Callable | None = None,
) -> int:
    """Bulk discover all wiki pages for a site using the appropriate adapter.

    Unified entry point replacing the per-platform bulk_discover_all_pages_*
    functions. Creates the adapter, runs discovery, and persists to graph.

    Args:
        facility: Facility ID
        site_type: Wiki platform (mediawiki, twiki, twiki_static, twiki_raw,
            static_html)
        base_url: Wiki base URL
        ssh_host: SSH host for proxied access
        auth_type: Authentication type (tequila, keycloak, basic, session, etc.)
        credential_service: Keyring service for credentials
        access_method: Access method ("direct" or "vpn")
        webs: TWiki web names (for twiki sites)
        data_path: TWiki data directory path (for twiki_raw)
        web_name: TWiki web name (for twiki_raw, default "Main")
        exclude_patterns: Topic name regex patterns to skip (for twiki_raw)
        exclude_prefixes: URL path prefixes to exclude (for static_html)
        on_progress: Progress callback(msg, stats)

    Returns:
        Number of pages discovered
    """
    from imas_codex.discovery.wiki.adapters import get_adapter

    logger.info(
        "Starting bulk page discovery: site_type=%s, base_url=%s, auth=%s",
        site_type,
        base_url,
        auth_type,
    )

    # Build adapter kwargs
    adapter_kwargs: dict[str, Any] = {
        "ssh_host": ssh_host,
        "credential_service": credential_service,
    }

    # MediaWiki auth session setup
    session = None
    close_session = False

    if site_type == "mediawiki":
        if auth_type == "tequila" and credential_service:
            # Tequila uses MediaWikiClient (handled by adapter)
            from imas_codex.discovery.wiki.mediawiki import MediaWikiClient

            wiki_client = MediaWikiClient(
                base_url=base_url,
                credential_service=credential_service,
                verify_ssl=False,
            )
            if not wiki_client.authenticate():
                logger.error("Tequila authentication failed for %s", base_url)
                return 0
            adapter_kwargs["wiki_client"] = wiki_client
        elif auth_type == "keycloak" and credential_service:
            from imas_codex.discovery.wiki.keycloak import KeycloakSession

            ks = KeycloakSession(credential_service)
            try:
                session = ks.login(f"{base_url}/")
                adapter_kwargs["session"] = session
                close_session = True
            except RuntimeError as e:
                logger.error("Keycloak authentication failed: %s", e)
                return 0
        elif auth_type == "basic" and credential_service:
            import requests

            from imas_codex.discovery.wiki.auth import CredentialManager

            cred_mgr = CredentialManager()
            creds = cred_mgr.get_credentials(credential_service)
            if not creds:
                logger.error("No credentials for %s", credential_service)
                return 0
            session = requests.Session()
            session.auth = (creds[0], creds[1])
            session.verify = False
            adapter_kwargs["session"] = session
            close_session = True

    # Platform-specific adapter kwargs
    if site_type == "twiki":
        adapter_kwargs["webs"] = webs or ["Main"]
        adapter_kwargs["base_url"] = base_url
    elif site_type == "twiki_static":
        adapter_kwargs["base_url"] = base_url
        adapter_kwargs["access_method"] = access_method
    elif site_type == "twiki_raw":
        adapter_kwargs["data_path"] = data_path or base_url
        adapter_kwargs["web_name"] = web_name
        adapter_kwargs["exclude_patterns"] = exclude_patterns
    elif site_type == "static_html":
        adapter_kwargs["base_url"] = base_url
        adapter_kwargs["access_method"] = access_method
        adapter_kwargs["exclude_prefixes"] = exclude_prefixes

    try:
        adapter = get_adapter(site_type, **adapter_kwargs)

        pages = adapter.bulk_discover_pages(facility, base_url, on_progress)

        if not pages:
            logger.warning("No pages discovered from %s (%s)", base_url, site_type)
            return 0

        logger.info("Discovered %d pages from %s", len(pages), base_url)
        return _persist_discovered_pages(facility, pages, on_progress)
    finally:
        if close_session and session is not None:
            session.close()


# =============================================================================
# Bulk Artifact Discovery via MediaWiki API
# =============================================================================


def bulk_discover_artifacts(
    facility: str,
    base_url: str,
    site_type: str = "mediawiki",
    ssh_host: str | None = None,
    wiki_client: Any = None,
    credential_service: str | None = None,
    access_method: str = "direct",
    data_path: str | None = None,
    pub_path: str | None = None,
    on_progress: Callable | None = None,
) -> tuple[int, dict[str, list[str]]]:
    """Bulk discover all wiki artifacts via platform API.

    This is much faster than scanning pages - uses dedicated APIs:
    - MediaWiki: list=allimages API (returns all files in one call)
    - TWiki: /pub/ directory listing
    - TWiki static: Parse topic pages for linked files
    - TWiki raw: List files in pub directory via SSH
    - Confluence: /rest/api/content/{id}/child/attachment

    Args:
        facility: Facility ID
        base_url: Wiki base URL
        site_type: Wiki platform type
        ssh_host: SSH host for proxied access
        wiki_client: Authenticated MediaWikiClient (for Tequila)
        credential_service: Keyring service name
        access_method: Access method ("direct" or "vpn")
        data_path: TWiki data directory path (for twiki_raw)
        pub_path: TWiki pub directory path (for twiki_raw)
        on_progress: Progress callback

    Returns:
        Tuple of (count, page_artifacts) where page_artifacts maps
        page names to lists of artifact filenames discovered on that page.
    """
    from imas_codex.discovery.wiki.adapters import get_adapter

    logger.debug(f"Starting bulk artifact discovery for {site_type}...")

    # Get the appropriate adapter
    adapter = get_adapter(
        site_type=site_type,
        ssh_host=ssh_host,
        wiki_client=wiki_client,
        credential_service=credential_service,
        base_url=base_url,  # Needed for static site adapters
        access_method=access_method,
        data_path=data_path,
        pub_path=pub_path,
    )

    # Discover artifacts
    artifacts = adapter.bulk_discover_artifacts(facility, base_url, on_progress)

    if not artifacts:
        logger.debug("No artifacts discovered")
        return 0, {}

    logger.debug(f"Discovered {len(artifacts)} artifacts")

    # Create artifact nodes in graph using shared helper (UNWIND + FACILITY_ID)
    # Include linked_pages for HAS_ARTIFACT relationship creation
    batch_data = [
        {
            "id": f"{facility}:{a.filename}",
            "filename": a.filename,
            "url": a.url,
            "artifact_type": a.artifact_type,
            "size_bytes": a.size_bytes,
            "mime_type": a.mime_type,
            "linked_pages": a.linked_pages,
        }
        for a in artifacts
    ]

    with GraphClient() as gc:
        created = _bulk_create_wiki_artifacts(
            gc, facility, batch_data, on_progress=on_progress
        )

    logger.info(f"Created/updated {created} artifact nodes in graph")
    if on_progress:
        on_progress(f"created {created} artifacts", None)

    # Build page → artifact filenames mapping for CLI display
    page_artifacts: dict[str, list[str]] = {}
    for a in artifacts:
        if a.linked_pages:
            for page in a.linked_pages:
                page_artifacts.setdefault(page, []).append(a.filename)
        else:
            page_artifacts.setdefault("(unlinked)", []).append(a.filename)

    return created, page_artifacts


# =============================================================================
def _is_artifact(link: str) -> bool:
    """Check if a link points to an artifact (PDF, image, etc.)."""
    artifact_extensions = {
        ".pdf",
        ".doc",
        ".docx",
        ".ppt",
        ".pptx",
        ".xls",
        ".xlsx",
        ".csv",
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".svg",
        ".zip",
        ".tar",
        ".gz",
        ".h5",
        ".hdf5",
        ".mat",
        ".ipynb",
    }
    link_lower = link.lower()
    return any(link_lower.endswith(ext) for ext in artifact_extensions)


def _get_artifact_type(link: str) -> str:
    """Get artifact type from link extension."""
    link_lower = link.lower()
    if link_lower.endswith(".pdf"):
        return "pdf"
    if link_lower.endswith((".doc", ".docx", ".odt", ".rtf")):
        return "document"
    if link_lower.endswith((".ppt", ".pptx", ".key")):
        return "presentation"
    if link_lower.endswith((".xls", ".xlsx", ".csv")):
        return "spreadsheet"
    if link_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".svg")):
        return "image"
    if link_lower.endswith((".ipynb", ".nb")):
        return "notebook"
    if link_lower.endswith((".h5", ".hdf5", ".mat", ".nc")):
        return "data"
    if link_lower.endswith((".zip", ".tar", ".gz", ".tgz")):
        return "archive"
    return "document"


# =============================================================================
# Async Workers
# =============================================================================


async def score_worker(
    state: WikiDiscoveryState,
    on_progress: Callable | None = None,
) -> None:
    """Score worker: Content-aware LLM scoring in single pass.

    Transitions: scanned → scored

    Fetches page content preview, then scores with LLM.
    Uses centralized LLM access via get_model_for_task().
    Cost is tracked from actual OpenRouter response.
    """
    from imas_codex.agentic.agents import get_model_for_task

    worker_id = id(asyncio.current_task())
    logger.info(f"score_worker started (task={worker_id})")

    # Use effective semaphore: SSH semaphore (4) for SSH-based sites,
    # HTTP semaphore (30) for direct HTTP. Prevents overwhelming remote
    # hosts with concurrent SSH subprocess calls.
    fetch_semaphore = state.effective_fetch_semaphore

    # Get shared async wiki client for Tequila auth (native async HTTP)
    logger.debug(f"score_worker {worker_id}: getting async wiki client")
    shared_async_wiki_client = (
        await state.get_async_wiki_client()
        if state.auth_type in ("tequila", "session")
        else None
    )
    # Get shared Keycloak client for keycloak auth (OIDC via oauth2-proxy)
    shared_keycloak_client = (
        await state.get_keycloak_client() if state.auth_type == "keycloak" else None
    )
    # Get shared HTTP Basic auth client (e.g. JET wikis)
    shared_basic_auth_client = (
        await state.get_basic_auth_client() if state.auth_type == "basic" else None
    )
    logger.debug(f"score_worker {worker_id}: got async wiki client")

    async def fetch_content_for_page(page: dict) -> dict:
        """Fetch content preview for a single page."""
        async with fetch_semaphore:
            url = page.get("url", "")
            try:
                preview = await _fetch_and_summarize(
                    url,
                    state.ssh_host,
                    auth_type=state.auth_type,
                    credential_service=state.credential_service,
                    max_chars=1500,  # Reduced for scoring
                    async_wiki_client=shared_async_wiki_client,  # Native async HTTP
                    keycloak_client=shared_keycloak_client,
                    basic_auth_client=shared_basic_auth_client,
                )
                return {
                    "id": page["id"],
                    "title": page.get("title", ""),
                    "url": url,
                    "preview_text": preview,
                    "fetch_error": None,
                }
            except Exception as e:
                logger.debug("Failed to fetch %s: %s", url, e)
                return {
                    "id": page["id"],
                    "title": page.get("title", ""),
                    "url": url,
                    "preview_text": "",
                    "fetch_error": str(e),
                }

    while not state.should_stop_scoring():
        # Increased batch size from 25 to 50 for better LLM throughput
        # Run blocking Neo4j call in thread pool to avoid blocking event loop
        logger.debug(f"score_worker {worker_id}: claiming pages...")
        try:
            pages = await asyncio.to_thread(claim_pages_for_scoring, state.facility, 50)
        except Exception as e:
            # Neo4j connection error - backoff and retry
            logger.warning("score_worker %s: claim failed: %s", worker_id, e)
            await asyncio.sleep(5.0)
            continue
        logger.debug(f"score_worker {worker_id}: claimed {len(pages)} pages")

        if not pages:
            state.score_idle_count += 1
            if on_progress:
                on_progress("idle", state.score_stats)
            await asyncio.sleep(1.0)
            continue

        state.score_idle_count = 0

        if on_progress:
            on_progress(f"fetching {len(pages)} pages", state.score_stats)

        # Step 1: Fetch content for all pages in parallel
        fetch_tasks = [fetch_content_for_page(page) for page in pages]
        fetched_pages = await asyncio.gather(*fetch_tasks)
        logger.debug(f"score_worker {worker_id}: fetched {len(fetched_pages)} pages")

        # Step 1b: Separate pages with content from those where fetch failed.
        # Pages without preview_text cannot be scored meaningfully — release
        # them back so they can be retried when the server is responsive.
        pages_with_content = [p for p in fetched_pages if p.get("preview_text")]
        pages_no_content = [p for p in fetched_pages if not p.get("preview_text")]

        if pages_no_content:
            logger.info(
                "score_worker %s: %d/%d pages had no content (fetch failed), "
                "releasing for retry",
                worker_id,
                len(pages_no_content),
                len(fetched_pages),
            )
            await asyncio.to_thread(
                _release_claimed_pages, [p["id"] for p in pages_no_content]
            )

        if not pages_with_content:
            logger.debug(
                f"score_worker {worker_id}: no pages with content, skipping LLM"
            )
            continue

        if on_progress:
            on_progress(f"scoring {len(pages_with_content)} pages", state.score_stats)

        try:
            # Step 2: Score batch with LLM (only pages that have content)
            model = get_model_for_task("discovery")
            logger.debug(f"score_worker {worker_id}: starting LLM scoring...")
            results, cost = await _score_pages_batch(
                pages_with_content, model, state.focus
            )
            logger.debug(
                f"score_worker {worker_id}: LLM scored {len(results)} pages, cost=${cost:.4f}"
            )

            # Add preview_text to results for persistence
            for r in results:
                matching_page = next(
                    (p for p in pages_with_content if p["id"] == r["id"]), {}
                )
                r["preview_text"] = matching_page.get("preview_text", "")
                r["score_cost"] = cost / len(results) if results else 0.0

            # Run blocking Neo4j call in thread pool to avoid blocking event loop
            await asyncio.to_thread(mark_pages_scored, state.facility, results)
            state.score_stats.processed += len(results)
            state.score_stats.cost += cost

            if on_progress:
                on_progress(
                    f"scored {len(results)} pages", state.score_stats, results=results
                )

        except ValueError as e:
            # LLM validation error (e.g., truncated JSON) - release pages and continue
            # Don't stop the whole process; pages will be reclaimed after timeout
            logger.warning(
                "LLM failed for batch of %d pages: %s. "
                "Pages reverted to scanned status for retry.",
                len(pages),
                e,
            )
            state.score_stats.errors = getattr(state.score_stats, "errors", 0) + 1
            # Release pages by clearing claimed_at (not marking as failed)
            # Run blocking Neo4j call in thread pool
            await asyncio.to_thread(_release_claimed_pages, [p["id"] for p in pages])
            # Continue processing - don't stop the whole discovery
            continue
        except Exception as e:
            logger.error("Error in scoring batch: %s", e)
            # Run blocking Neo4j calls in thread pool
            for page in pages:
                await asyncio.to_thread(
                    mark_page_failed,
                    page["id"],
                    str(e),
                    WikiPageStatus.scanned.value,
                )


async def ingest_worker(
    state: WikiDiscoveryState,
    on_progress: Callable | None = None,
    min_score: float = 0.5,
) -> None:
    """Ingest worker: Chunk and embed high-value scored pages.

    Transitions: scored → ingested

    Claims scored pages with score >= min_score, fetches full content,
    chunks it, and creates embeddings.

    The ingest worker continues running even after cost limit is reached
    to drain the ingest queue. This ensures all scored content gets ingested.

    PERF: Pages are processed in parallel using asyncio.gather() with a
    semaphore to limit concurrency. This provides ~5x speedup over sequential.
    """
    # Get shared async wiki client for Tequila auth (native async HTTP)
    shared_async_wiki_client = (
        await state.get_async_wiki_client()
        if state.auth_type in ("tequila", "session")
        else None
    )
    # Get shared Keycloak client for keycloak auth (OIDC via oauth2-proxy)
    shared_keycloak_client = (
        await state.get_keycloak_client() if state.auth_type == "keycloak" else None
    )
    # Get shared HTTP Basic auth client (e.g. JET wikis)
    shared_basic_auth_client = (
        await state.get_basic_auth_client() if state.auth_type == "basic" else None
    )

    # Use effective semaphore: SSH semaphore (4) for SSH-based sites,
    # HTTP semaphore (30) for direct HTTP. Prevents overwhelming remote
    # hosts with concurrent SSH subprocess calls.
    http_semaphore = state.effective_fetch_semaphore

    async def process_single_page(page: dict) -> dict | None:
        """Process a single page with semaphore-limited concurrency."""
        page_id = page["id"]
        url = page.get("url", "")

        async with http_semaphore:
            try:
                chunk_count = await _ingest_page(
                    url=url,
                    page_id=page_id,
                    facility=state.facility,
                    site_type=state.site_type,
                    ssh_host=state.ssh_host,
                    auth_type=state.auth_type,
                    credential_service=state.credential_service,
                    async_wiki_client=shared_async_wiki_client,
                    keycloak_client=shared_keycloak_client,
                    basic_auth_client=shared_basic_auth_client,
                )
                return {
                    "id": page_id,
                    "chunk_count": chunk_count,
                    "score": page.get("score"),
                    "description": page.get("description", ""),
                    "physics_domain": page.get("physics_domain"),
                }
            except Exception as e:
                logger.warning("Error ingesting %s: %s", page_id, e)
                # Run blocking Neo4j call in thread pool
                await asyncio.to_thread(
                    mark_page_failed, page_id, str(e), WikiPageStatus.scored.value
                )
                return None

    while not state.should_stop_ingesting():
        # Increased batch size from 10 to 20 for better embedding throughput
        # Run blocking Neo4j call in thread pool to avoid blocking event loop
        try:
            pages = await asyncio.to_thread(
                claim_pages_for_ingesting, state.facility, min_score, 20
            )
        except Exception as e:
            # Neo4j connection error - backoff and retry
            logger.warning("ingest_worker: claim failed: %s", e)
            await asyncio.sleep(5.0)
            continue

        if not pages:
            state.ingest_idle_count += 1
            if on_progress:
                on_progress("idle", state.ingest_stats)
            await asyncio.sleep(1.0)
            continue

        state.ingest_idle_count = 0

        if on_progress:
            on_progress(f"ingesting {len(pages)} pages", state.ingest_stats)

        # Process all pages in parallel with semaphore-limited concurrency
        tasks = [process_single_page(page) for page in pages]
        results_raw = await asyncio.gather(*tasks)

        # Filter out None results (failed pages)
        results = [r for r in results_raw if r is not None]

        # Run blocking Neo4j call in thread pool
        await asyncio.to_thread(mark_pages_ingested, state.facility, results)
        state.ingest_stats.processed += len(results)

        if on_progress:
            on_progress(
                f"ingested {len(results)} pages", state.ingest_stats, results=results
            )


async def artifact_worker(
    state: WikiDiscoveryState,
    on_progress: Callable | None = None,
    max_size_mb: float = 5.0,
) -> None:
    """Artifact worker: Download and ingest wiki artifacts.

    Transitions: discovered → ingesting → ingested

    Claims discovered artifacts with supported types (pdf, docx, pptx, xlsx, ipynb),
    downloads content, and extracts text.
    Unsupported artifact types are marked as deferred.

    Args:
        state: Shared discovery state
        on_progress: Progress callback (msg, stats, results=None)
        max_size_mb: Maximum artifact size in MB
    """
    from imas_codex.discovery.wiki.pipeline import (
        WikiArtifactPipeline,
        fetch_artifact_content,
        fetch_artifact_size,
    )

    max_size_bytes = int(max_size_mb * 1024 * 1024)
    pipeline = WikiArtifactPipeline(
        facility_id=state.facility,
        max_size_mb=max_size_mb,
        use_rich=False,
    )

    while not state.should_stop_artifact_worker():
        # Increased batch size from 3 to 8 for better throughput
        # Run blocking Neo4j call in thread pool to avoid blocking event loop
        try:
            artifacts = await asyncio.to_thread(
                claim_artifacts_for_ingesting, state.facility, limit=8
            )
        except Exception as e:
            # Neo4j connection error - backoff and retry
            logger.warning("artifact_worker: claim failed: %s", e)
            await asyncio.sleep(5.0)
            continue

        if not artifacts:
            state.artifact_idle_count += 1
            if on_progress:
                on_progress("idle", state.artifact_stats)
            await asyncio.sleep(1.0)
            continue

        state.artifact_idle_count = 0

        if on_progress:
            on_progress(f"ingesting {len(artifacts)} artifacts", state.artifact_stats)

        results = []
        for artifact in artifacts:
            artifact_id = artifact["id"]
            artifact_type = artifact.get("artifact_type", "unknown")
            url = artifact.get("url", "")
            filename = artifact.get("filename", "unknown")

            try:
                # Check size before downloading
                size_bytes = fetch_artifact_size(url, facility=state.facility)

                if size_bytes is not None and size_bytes > max_size_bytes:
                    size_mb = size_bytes / (1024 * 1024)
                    reason = (
                        f"File size {size_mb:.1f} MB exceeds limit {max_size_mb:.1f} MB"
                    )
                    logger.info("Deferring oversized artifact %s: %s", filename, reason)
                    # Run blocking Neo4j call in thread pool
                    await asyncio.to_thread(mark_artifact_deferred, artifact_id, reason)
                    continue

                # Check if type is supported
                if artifact_type.lower() not in SUPPORTED_ARTIFACT_TYPES:
                    reason = f"Artifact type '{artifact_type}' not supported"
                    # Run blocking Neo4j call in thread pool
                    await asyncio.to_thread(mark_artifact_deferred, artifact_id, reason)
                    continue

                # Download and ingest
                _, content = await fetch_artifact_content(url, facility=state.facility)
                stats = await pipeline.ingest_artifact(
                    artifact_id, content, artifact_type
                )
                results.append(
                    {
                        "id": artifact_id,
                        "chunk_count": stats["chunks"],
                        "filename": filename,
                        "artifact_type": artifact_type,
                        "score": artifact.get("score"),
                        "physics_domain": artifact.get("physics_domain"),
                        "description": artifact.get("description", ""),
                    }
                )

            except Exception as e:
                # Unwrap common error wrappers to get the actual error message
                error_msg = str(e)
                if hasattr(e, "__cause__") and e.__cause__:
                    error_msg = str(e.__cause__)
                elif "RetryError" in type(e).__name__ and hasattr(e, "last_attempt"):
                    # tenacity RetryError - get the underlying exception
                    try:
                        error_msg = str(e.last_attempt.exception())
                    except Exception:
                        pass
                logger.warning(
                    "Error ingesting artifact %s: %s", artifact_id, error_msg
                )
                # Run blocking Neo4j call in thread pool
                await asyncio.to_thread(mark_artifact_failed, artifact_id, error_msg)

        # Run blocking Neo4j call in thread pool
        await asyncio.to_thread(mark_artifacts_ingested, state.facility, results)
        state.artifact_stats.processed += len(results)

        if on_progress:
            on_progress(
                f"ingested {len(results)} artifacts",
                state.artifact_stats,
                results=results,
            )


async def artifact_score_worker(
    state: WikiDiscoveryState,
    on_progress: Callable | None = None,
) -> None:
    """Artifact score worker: LLM scoring with text preview extraction.

    Transitions: discovered → scored

    Claims discovered artifacts with supported types, downloads a small
    portion of content to extract a text preview, then scores with LLM.
    After scoring, artifacts with score >= 0.5 become eligible for full ingestion.

    Uses the same scoring dimensions as wiki page scoring for consistency.
    """
    from imas_codex.agentic.agents import get_model_for_task

    worker_id = id(asyncio.current_task())
    logger.info(f"artifact_score_worker started (task={worker_id})")

    while not state.should_stop_artifact_scoring():
        # Claim artifacts for scoring
        try:
            artifacts = await asyncio.to_thread(
                claim_artifacts_for_scoring, state.facility, 10
            )
        except Exception as e:
            logger.warning("artifact_score_worker %s: claim failed: %s", worker_id, e)
            await asyncio.sleep(5.0)
            continue

        if not artifacts:
            state.artifact_score_idle_count += 1
            if on_progress:
                on_progress("idle", state.artifact_score_stats)
            await asyncio.sleep(1.0)
            continue

        state.artifact_score_idle_count = 0

        if on_progress:
            on_progress(
                f"extracting text from {len(artifacts)} artifacts",
                state.artifact_score_stats,
            )

        # Step 1: Extract text preview from each artifact
        artifacts_with_text = []
        for artifact in artifacts:
            artifact_id = artifact["id"]
            url = artifact.get("url", "")
            filename = artifact.get("filename", "")
            artifact_type = artifact.get("artifact_type", "unknown")

            try:
                preview_text = await _extract_artifact_preview(
                    url=url,
                    artifact_type=artifact_type,
                    facility=state.facility,
                    max_chars=1500,
                )
                artifacts_with_text.append(
                    {
                        "id": artifact_id,
                        "filename": filename,
                        "url": url,
                        "artifact_type": artifact_type,
                        "size_bytes": artifact.get("size_bytes"),
                        "preview_text": preview_text,
                    }
                )
            except Exception as e:
                logger.debug("Failed to extract preview for %s: %s", filename, e)
                # Track failed extractions for release
                artifacts_with_text.append(
                    {
                        "id": artifact_id,
                        "filename": filename,
                        "url": url,
                        "artifact_type": artifact_type,
                        "size_bytes": artifact.get("size_bytes"),
                        "preview_text": "",
                    }
                )

        # Separate artifacts with content from those where extraction failed.
        # Artifacts without preview_text are released for retry.
        artifacts_to_score = [a for a in artifacts_with_text if a.get("preview_text")]
        artifacts_no_content = [
            a for a in artifacts_with_text if not a.get("preview_text")
        ]

        if artifacts_no_content:
            logger.info(
                "artifact_score_worker %s: %d/%d artifacts had no content, "
                "releasing for retry",
                worker_id,
                len(artifacts_no_content),
                len(artifacts_with_text),
            )
            try:
                with GraphClient() as gc:
                    gc.query(
                        """
                        UNWIND $ids AS id
                        MATCH (wa:WikiArtifact {id: id})
                        SET wa.claimed_at = null
                        """,
                        ids=[a["id"] for a in artifacts_no_content],
                    )
            except Exception:
                pass

        if not artifacts_to_score:
            logger.debug(
                "artifact_score_worker %s: no artifacts with content, skipping LLM",
                worker_id,
            )
            continue

        if on_progress:
            on_progress(
                f"scoring {len(artifacts_to_score)} artifacts",
                state.artifact_score_stats,
            )

        try:
            # Step 2: Score batch with LLM (only artifacts that have content)
            model = get_model_for_task("discovery")
            results, cost = await _score_artifacts_batch(
                artifacts_to_score, model, state.focus
            )

            # Add preview_text to results for persistence
            for r in results:
                matching = next(
                    (a for a in artifacts_to_score if a["id"] == r["id"]), {}
                )
                r["preview_text"] = matching.get("preview_text", "")[:500]
                r["score_cost"] = cost / len(results) if results else 0.0

            # Persist scores to graph
            await asyncio.to_thread(mark_artifacts_scored, state.facility, results)
            state.artifact_score_stats.processed += len(results)
            state.artifact_score_stats.cost += cost

            if on_progress:
                on_progress(
                    f"scored {len(results)} artifacts",
                    state.artifact_score_stats,
                    results=results,
                )

        except ValueError as e:
            logger.warning(
                "LLM failed for artifact batch of %d: %s. "
                "Artifacts reverted to discovered status for retry.",
                len(artifacts),
                e,
            )
            state.artifact_score_stats.errors = (
                getattr(state.artifact_score_stats, "errors", 0) + 1
            )
            # Release artifacts by clearing claimed_at
            try:
                with GraphClient() as gc:
                    gc.query(
                        """
                        UNWIND $ids AS id
                        MATCH (wa:WikiArtifact {id: id})
                        SET wa.claimed_at = null
                        """,
                        ids=[a["id"] for a in artifacts],
                    )
            except Exception:
                pass
            continue
        except Exception as e:
            logger.error("Error in artifact scoring batch: %s", e)
            for artifact in artifacts:
                await asyncio.to_thread(mark_artifact_failed, artifact["id"], str(e))


async def _extract_artifact_preview(
    url: str,
    artifact_type: str,
    facility: str,
    max_chars: int = 1500,
) -> str:
    """Extract a text preview from an artifact for LLM scoring.

    Downloads the artifact content and extracts text from the first
    portion. This is a lightweight extraction for scoring purposes only.

    For PDFs: extracts text from first few pages.
    For documents: extracts text paragraphs.
    For presentations: extracts slide text.
    For notebooks: extracts cell content.

    Args:
        url: Artifact download URL
        artifact_type: Type of artifact (pdf, docx, pptx, xlsx, ipynb)
        facility: Facility ID (for SSH proxy)
        max_chars: Maximum characters to extract

    Returns:
        Extracted text preview or empty string on failure
    """
    from imas_codex.discovery.wiki.pipeline import fetch_artifact_content

    try:
        _, content_bytes = await fetch_artifact_content(url, facility=facility)
    except Exception as e:
        logger.debug("Failed to download %s: %s", url, e)
        return ""

    try:
        text = _extract_text_from_bytes(content_bytes, artifact_type)
        return text[:max_chars] if text else ""
    except Exception as e:
        logger.debug("Failed to extract text from %s: %s", url, e)
        return ""


def _extract_text_from_bytes(content_bytes: bytes, artifact_type: str) -> str:
    """Extract text from artifact bytes based on type.

    Lightweight extraction for scoring preview. Does not use LlamaIndex
    to avoid heavyweight dependencies in the scoring path.
    """
    import tempfile
    from pathlib import Path

    at = artifact_type.lower()

    if at == "pdf":
        # Extract text from first few pages
        if b"%PDF" not in content_bytes[:1024]:
            return ""
        try:
            import logging as _logging

            import pypdf

            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                f.write(content_bytes)
                temp_path = Path(f.name)
            try:
                # Suppress pypdf's verbose warnings about corrupt PDF objects
                # (e.g. "Ignoring wrong pointing object") which are benign
                pypdf_logger = _logging.getLogger("pypdf")
                original_level = pypdf_logger.level
                pypdf_logger.setLevel(_logging.ERROR)
                try:
                    reader = pypdf.PdfReader(temp_path)
                    text_parts = []
                    for page in reader.pages[:5]:  # First 5 pages
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                finally:
                    pypdf_logger.setLevel(original_level)
                return "\n\n".join(text_parts)
            finally:
                temp_path.unlink(missing_ok=True)
        except Exception:
            return ""

    elif at in ("docx", "doc"):
        try:
            from docx import Document as DocxDocument

            with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
                f.write(content_bytes)
                temp_path = Path(f.name)
            try:
                doc = DocxDocument(temp_path)
                paragraphs = [p.text for p in doc.paragraphs[:50] if p.text.strip()]
                return "\n\n".join(paragraphs)
            finally:
                temp_path.unlink(missing_ok=True)
        except Exception:
            return ""

    elif at in ("pptx", "ppt"):
        try:
            from pptx import Presentation

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                f.write(content_bytes)
                temp_path = Path(f.name)
            try:
                prs = Presentation(temp_path)
                text_parts = []
                for slide in list(prs.slides)[:10]:  # First 10 slides
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text.strip())
                return "\n".join(text_parts)
            finally:
                temp_path.unlink(missing_ok=True)
        except Exception:
            return ""

    elif at in ("xlsx", "xls"):
        try:
            from imas_codex.discovery.wiki.excel import extract_excel_preview

            return extract_excel_preview(content_bytes)
        except Exception:
            return ""

    elif at == "ipynb":
        try:
            import json

            nb = json.loads(content_bytes.decode("utf-8"))
            cells = nb.get("cells", [])
            text_parts = []
            for cell in cells[:20]:  # First 20 cells
                source = "".join(cell.get("source", []))
                if source.strip():
                    text_parts.append(source)
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    return ""


async def _score_artifacts_batch(
    artifacts: list[dict[str, Any]],
    model: str,
    focus: str | None = None,
) -> tuple[list[dict[str, Any]], float]:
    """Score a batch of artifacts using LLM with structured output.

    Uses litellm.acompletion with ArtifactScoreBatch Pydantic model for
    structured output. Content-based scoring with per-dimension scores.

    Args:
        artifacts: List of artifact dicts with id, filename, preview_text, etc.
        model: Model identifier from get_model_for_task()
        focus: Optional focus area for scoring

    Returns:
        (results, cost) tuple where cost is actual LLM cost from OpenRouter.
    """
    import os
    import re

    import litellm

    from imas_codex.agentic.prompt_loader import render_prompt
    from imas_codex.discovery.wiki.models import (
        ArtifactScoreBatch,
        grounded_artifact_score,
    )

    api_key = os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        raise ValueError("OPENROUTER_API_KEY environment variable not set.")

    model_id = model
    if not model_id.startswith("openrouter/"):
        model_id = f"openrouter/{model_id}"

    # Build system prompt using artifact-scorer template
    context: dict[str, Any] = {}
    if focus:
        context["focus"] = focus

    system_prompt = render_prompt("wiki/artifact-scorer", context)

    # Build user prompt with artifact content
    lines = [
        f"Score these {len(artifacts)} wiki artifacts based on their content.",
        "(Use the preview text to assess value for the IMAS knowledge graph.)\n",
    ]

    for i, a in enumerate(artifacts, 1):
        lines.append(f"\n## Artifact {i}")
        lines.append(f"ID: {a['id']}")
        lines.append(f"Filename: {a.get('filename', 'Unknown')}")
        lines.append(f"Type: {a.get('artifact_type', 'unknown')}")

        if a.get("size_bytes"):
            size_mb = a["size_bytes"] / (1024 * 1024)
            lines.append(f"Size: {size_mb:.1f} MB")

        preview = a.get("preview_text", "")
        if preview:
            lines.append(f"Content Preview:\n{preview[:800]}")
        else:
            lines.append("Content Preview: (not available - score from filename/type)")

        if a.get("url"):
            lines.append(f"URL: {a['url']}")

    lines.append(
        "\n\nReturn results for each artifact in order. "
        "The response format is enforced by the schema."
    )

    user_prompt = "\n".join(lines)

    # Retry loop
    max_retries = 5
    retry_base_delay = 4.0
    last_error = None
    total_cost = 0.0

    for attempt in range(max_retries):
        try:
            response = await litellm.acompletion(
                model=model_id,
                api_key=api_key,
                response_format=ArtifactScoreBatch,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
                temperature=0.3,
                max_tokens=16000,
                timeout=120,
            )

            input_tokens = response.usage.prompt_tokens
            output_tokens = response.usage.completion_tokens

            if (
                hasattr(response, "_hidden_params")
                and "response_cost" in response._hidden_params
            ):
                cost = response._hidden_params["response_cost"]
            else:
                cost = (input_tokens * 3 + output_tokens * 15) / 1_000_000

            total_cost += cost

            content = response.choices[0].message.content
            if not content:
                logger.warning("LLM returned empty response for artifacts")
                return [], total_cost

            content = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", content)
            content = content.encode("utf-8", errors="surrogateescape").decode(
                "utf-8", errors="replace"
            )

            batch = ArtifactScoreBatch.model_validate_json(content)
            llm_results = batch.results
            break

        except Exception as e:
            last_error = e
            error_msg = str(e).lower()

            is_retryable = any(
                x in error_msg
                for x in [
                    "overloaded",
                    "rate",
                    "429",
                    "503",
                    "timeout",
                    "eof",
                    "json",
                    "truncated",
                    "validation",
                ]
            )

            if is_retryable and attempt < max_retries - 1:
                delay = retry_base_delay * (2**attempt)
                logger.debug(
                    "LLM error (attempt %d/%d): %s. Retrying in %.1fs...",
                    attempt + 1,
                    max_retries,
                    str(e)[:100],
                    delay,
                )
                await asyncio.sleep(delay)
            else:
                logger.error(
                    "LLM failed for artifact batch of %d after %d attempts: %s. "
                    "Artifacts reverted to discovered status for retry.",
                    len(artifacts),
                    attempt + 1,
                    e,
                )
                raise ValueError(
                    f"LLM failed after {attempt + 1} attempts: {e}. "
                    f"Artifacts reverted to discovered status."
                ) from e
    else:
        raise last_error  # type: ignore[misc]

    # Convert to result dicts
    cost_per_artifact = total_cost / len(artifacts) if artifacts else 0.0
    results = []

    for r in llm_results[: len(artifacts)]:
        scores = {
            "score_data_documentation": r.score_data_documentation,
            "score_physics_content": r.score_physics_content,
            "score_code_documentation": r.score_code_documentation,
            "score_data_access": r.score_data_access,
            "score_calibration": r.score_calibration,
            "score_imas_relevance": r.score_imas_relevance,
        }

        combined_score = grounded_artifact_score(scores, r.artifact_purpose)

        # Find the matching artifact for filename
        matching = next((a for a in artifacts if a["id"] == r.id), {})

        results.append(
            {
                "id": r.id,
                "score": combined_score,
                "artifact_purpose": r.artifact_purpose.value,
                "description": r.description[:150],
                "reasoning": r.reasoning[:80],
                "keywords": r.keywords[:5],
                "physics_domain": r.physics_domain.value if r.physics_domain else None,
                "should_ingest": r.should_ingest,
                "skip_reason": r.skip_reason or None,
                "score_data_documentation": r.score_data_documentation,
                "score_physics_content": r.score_physics_content,
                "score_code_documentation": r.score_code_documentation,
                "score_data_access": r.score_data_access,
                "score_calibration": r.score_calibration,
                "score_imas_relevance": r.score_imas_relevance,
                "score_cost": cost_per_artifact,
                # Pass through filename for display
                "filename": matching.get("filename", ""),
                "artifact_type": matching.get("artifact_type", ""),
            }
        )

    return results, total_cost


# =============================================================================
# Image Score Worker (VLM caption + scoring in single pass)
# =============================================================================


async def image_score_worker(
    state: WikiDiscoveryState,
    on_progress: Callable | None = None,
) -> None:
    """Image score worker: VLM captioning + scoring in single pass.

    Transitions: ingested → captioned

    Claims images that have been ingested (image_data stored) but not yet
    captioned/scored. Sends image bytes + context to VLM, receives
    caption + scoring in one pass.
    """
    from imas_codex.agentic.agents import get_model_for_task

    worker_id = id(asyncio.current_task())
    logger.info(f"image_score_worker started (task={worker_id})")

    while not state.should_stop_image_scoring():
        try:
            images = await asyncio.to_thread(
                claim_images_for_scoring, state.facility, 10
            )
        except Exception as e:
            logger.warning("image_score_worker %s: claim failed: %s", worker_id, e)
            await asyncio.sleep(5.0)
            continue

        if not images:
            state.image_idle_count += 1
            if on_progress:
                on_progress("idle", state.image_stats)
            await asyncio.sleep(2.0)
            continue

        state.image_idle_count = 0

        # Filter out images without stored image_data
        images_with_data = [img for img in images if img.get("image_data")]
        images_no_data = [img for img in images if not img.get("image_data")]

        if images_no_data:
            logger.info(
                "image_score_worker: %d/%d images have no image_data, releasing",
                len(images_no_data),
                len(images),
            )
            await asyncio.to_thread(
                _release_claimed_images, [img["id"] for img in images_no_data]
            )

        if not images_with_data:
            continue

        if on_progress:
            on_progress(f"scoring {len(images_with_data)} images", state.image_stats)

        try:
            model = get_model_for_task("vlm")
            results, cost = await _score_images_batch(
                images_with_data, model, state.focus
            )

            # Persist to graph
            await asyncio.to_thread(mark_images_scored, state.facility, results)
            state.image_stats.processed += len(results)
            state.image_stats.cost += cost

            if on_progress:
                on_progress(
                    f"scored {len(results)} images",
                    state.image_stats,
                    results=results,
                )

        except ValueError as e:
            logger.warning(
                "VLM failed for batch of %d images: %s. Images released for retry.",
                len(images_with_data),
                e,
            )
            state.image_stats.errors = getattr(state.image_stats, "errors", 0) + 1
            await asyncio.to_thread(
                _release_claimed_images,
                [img["id"] for img in images_with_data],
            )


async def _score_images_batch(
    images: list[dict[str, Any]],
    model: str,
    focus: str | None = None,
) -> tuple[list[dict[str, Any]], float]:
    """Score a batch of images using VLM with structured output.

    Sends image bytes + context to VLM and receives caption + scoring
    in a single pass. Uses ImageScoreBatch Pydantic model.

    Returns:
        (results, cost) tuple
    """
    import os
    import re

    import litellm

    from imas_codex.agentic.prompt_loader import render_prompt
    from imas_codex.discovery.wiki.models import (
        ImageScoreBatch,
        grounded_image_score,
    )

    api_key = os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        raise ValueError("OPENROUTER_API_KEY environment variable not set.")

    model_id = model
    if not model_id.startswith("openrouter/"):
        model_id = f"openrouter/{model_id}"

    # Build system prompt
    context: dict[str, Any] = {}
    if focus:
        context["focus"] = focus
    system_prompt = render_prompt("wiki/image-captioner", context)

    # Build user message with image content
    user_content: list[dict[str, Any]] = [
        {
            "type": "text",
            "text": f"Score and caption these {len(images)} images from fusion facility documentation.\n",
        }
    ]

    for i, img in enumerate(images, 1):
        # Add text context for this image
        context_parts = [f"\n## Image {i}", f"ID: {img['id']}"]
        if img.get("page_title"):
            context_parts.append(f"Page: {img['page_title']}")
        if img.get("section"):
            context_parts.append(f"Section: {img['section']}")
        if img.get("surrounding_text"):
            context_parts.append(f"Context: {img['surrounding_text'][:500]}")
        if img.get("alt_text"):
            context_parts.append(f"Alt text: {img['alt_text']}")

        user_content.append({"type": "text", "text": "\n".join(context_parts)})

        # Add image data
        img_format = img.get("image_format", "webp")
        mime_type = f"image/{img_format}"
        user_content.append(
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:{mime_type};base64,{img['image_data']}",
                },
            }
        )

    user_content.append(
        {
            "type": "text",
            "text": "\n\nReturn results for each image in order. "
            "The response format is enforced by the schema.",
        }
    )

    # Retry loop
    max_retries = 5
    retry_base_delay = 4.0
    last_error = None
    total_cost = 0.0

    for attempt in range(max_retries):
        try:
            response = await litellm.acompletion(
                model=model_id,
                api_key=api_key,
                response_format=ImageScoreBatch,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_content},
                ],
                temperature=0.3,
                max_tokens=32000,
                timeout=180,  # VLM may need more time for images
            )

            input_tokens = response.usage.prompt_tokens
            output_tokens = response.usage.completion_tokens

            if (
                hasattr(response, "_hidden_params")
                and "response_cost" in response._hidden_params
            ):
                cost = response._hidden_params["response_cost"]
            else:
                # Fallback VLM rates
                cost = (input_tokens * 3 + output_tokens * 15) / 1_000_000

            total_cost += cost

            content = response.choices[0].message.content
            if not content:
                return [], total_cost

            content = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", content)
            content = content.encode("utf-8", errors="surrogateescape").decode(
                "utf-8", errors="replace"
            )

            batch = ImageScoreBatch.model_validate_json(content)
            llm_results = batch.results
            break

        except Exception as e:
            last_error = e
            error_msg = str(e).lower()
            is_retryable = any(
                x in error_msg
                for x in [
                    "overloaded",
                    "rate",
                    "429",
                    "503",
                    "timeout",
                    "eof",
                    "json",
                    "truncated",
                    "validation",
                ]
            )
            if is_retryable and attempt < max_retries - 1:
                delay = retry_base_delay * (2**attempt)
                logger.debug(
                    "VLM error (attempt %d/%d): %s. Retrying in %.1fs...",
                    attempt + 1,
                    max_retries,
                    str(e)[:100],
                    delay,
                )
                await asyncio.sleep(delay)
            else:
                raise ValueError(f"VLM failed after {attempt + 1} attempts: {e}") from e
    else:
        raise last_error  # type: ignore[misc]

    # Convert to result dicts with grounded scoring
    cost_per_image = total_cost / len(images) if images else 0.0
    results = []

    for r in llm_results[: len(images)]:
        scores = {
            "score_data_documentation": r.score_data_documentation,
            "score_physics_content": r.score_physics_content,
            "score_code_documentation": r.score_code_documentation,
            "score_data_access": r.score_data_access,
            "score_calibration": r.score_calibration,
            "score_imas_relevance": r.score_imas_relevance,
        }
        combined_score = grounded_image_score(scores, r.purpose)

        results.append(
            {
                "id": r.id,
                "caption": r.caption,
                "ocr_text": r.ocr_text,
                "purpose": r.purpose.value,
                "description": r.description[:150],
                "score": combined_score,
                "score_data_documentation": r.score_data_documentation,
                "score_physics_content": r.score_physics_content,
                "score_code_documentation": r.score_code_documentation,
                "score_data_access": r.score_data_access,
                "score_calibration": r.score_calibration,
                "score_imas_relevance": r.score_imas_relevance,
                "reasoning": r.reasoning[:80],
                "keywords": r.keywords[:5],
                "physics_domain": r.physics_domain.value if r.physics_domain else None,
                "should_ingest": r.should_ingest,
                "skip_reason": r.skip_reason or None,
                "score_cost": cost_per_image,
            }
        )

    return results, total_cost


# =============================================================================
# Worker Helpers
# =============================================================================


async def _fetch_html(
    url: str,
    ssh_host: str | None,
    auth_type: str | None = None,
    credential_service: str | None = None,
    async_wiki_client: Any = None,
    keycloak_client: Any = None,
    basic_auth_client: Any = None,
) -> str:
    """Fetch HTML content from URL.

    Args:
        url: Page URL
        ssh_host: Optional SSH host for proxied fetching
        auth_type: Authentication type (tequila, session, keycloak, basic, etc.)
        credential_service: Keyring service for credentials
        async_wiki_client: Shared AsyncMediaWikiClient for Tequila auth
        keycloak_client: Shared httpx.AsyncClient for Keycloak auth
        basic_auth_client: Shared httpx.AsyncClient with HTTP Basic auth

    Returns:
        HTML content string or empty string on error
    """

    def _ssh_fetch() -> str:
        """Fetch via SSH proxy."""
        cmd = f'curl -sk --noproxy "*" "{url}" 2>/dev/null'
        try:
            result = subprocess.run(
                ["ssh", ssh_host, cmd],
                capture_output=True,
                timeout=30,
            )
            if result.returncode == 0:
                # Wiki pages may be ISO-8859 encoded despite claiming UTF-8
                return result.stdout.decode("utf-8", errors="replace")
            return ""
        except Exception as e:
            logger.warning("SSH fetch failed for %s: %s", url, e)
            return ""

    async def _async_tequila_fetch() -> str:
        """Fetch with Tequila authentication using async client."""
        import urllib.parse as urlparse

        # Extract page name from URL
        page_name = url.split("/wiki/")[-1] if "/wiki/" in url else url.split("/")[-1]
        if "?" in page_name:
            parsed = urlparse.parse_qs(urlparse.urlparse(url).query)
            page_name = parsed.get("title", [page_name])[0]
        page_name = urlparse.unquote(page_name)

        # Use provided async client
        if async_wiki_client is not None:
            try:
                page = await async_wiki_client.get_page(page_name)
                if page:
                    return page.content_html
                return ""
            except Exception as e:
                logger.debug("Async client fetch failed for %s: %s", url, e)
                return ""

        # No client provided - create a new async client
        from imas_codex.discovery.wiki.mediawiki import AsyncMediaWikiClient

        base_url_local = (
            url.rsplit("/", 1)[0] if "/wiki/" in url else url.rsplit("/", 1)[0]
        )
        if "/wiki" in base_url_local:
            base_url_local = base_url_local.rsplit("/wiki", 1)[0] + "/wiki"

        async with AsyncMediaWikiClient(
            base_url=base_url_local,
            credential_service=credential_service or "tcv",
            verify_ssl=False,
        ) as client:
            if not await client.authenticate():
                logger.warning("Tequila auth failed for %s", url)
                return ""
            page = await client.get_page(page_name)
            return page.content_html if page else ""

    async def _async_keycloak_fetch() -> str:
        """Fetch with Keycloak auth using shared httpx.AsyncClient."""
        if keycloak_client is None:
            logger.warning("No Keycloak client available for %s", url)
            return ""
        try:
            response = await keycloak_client.get(url)
            if response.status_code == 200:
                return response.text
            logger.debug(
                "Keycloak fetch returned HTTP %d for %s", response.status_code, url
            )
            return ""
        except Exception as e:
            logger.debug("Keycloak fetch failed for %s: %s", url, e)
            return ""

    async def _async_basic_auth_fetch() -> str:
        """Fetch with HTTP Basic auth using shared httpx.AsyncClient."""
        if basic_auth_client is None:
            logger.warning("No HTTP Basic auth client available for %s", url)
            return ""
        try:
            response = await basic_auth_client.get(url)
            if response.status_code == 200:
                return response.text
            logger.debug(
                "Basic auth fetch returned HTTP %d for %s", response.status_code, url
            )
            return ""
        except Exception as e:
            logger.debug("Basic auth fetch failed for %s: %s", url, e)
            return ""

    # Determine fetch strategy - prefer direct HTTP over SSH when credentials available
    if auth_type in ("tequila", "session"):
        return await _async_tequila_fetch()
    elif auth_type == "keycloak" and keycloak_client:
        return await _async_keycloak_fetch()
    elif auth_type == "basic" and basic_auth_client:
        return await _async_basic_auth_fetch()
    elif ssh_host:
        return await asyncio.to_thread(_ssh_fetch)
    else:
        # Direct HTTP fetch (no auth) - already async
        from imas_codex.discovery.wiki.prefetch import fetch_page_content

        html, error = await fetch_page_content(url)
        if html:
            return html
        if error:
            logger.debug("HTTP fetch failed for %s: %s", url, error)
        return ""


async def _fetch_and_summarize(
    url: str,
    ssh_host: str | None,
    auth_type: str | None = None,
    credential_service: str | None = None,
    max_chars: int = 2000,
    async_wiki_client: Any = None,
    keycloak_client: Any = None,
    basic_auth_client: Any = None,
) -> str:
    """Fetch page content and extract text preview.

    No LLM is used here - prefetch extracts text deterministically.
    The summary is just cleaned text for the scorer to evaluate.

    Delegates HTML fetching to _fetch_html and applies text extraction.
    """
    from imas_codex.discovery.wiki.prefetch import extract_text_from_html

    html = await _fetch_html(
        url,
        ssh_host,
        auth_type=auth_type,
        credential_service=credential_service,
        async_wiki_client=async_wiki_client,
        keycloak_client=keycloak_client,
        basic_auth_client=basic_auth_client,
    )

    if html:
        return extract_text_from_html(html, max_chars=max_chars)
    return ""


async def _score_pages_batch(
    pages: list[dict[str, Any]],
    model: str,
    focus: str | None = None,
) -> tuple[list[dict[str, Any]], float]:
    """Score a batch of pages using LLM with structured output.

    Uses litellm.acompletion with WikiScoreBatch Pydantic model for
    structured output. Content-based scoring with per-dimension scores.

    The retry loop includes response parsing because JSON/validation
    errors from truncated responses are retryable — a fresh LLM call
    often returns valid output. Cost is accumulated across all attempts
    since API calls are billed regardless of parsing success.

    Args:
        pages: List of page dicts with id, title, summary, preview_text, etc.
        model: Model identifier from get_model_for_task()
        focus: Optional focus area for scoring

    Returns:
        (results, cost) tuple where cost is actual LLM cost from OpenRouter.
    """
    from imas_codex.agentic.prompt_loader import render_prompt
    from imas_codex.discovery.base.llm import acall_llm_structured
    from imas_codex.discovery.wiki.models import (
        WikiScoreBatch,
        grounded_wiki_score,
    )

    # Build system prompt using dynamic template with schema injection
    context: dict[str, Any] = {}
    if focus:
        context["focus"] = focus

    system_prompt = render_prompt("wiki/scorer", context)

    # Build user prompt with page content (not graph metrics)
    lines = [
        f"Score these {len(pages)} wiki pages based on their content.",
        "(Use the preview text to infer value - graph metrics like in_degree are NOT indicators.)\n",
    ]

    for i, p in enumerate(pages, 1):
        lines.append(f"\n## Page {i}")
        lines.append(f"ID: {p['id']}")
        lines.append(f"Title: {p.get('title', 'Unknown')}")

        # Use preview_text for content-based scoring (preferred over summary)
        preview = p.get("preview_text") or p.get("summary") or ""
        if preview:
            lines.append(f"Preview: {preview[:800]}")

        # Include URL for context (Confluence vs MediaWiki structure hints)
        url = p.get("url")
        if url:
            lines.append(f"URL: {url}")

    lines.append(
        "\n\nReturn results for each page in order. "
        "The response format is enforced by the schema."
    )

    user_prompt = "\n".join(lines)

    # Shared retry+parse loop handles both API errors and JSON/validation
    # errors from truncated responses. Cost accumulated across retries.
    # Model-aware token limits + timeout applied automatically.
    batch, total_cost, _tokens = await acall_llm_structured(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        response_model=WikiScoreBatch,
        temperature=0.3,
    )

    llm_results = batch.results

    # Convert to result dicts, computing combined scores
    cost_per_page = total_cost / len(pages) if pages else 0.0
    results = []

    for r in llm_results[: len(pages)]:
        # Build per-dimension scores dict
        scores = {
            "score_data_documentation": r.score_data_documentation,
            "score_physics_content": r.score_physics_content,
            "score_code_documentation": r.score_code_documentation,
            "score_data_access": r.score_data_access,
            "score_calibration": r.score_calibration,
            "score_imas_relevance": r.score_imas_relevance,
        }

        # Compute combined score using grounded function
        combined_score = grounded_wiki_score(scores, r.page_purpose)

        results.append(
            {
                "id": r.id,
                "score": combined_score,
                "purpose": r.page_purpose.value,
                "description": r.description[:150],
                "reasoning": r.reasoning[:80],
                "keywords": r.keywords[:5],
                "physics_domain": r.physics_domain.value if r.physics_domain else None,
                "should_ingest": r.should_ingest,
                "skip_reason": r.skip_reason or None,
                # Per-dimension scores
                "score_data_documentation": r.score_data_documentation,
                "score_physics_content": r.score_physics_content,
                "score_code_documentation": r.score_code_documentation,
                "score_data_access": r.score_data_access,
                "score_calibration": r.score_calibration,
                "score_imas_relevance": r.score_imas_relevance,
                # Legacy fields for compatibility
                "page_type": r.page_purpose.value,
                "is_physics": r.physics_domain is not None
                and r.physics_domain.value != "general",
                "score_cost": cost_per_page,
            }
        )

    return results, total_cost


def _score_pages_heuristic(pages: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Heuristic fallback scoring based on keywords. Zero cost."""
    results = []
    for page in pages:
        title = page.get("title", "").lower()
        summary = page.get("summary", "") or ""

        score = 0.5
        reasoning = "Default score"

        physics_keywords = [
            "thomson",
            "liuqe",
            "equilibrium",
            "mhd",
            "plasma",
            "diagnostic",
            "calibration",
            "signal",
            "node",
        ]
        low_value_keywords = [
            "meeting",
            "workshop",
            "todo",
            "draft",
            "notes",
            "personal",
            "test",
            "sandbox",
        ]

        for kw in physics_keywords:
            if kw in title or kw in summary.lower():
                score = min(score + 0.15, 1.0)
                reasoning = f"Contains physics keyword: {kw}"

        for kw in low_value_keywords:
            if kw in title:
                score = max(score - 0.2, 0.0)
                reasoning = f"Contains low-value keyword: {kw}"

        results.append(
            {
                "id": page["id"],
                "score": score,
                "reasoning": reasoning,
                "page_type": "documentation",
                "is_physics": score >= 0.6,
            }
        )

    return results


async def _ingest_page(
    url: str,
    page_id: str,
    facility: str,
    site_type: str,
    ssh_host: str | None,
    auth_type: str | None = None,
    credential_service: str | None = None,
    async_wiki_client: Any = None,
    keycloak_client: Any = None,
    basic_auth_client: Any = None,
) -> int:
    """Ingest a page: fetch content, chunk, and embed.

    Uses the WikiIngestionPipeline for proper chunking and embedding.

    Args:
        url: Page URL to fetch
        page_id: Unique page identifier
        facility: Facility ID (e.g., 'tcv', 'jet')
        site_type: Site type ('mediawiki', 'confluence', 'twiki')
        ssh_host: Optional SSH host for proxied fetching
        auth_type: Authentication type (tequila, session, keycloak, basic, etc.)
        credential_service: Keyring service for credentials
        async_wiki_client: Optional shared AsyncMediaWikiClient for Tequila auth
        keycloak_client: Optional shared httpx.AsyncClient for Keycloak auth
        basic_auth_client: Optional shared httpx.AsyncClient with HTTP Basic auth

    Returns:
        Number of chunks created
    """
    from imas_codex.discovery.wiki.pipeline import WikiIngestionPipeline
    from imas_codex.discovery.wiki.scraper import WikiPage

    # Extract page name from URL or page_id
    page_name = page_id.split(":", 1)[1] if ":" in page_id else page_id

    # Handle TWiki raw content (ssh:// URLs → read file via SSH, convert markup)
    if url and url.startswith("ssh://"):
        import asyncio

        from imas_codex.discovery.wiki.adapters import fetch_twiki_raw_content
        from imas_codex.discovery.wiki.pipeline import twiki_markup_to_html

        # Parse ssh://host/path
        parts = url[len("ssh://") :]
        slash_idx = parts.index("/")
        raw_ssh_host = parts[:slash_idx]
        filepath = parts[slash_idx:]

        # Read raw TWiki markup via SSH (blocking I/O → thread pool)
        raw_markup = await asyncio.to_thread(
            fetch_twiki_raw_content, raw_ssh_host, filepath
        )
        if not raw_markup or len(raw_markup) < 50:
            logger.warning("Insufficient TWiki content for %s", page_id)
            return 0

        # Convert TWiki markup to minimal HTML
        html = twiki_markup_to_html(raw_markup)
    else:
        # Fetch HTML content with auth (standard HTTP path)
        html = await _fetch_html(
            url,
            ssh_host,
            auth_type=auth_type,
            credential_service=credential_service,
            async_wiki_client=async_wiki_client,
            keycloak_client=keycloak_client,
            basic_auth_client=basic_auth_client,
        )
    if not html or len(html) < 100:
        logger.warning("Insufficient content for %s", page_id)
        return 0

    # Extract title from HTML
    import re

    title_match = re.search(r"<title>([^<]+)</title>", html)
    title = title_match.group(1) if title_match else page_name

    # Clean up title (remove wiki suffix)
    for suffix in [" - SPCwiki", " - Wikipedia", " - Confluence"]:
        if title.endswith(suffix):
            title = title[: -len(suffix)]

    # Create WikiPage object (fields from dataclass in scraper.py)
    page = WikiPage(
        url=url,
        title=title,
        content_html=html,
        content_text="",  # Will be extracted by pipeline
        sections={},
        mdsplus_paths=[],  # Will be extracted by pipeline
        imas_paths=[],
        units=[],
        conventions=[],
    )

    # Use the ingestion pipeline
    pipeline = WikiIngestionPipeline(
        facility_id=facility,
        use_rich=False,  # No progress display in worker
    )

    try:
        stats = await pipeline.ingest_page(page)
        return stats.get("chunks", 0)
    except Exception as e:
        logger.warning("Failed to ingest %s: %s", page_id, e)
        return 0


# =============================================================================
# Main Entry Point
# =============================================================================


async def run_parallel_wiki_discovery(
    facility: str,
    site_type: str,
    base_url: str,
    portal_page: str,
    ssh_host: str | None = None,
    auth_type: str | None = None,
    credential_service: str | None = None,
    cost_limit: float = 10.0,
    page_limit: int | None = None,
    max_depth: int | None = None,
    focus: str | None = None,
    num_scan_workers: int = 1,
    num_score_workers: int = 1,
    num_ingest_workers: int = 4,
    scan_only: bool = False,
    score_only: bool = False,
    bulk_discover: bool = True,
    ingest_artifacts: bool = True,
    skip_reset: bool = False,
    on_scan_progress: Callable | None = None,
    on_score_progress: Callable | None = None,
    on_ingest_progress: Callable | None = None,
    on_artifact_progress: Callable | None = None,
    on_artifact_score_progress: Callable | None = None,
    on_image_progress: Callable | None = None,
    on_worker_status: Callable[[SupervisedWorkerGroup], None] | None = None,
) -> dict[str, Any]:
    """Run parallel wiki discovery with async workers.

    Args:
        ssh_host: SSH host for proxied access (for non-public wikis)
        auth_type: Authentication type (tequila, session, ssh_proxy, or None)
        credential_service: Keyring service name for credentials
        bulk_discover: If True (default), use Special:AllPages for fast
            discovery of all pages upfront. This is 100-300x faster than
            crawling links page-by-page. Only works for MediaWiki sites.
        ingest_artifacts: If True (default), start artifact worker to
            download and ingest artifacts discovered during scanning.
            Artifacts go through score→ingest pipeline:
            discovered → scored (LLM) → ingested (embed)
        skip_reset: If True, skip reset_transient_pages on entry. Set this
            when called from multi-site loop where the CLI has already
            reset orphans once at startup.
        on_artifact_progress: Progress callback for artifact ingest worker.
        on_artifact_score_progress: Progress callback for artifact score worker.
        on_worker_status: Callback for worker status changes. Called with
            SupervisedWorkerGroup for live status display.

    Returns:
        Dict with discovery statistics
    """
    start_time = time.time()

    # Reset orphans from previous runs (skip when called from multi-site
    # loop since the CLI already resets once at the start)
    if not skip_reset:
        reset_transient_pages(facility)

    # Initialize state with auth info
    state = WikiDiscoveryState(
        facility=facility,
        site_type=site_type,
        base_url=base_url,
        portal_page=portal_page,
        ssh_host=ssh_host,
        auth_type=auth_type,
        credential_service=credential_service,
        cost_limit=cost_limit,
        page_limit=page_limit,
        max_depth=max_depth,
        focus=focus,
    )

    # Pre-warm SSH ControlMaster if using SSH transport.
    # Ensures the master connection is established before concurrent workers
    # spawn multiple ssh subprocesses (which would race to create it).
    if ssh_host:
        logger.info("Pre-warming SSH ControlMaster to %s...", ssh_host)
        try:
            await asyncio.to_thread(
                subprocess.run,
                ["ssh", "-O", "check", ssh_host],
                capture_output=True,
                timeout=10,
            )
            logger.info("SSH ControlMaster active for %s", ssh_host)
        except Exception:
            # No master yet — establish one with a quick command
            try:
                await asyncio.to_thread(
                    subprocess.run,
                    ["ssh", ssh_host, "true"],
                    capture_output=True,
                    timeout=30,
                )
                logger.info("SSH ControlMaster established for %s", ssh_host)
            except Exception as e:
                logger.warning("Failed to pre-warm SSH to %s: %s", ssh_host, e)

    # Create worker group for status tracking
    worker_group = SupervisedWorkerGroup()

    # Bulk discovery: use adapter-based APIs to find all pages instantly
    # This replaces the slow scan phase (crawling links page-by-page)
    bulk_discovered = 0
    if bulk_discover and not score_only:

        def bulk_progress(msg, _stats):
            if on_scan_progress:
                on_scan_progress(f"bulk: {msg}", state.scan_stats)

        # Extract data_path for twiki_raw from base_url
        data_path = None
        if site_type == "twiki_raw":
            data_path = base_url
            if data_path.startswith("ssh://"):
                data_path = data_path.split("/", 3)[-1]
                if not data_path.startswith("/"):
                    data_path = "/" + data_path

        # Compute exclude prefixes for static_html
        exclude_prefixes = None
        if site_type == "static_html":
            exclude_prefixes = _get_exclude_prefixes(facility, base_url)

        bulk_discovered = await asyncio.to_thread(
            bulk_discover_pages,
            facility,
            site_type,
            base_url,
            ssh_host=ssh_host,
            auth_type=state.auth_type,
            credential_service=state.credential_service,
            data_path=data_path,
            exclude_prefixes=exclude_prefixes,
            on_progress=bulk_progress,
        )

        if bulk_discovered:
            logger.info(f"Bulk discovery found {bulk_discovered} pages")
        state.scan_stats.processed = bulk_discovered

    # Bulk artifact discovery: use platform API to find all artifacts
    # This is much faster than scanning each page for links
    bulk_artifacts_discovered = 0
    if bulk_discover and ingest_artifacts and not score_only:
        logger.info("Starting bulk artifact discovery via API...")

        def artifact_progress(msg, _stats):
            if on_artifact_progress:
                on_artifact_progress(f"bulk: {msg}", state.artifact_stats)

        # Bulk discovery uses SSH when available, no async wiki client support yet
        bulk_artifacts_discovered = await asyncio.to_thread(
            bulk_discover_artifacts,
            facility,
            base_url,
            site_type,
            ssh_host,
            None,  # wiki_client - removed, use SSH path
            state.credential_service,
            "vpn" if ssh_host else "direct",
            artifact_progress,
        )

        if bulk_artifacts_discovered:
            logger.info(f"Bulk discovery found {bulk_artifacts_discovered} artifacts")
            state.artifact_stats.processed = bulk_artifacts_discovered

    # Create portal page if not exists (may already exist from bulk discovery)
    _seed_portal_page(facility, portal_page, base_url, site_type)

    # Start supervised workers with status tracking
    # Note: scan_worker was removed — bulk_discover_pages() replaces link-crawling.
    # num_scan_workers and scan_only params are kept for CLI backwards compat.

    if not score_only:
        for i in range(num_score_workers):
            worker_name = f"score_worker_{i}"
            status = worker_group.create_status(worker_name)
            worker_group.add_task(
                asyncio.create_task(
                    supervised_worker(
                        score_worker,
                        worker_name,
                        state,
                        state.should_stop_scoring,
                        on_progress=on_score_progress,
                        status_tracker=status,
                    )
                )
            )

        for i in range(num_ingest_workers):
            worker_name = f"ingest_worker_{i}"
            ingest_status = worker_group.create_status(worker_name)
            worker_group.add_task(
                asyncio.create_task(
                    supervised_worker(
                        ingest_worker,
                        worker_name,
                        state,
                        state.should_stop_ingesting,
                        on_progress=on_ingest_progress,
                        status_tracker=ingest_status,
                    )
                )
            )

        # Artifact workers: score→ingest pipeline for discovered artifacts
        if ingest_artifacts:
            # Artifact score worker: LLM scoring with text preview extraction
            artifact_score_status = worker_group.create_status("artifact_score_worker")
            worker_group.add_task(
                asyncio.create_task(
                    supervised_worker(
                        artifact_score_worker,
                        "artifact_score_worker",
                        state,
                        state.should_stop_artifact_scoring,
                        on_progress=on_artifact_score_progress,
                        status_tracker=artifact_score_status,
                    )
                )
            )

            # Artifact ingest worker: download and embed scored artifacts
            artifact_ingest_status = worker_group.create_status("artifact_worker")
            worker_group.add_task(
                asyncio.create_task(
                    supervised_worker(
                        artifact_worker,
                        "artifact_worker",
                        state,
                        state.should_stop_artifact_worker,
                        on_progress=on_artifact_progress,
                        status_tracker=artifact_ingest_status,
                    )
                )
            )

        # Image score worker: VLM caption + scoring for ingested images
        image_score_status = worker_group.create_status("image_score_worker")
        worker_group.add_task(
            asyncio.create_task(
                supervised_worker(
                    image_score_worker,
                    "image_score_worker",
                    state,
                    state.should_stop_image_scoring,
                    on_progress=on_image_progress,
                    status_tracker=image_score_status,
                )
            )
        )

    logger.info(
        f"Started {worker_group.get_active_count()} workers: "
        f"score_only={score_only}, scan_only={scan_only}, "
        f"ingest_artifacts={ingest_artifacts}"
    )

    # Send initial worker status update immediately so display shows workers
    if on_worker_status:
        try:
            on_worker_status(worker_group)
        except Exception as e:
            logger.warning("Initial worker status callback failed: %s", e)

    # Wait for termination condition with periodic orphan recovery
    orphan_check_interval = 60  # Check every 60 seconds
    last_orphan_check = time.time()
    status_update_interval = 0.5  # Update status every 0.5 seconds
    last_status_update = time.time()

    while not state.should_stop():
        await asyncio.sleep(0.25)

        # Update worker status for display
        if (
            on_worker_status
            and time.time() - last_status_update > status_update_interval
        ):
            try:
                on_worker_status(worker_group)
            except Exception as e:
                logger.warning("Worker status callback failed: %s", e)
            last_status_update = time.time()

        # Periodically release orphaned claims (from crashed workers)
        if time.time() - last_orphan_check > orphan_check_interval:
            try:
                released = release_orphaned_claims(facility)
                if released.get("released_pages", 0) or released.get(
                    "released_artifacts", 0
                ):
                    logger.info(
                        "Recovered %d orphaned page claims, %d artifact claims",
                        released.get("released_pages", 0),
                        released.get("released_artifacts", 0),
                    )
            except Exception as e:
                logger.debug("Orphan recovery check failed: %s", e)
            last_orphan_check = time.time()

    # Stop workers
    state.stop_requested = True
    await worker_group.cancel_all()

    # Clean up async wiki client
    await state.close_async_wiki_client()
    await state.close_keycloak_client()
    await state.close_basic_auth_client()

    elapsed = time.time() - start_time

    return {
        "scanned": state.scan_stats.processed,
        "scored": state.score_stats.processed,
        "ingested": state.ingest_stats.processed,
        "artifacts": state.artifact_stats.processed,
        "images_scored": state.image_stats.processed,
        "cost": state.total_cost,
        "elapsed_seconds": elapsed,
        "scan_rate": state.scan_stats.rate,
        "score_rate": state.score_stats.rate,
    }


def _get_exclude_prefixes(facility: str, current_base_url: str) -> list[str]:
    """Get URL path prefixes to exclude from crawling.

    Combines two sources of exclusions:
    1. Paths of OTHER wiki_sites on the same origin (auto-computed)
    2. Explicit exclude_prefixes from the site config (e.g., vendor docs)

    For example, jt60sa has:
      - https://nakasvr23.iferc.org/twiki_html (twiki_static)
      - https://nakasvr23.iferc.org (static_html with exclude_prefixes)

    When crawling the root site, /twiki_html is auto-excluded, and
    /matlab, /idl, etc. are excluded via config.

    Args:
        facility: Facility ID
        current_base_url: Base URL of the site being crawled

    Returns:
        List of URL path prefixes to exclude
    """
    import urllib.parse as urlparse

    try:
        from imas_codex.discovery.base.facility import get_facility

        config = get_facility(facility)
        wiki_sites = config.get("wiki_sites", [])
    except Exception:
        return []

    current_parsed = urlparse.urlparse(current_base_url.rstrip("/"))
    current_origin = f"{current_parsed.scheme}://{current_parsed.netloc}"

    prefixes = []

    # Find the current site config to get explicit exclude_prefixes
    current_site_config = None
    for site in wiki_sites:
        site_url = site.get("url", "").rstrip("/")
        if site_url == current_base_url.rstrip("/"):
            current_site_config = site
            break

    # Add explicit exclude_prefixes from site config (vendor docs, etc.)
    if current_site_config:
        explicit_excludes = current_site_config.get("exclude_prefixes", [])
        prefixes.extend(explicit_excludes)

    # Add paths of other wiki_sites on same origin
    for site in wiki_sites:
        site_url = site.get("url", "").rstrip("/")
        if not site_url or site_url == current_base_url.rstrip("/"):
            continue

        site_parsed = urlparse.urlparse(site_url)
        site_origin = f"{site_parsed.scheme}://{site_parsed.netloc}"

        # Same origin, different path → exclude that path
        if site_origin == current_origin and site_parsed.path:
            prefixes.append(site_parsed.path)

    if prefixes:
        logger.info(
            "Excluding paths from crawl: %s (belong to other wiki_sites)", prefixes
        )

    return prefixes


def _seed_portal_page(
    facility: str,
    portal_page: str,
    base_url: str,
    site_type: str,
) -> None:
    """Create the portal page as initial seed if it doesn't exist."""
    from imas_codex.discovery.wiki.scraper import canonical_page_id

    page_id = canonical_page_id(portal_page, facility)

    # Build URL based on site type
    if site_type == "twiki":
        if "/" not in portal_page:
            portal_page = f"Main/{portal_page}"
        url = f"{base_url}/bin/view/{portal_page}"
    elif site_type == "confluence":
        # portal_page is a space key (e.g. "IMP"), not a numeric page ID
        url = f"{base_url}/spaces/{portal_page}/overview"
    else:
        url = f"{base_url}/{urllib.parse.quote(portal_page, safe='/')}"

    with GraphClient() as gc:
        gc.query(
            """
            MERGE (wp:WikiPage {id: $id})
            ON CREATE SET wp.title = $title,
                          wp.url = $url,
                          wp.facility_id = $facility,
                          wp.status = $scanned,
                          wp.link_depth = 0,
                          wp.discovered_at = datetime()
            """,
            id=page_id,
            title=portal_page,
            url=url,
            facility=facility,
            scanned=WikiPageStatus.scanned.value,
        )


# =============================================================================
# Stats Query
# =============================================================================


def get_wiki_discovery_stats(facility: str) -> dict[str, int | float]:
    """Get wiki discovery statistics from graph.

    Returns stats needed for progress display:
    - total: Total wiki pages for this facility
    - scanned: Pages with status=scanned (awaiting scoring)
    - scored: Pages with status=scored (awaiting ingest or skipped)
    - ingested: Pages with status=ingested (final state)
    - pending_score: Same as scanned count (for progress display)
    - pending_ingest: Scored pages with score >= 0.5 awaiting ingestion
    - accumulated_cost: Total score_cost from all scored/ingested pages
    - skipped: Pages with status=skipped or score < 0.5
    """
    with GraphClient() as gc:
        # Get status counts
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WITH wp.status AS status
            RETURN status, count(*) AS count
            """,
            facility=facility,
        )

        stats: dict[str, int | float] = {
            "total": 0,
            "discovered": 0,
            "scanning": 0,
            "scanned": 0,
            "scoring": 0,
            "scored": 0,
            "ingesting": 0,
            "ingested": 0,
            "skipped": 0,
            "failed": 0,
        }

        for r in result:
            status = r["status"]
            count = r["count"]
            if status in stats:
                stats[status] = count
            stats["total"] += count

        # Get pending ingest count (scored pages with score >= 0.5)
        ingest_result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE wp.status = $scored AND wp.score >= 0.5
            RETURN count(wp) AS pending_ingest
            """,
            facility=facility,
            scored=WikiPageStatus.scored.value,
        )
        stats["pending_score"] = stats["scanned"]
        stats["pending_ingest"] = (
            ingest_result[0]["pending_ingest"] if ingest_result else 0
        )

        # Get accumulated cost from all pages
        cost_result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility})
            WHERE wp.score_cost IS NOT NULL
            RETURN sum(wp.score_cost) AS total_cost
            """,
            facility=facility,
        )
        stats["accumulated_cost"] = (
            cost_result[0]["total_cost"]
            if cost_result and cost_result[0]["total_cost"]
            else 0.0
        )

        # Add artifact stats (total + pending counts per worker)
        artifact_result = gc.query(
            """
            MATCH (wa:WikiArtifact {facility_id: $facility})
            WITH wa.status AS status, wa.score AS score,
                 wa.artifact_type AS atype
            RETURN status, score, atype, count(*) AS cnt
            """,
            facility=facility,
        )
        total_artifacts = 0
        pending_artifact_score = 0
        pending_artifact_ingest = 0
        supported = list(SUPPORTED_ARTIFACT_TYPES)
        for r in artifact_result:
            total_artifacts += r["cnt"]
            st = r["status"]
            atype = r["atype"]
            if st == WikiArtifactStatus.discovered.value and atype in supported:
                pending_artifact_score += r["cnt"]
            elif (
                st == WikiArtifactStatus.scored.value
                and atype in supported
                and r["score"] is not None
                and r["score"] >= 0.5
            ):
                pending_artifact_ingest += r["cnt"]

        # Count ingested artifacts
        artifacts_ingested = 0
        artifacts_scored = 0
        for r in artifact_result:
            st = r["status"]
            if st == WikiArtifactStatus.ingested.value:
                artifacts_ingested += r["cnt"]
            elif st == WikiArtifactStatus.scored.value:
                artifacts_scored += r["cnt"]

        stats["total_artifacts"] = total_artifacts
        stats["artifacts_ingested"] = artifacts_ingested
        stats["artifacts_scored"] = artifacts_scored
        stats["pending_artifact_score"] = pending_artifact_score
        stats["pending_artifact_ingest"] = pending_artifact_ingest

        return stats
