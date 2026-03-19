"""Scoring functions for wiki discovery.

LLM-based scoring for wiki pages, documents, and images:
- Page scoring: content preview + LLM structured output
- Document scoring: text extraction + LLM structured output
- Image scoring: VLM captioning + structured output
- HTML fetching: SSH proxy, Tequila, Keycloak, HTTP Basic auth
- Heuristic fallback scoring for non-LLM path

All scoring functions return (results, cost) tuples where cost
is the actual LLM cost from OpenRouter API.
"""

from __future__ import annotations

import asyncio
import io
import logging
import shlex
import subprocess
import tarfile
import time
from typing import Any

from imas_codex.discovery.base.scoring import CONTENT_SCORE_DIMENSIONS
from imas_codex.graph import GraphClient

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Dynamic calibration (same architecture as code/scorer.py)
# ---------------------------------------------------------------------------

_wiki_calibration_cache: dict[str, tuple[float, dict]] = {}
_WIKI_CALIBRATION_TTL_SECONDS = (
    300.0  # 5 minutes — matches LLM provider ephemeral cache TTL
)


def sample_wiki_page_calibration(
    facility: str | None = None,
    per_level: int = 2,
) -> dict[str, dict[str, list[dict[str, Any]]]]:
    """Sample calibration examples from scored WikiPage nodes.

    Same pattern as code calibration: deterministic ordering,
    facility preference, in-process TTL cache.

    Returns:
        Nested dict: dimension -> level -> list of examples.
        Each example: path (title), facility, score, purpose, description.
    """
    cache_key = f"page:{facility}:{per_level}"
    now = time.monotonic()

    if cache_key in _wiki_calibration_cache:
        cached_time, cached_data = _wiki_calibration_cache[cache_key]
        if (now - cached_time) < _WIKI_CALIBRATION_TTL_SECONDS:
            return cached_data

    samples = _fetch_wiki_calibration("WikiPage", facility, per_level)
    _wiki_calibration_cache[cache_key] = (now, samples)
    return samples


def sample_wiki_document_calibration(
    facility: str | None = None,
    per_level: int = 2,
) -> dict[str, dict[str, list[dict[str, Any]]]]:
    """Sample calibration examples from scored Document nodes.

    Same pattern as page calibration but queries Document nodes.

    Returns:
        Nested dict: dimension -> level -> list of examples.
        Each example: path (filename), facility, score, purpose, description.
    """
    cache_key = f"document:{facility}:{per_level}"
    now = time.monotonic()

    if cache_key in _wiki_calibration_cache:
        cached_time, cached_data = _wiki_calibration_cache[cache_key]
        if (now - cached_time) < _WIKI_CALIBRATION_TTL_SECONDS:
            return cached_data

    samples = _fetch_wiki_calibration("Document", facility, per_level)
    _wiki_calibration_cache[cache_key] = (now, samples)
    return samples


def _fetch_wiki_calibration(
    node_label: str,
    facility: str | None,
    per_level: int,
) -> dict[str, dict[str, list[dict[str, Any]]]]:
    """Fetch dimension calibration from WikiPage or Document nodes (uncached).

    Queries scored nodes at 5 score levels per dimension using deterministic
    ordering: same-facility preference → closest to bucket midpoint → id.
    """
    if node_label == "WikiPage":
        status_clause = "n.status IN ['scored', 'ingested']"
        title_prop = "n.title"
        desc_prop = "n.description"
    else:
        status_clause = "n.status IN ['scored', 'ingested']"
        title_prop = "n.filename"
        desc_prop = "n.description"

    buckets: list[tuple[str, float, float]] = [
        ("lowest", 0.0, 0.15),
        ("low", 0.10, 0.30),
        ("medium", 0.40, 0.60),
        ("high", 0.70, 0.90),
        ("highest", 0.90, 1.01),
    ]

    samples: dict[str, dict[str, list[dict[str, Any]]]] = {}

    with GraphClient() as gc:
        for dim in CONTENT_SCORE_DIMENSIONS:
            samples[dim] = {}

            for level_name, min_score, max_score in buckets:
                target = (min_score + max_score) / 2
                result = gc.query(
                    f"""
                    MATCH (n:{node_label})
                    WHERE {status_clause}
                        AND n.{dim} >= $min_score
                        AND n.{dim} < $max_score
                        AND n.{dim} IS NOT NULL
                    RETURN {title_prop} AS path,
                           n.facility_id AS facility,
                           n.{dim} AS score,
                           {desc_prop} AS description
                    ORDER BY
                        CASE WHEN n.facility_id = $facility
                             THEN 0 ELSE 1 END,
                        abs(n.{dim} - $target) ASC,
                        n.id ASC
                    LIMIT $limit
                    """,
                    min_score=min_score,
                    max_score=max_score,
                    target=target,
                    facility=facility or "",
                    limit=per_level,
                )

                samples[dim][level_name] = [
                    {
                        "path": r["path"] or "",
                        "facility": r["facility"],
                        "score": round(r["score"], 2),
                        "purpose": "wiki page"
                        if node_label == "WikiPage"
                        else "document",
                        "description": r["description"] or "",
                    }
                    for r in result
                ]

    return samples


def batch_ssh_read_files(
    ssh_urls: list[str],
) -> dict[str, str]:
    """Read multiple raw files in a single SSH call.

    Groups URLs by SSH host, then reads all files for each host
    in one SSH command with server-side concatenation. Returns a
    dict mapping ssh:// URL → file content.

    This amortises the ~3-5 second SSH connection overhead over
    the entire batch instead of paying it per file.

    Args:
        ssh_urls: List of ssh://host/path URLs

    Returns:
        Dict mapping URL → raw file content (empty string on failure)
    """
    if not ssh_urls:
        return {}

    # Group by SSH host
    by_host: dict[str, list[tuple[str, str]]] = {}
    for url in ssh_urls:
        parts = url[len("ssh://") :]
        slash_idx = parts.index("/")
        host = parts[:slash_idx]
        filepath = parts[slash_idx:]
        by_host.setdefault(host, []).append((url, filepath))

    results: dict[str, str] = {}

    for host, url_paths in by_host.items():
        # Use a tar stream instead of textual separators.
        # This avoids delimiter collisions and is insensitive to file content.
        quoted_paths = " ".join(shlex.quote(path) for _url, path in url_paths)
        remote_cmd = f"tar -cf - -- {quoted_paths} 2>/dev/null"

        try:
            result = subprocess.run(
                ["ssh", "-o", "ClearAllForwardings=yes", host, remote_cmd],
                capture_output=True,
                timeout=max(30, len(url_paths) * 2),
            )
            if result.returncode not in (0, 1):
                logger.warning(
                    "Batch SSH read failed for %s (%d files): exit %d",
                    host,
                    len(url_paths),
                    result.returncode,
                )
                for url, _ in url_paths:
                    results[url] = ""
                continue

            # Parse tar archive members by path and decode as UTF-8.
            url_by_path = {path: url for url, path in url_paths}
            stream = io.BytesIO(result.stdout)
            with tarfile.open(fileobj=stream, mode="r:") as archive:
                for member in archive:
                    if not member.isfile():
                        continue

                    member_path = "/" + member.name.lstrip("/")
                    url = url_by_path.get(member_path)
                    if not url:
                        continue

                    extracted = archive.extractfile(member)
                    if extracted is None:
                        continue
                    data = extracted.read()
                    results[url] = data.decode("utf-8", errors="replace")

        except subprocess.TimeoutExpired:
            logger.warning(
                "Batch SSH read timed out for %s (%d files)",
                host,
                len(url_paths),
            )
            for url, _ in url_paths:
                results[url] = ""
        except tarfile.ReadError as e:
            logger.warning("Batch SSH tar parse error for %s: %s", host, e)
            for url, _ in url_paths:
                results[url] = ""
        except Exception as e:
            logger.warning("Batch SSH read error for %s: %s", host, e)
            for url, _ in url_paths:
                results[url] = ""

    # Fill missing URLs with empty string
    for url in ssh_urls:
        if url not in results:
            results[url] = ""

    return results


async def _extract_document_preview(
    url: str,
    document_type: str,
    facility: str,
    max_chars: int = 1500,
    session: Any = None,
) -> str:
    """Extract a text preview from an document for LLM scoring.

    Downloads the document content and extracts text from the first
    portion. This is a lightweight extraction for scoring purposes only.

    For PDFs: extracts text from first few pages.
    For documents: extracts text paragraphs.
    For presentations: extracts slide text.
    For notebooks: extracts cell content.

    Args:
        url: Document download URL
        document_type: Type of document (pdf, docx, pptx, xlsx, ipynb)
        facility: Facility ID (for SSH proxy)
        max_chars: Maximum characters to extract
        session: Optional authenticated requests.Session (bypasses SSH)

    Returns:
        Extracted text preview or empty string on failure
    """
    from imas_codex.discovery.wiki.pipeline import fetch_document_content

    try:
        _, content_bytes = await fetch_document_content(
            url, facility=facility, session=session
        )
    except Exception as e:
        logger.debug("Failed to download %s: %s", url, e)
        return ""

    try:
        text = _extract_text_from_bytes(content_bytes, document_type)
        return text[:max_chars] if text else ""
    except Exception as e:
        logger.debug("Failed to extract text from %s: %s", url, e)
        return ""


def _extract_text_from_bytes(content_bytes: bytes, document_type: str) -> str:
    """Extract text from document bytes based on type.

    Lightweight extraction for scoring preview. Does not use LlamaIndex
    to avoid heavyweight dependencies in the scoring path.

    Uses semantic document type names matching DocumentType enum values.
    """
    import io

    at = document_type.lower()

    if at == "pdf":
        # Extract text from first few pages
        if b"%PDF" not in content_bytes[:1024]:
            return ""
        try:
            import logging as _logging

            import pypdf

            # Suppress pypdf's verbose warnings and errors.
            # pypdf._cmap uses logger_error() for benign "Advanced
            # encoding" messages at ERROR level — CRITICAL silences them.
            pypdf_logger = _logging.getLogger("pypdf")
            original_level = pypdf_logger.level
            pypdf_logger.setLevel(_logging.CRITICAL)
            try:
                reader = pypdf.PdfReader(io.BytesIO(content_bytes))
                text_parts = []
                for page in reader.pages[:5]:  # First 5 pages
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
            finally:
                pypdf_logger.setLevel(original_level)
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    elif at == "text_document":
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(io.BytesIO(content_bytes))
            paragraphs = [p.text for p in doc.paragraphs[:50] if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception:
            return ""

    elif at == "presentation":
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content_bytes))
            text_parts = []
            for slide in list(prs.slides)[:10]:  # First 10 slides
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
            return "\n".join(text_parts)
        except Exception:
            return ""

    elif at == "spreadsheet":
        try:
            from imas_codex.discovery.wiki.excel import extract_excel_preview

            return extract_excel_preview(content_bytes)
        except Exception:
            return ""

    elif at == "notebook":
        try:
            import json

            nb = json.loads(content_bytes.decode("utf-8"))
            cells = nb.get("cells", [])
            text_parts = []
            for cell in cells[:20]:  # First 20 cells
                source = "".join(cell.get("source", []))
                if source.strip():
                    text_parts.append(source)
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    elif at == "json":
        try:
            import json

            text = content_bytes.decode("utf-8", errors="replace")
            # Validate it's parseable JSON
            data = json.loads(text)
            # For large JSON, extract a structured preview
            preview = json.dumps(data, indent=2, ensure_ascii=False, default=str)
            return preview[:5000]  # Cap preview at 5000 chars
        except Exception:
            return ""

    return ""


async def _score_documents_batch(
    documents: list[dict[str, Any]],
    model: str,
    focus: str | None = None,
    data_access_patterns: dict[str, Any] | None = None,
    facility: str | None = None,
) -> tuple[list[dict[str, Any]], float]:
    """Score a batch of documents using LLM with structured output.

    Uses acall_llm_structured with DocumentScoreBatch Pydantic model for
    structured output. Content-based scoring with per-dimension scores.

    Args:
        documents: List of document dicts with id, filename, preview_text, etc.
        model: Model identifier from get_model()
        focus: Optional focus area for scoring
        data_access_patterns: Optional facility-specific data access patterns
            from facility config. When provided, injected into the prompt
            template so the LLM can boost documents matching facility tools/APIs.

    Returns:
        (results, cost) tuple where cost is actual LLM cost from OpenRouter.
    """
    from imas_codex.discovery.base.llm import acall_llm_structured
    from imas_codex.discovery.wiki.models import (
        DocumentScoreBatch,
        grounded_document_score,
    )
    from imas_codex.llm.prompt_loader import render_prompt

    # Build system prompt using document-scorer template
    context: dict[str, Any] = {}
    if focus:
        context["focus"] = focus
    if data_access_patterns:
        context["data_access_patterns"] = data_access_patterns

    # Inject calibration examples from scored Document nodes
    dimension_calibration = sample_wiki_document_calibration(
        facility=facility, per_level=2
    )
    has_calibration = any(
        any(examples for examples in levels.values())
        for levels in dimension_calibration.values()
    )
    if has_calibration:
        context["dimension_calibration"] = dimension_calibration

    system_prompt = render_prompt("wiki/document-scorer", context)

    # Build user prompt with document content
    lines = [
        f"Score these {len(documents)} wiki documents based on their content.",
        "(Use the preview text to assess value for the IMAS knowledge graph.)\n",
    ]

    for i, a in enumerate(documents, 1):
        lines.append(f"\n## Document {i}")
        lines.append(f"ID: {a['id']}")
        lines.append(f"Filename: {a.get('filename', 'Unknown')}")
        lines.append(f"Type: {a.get('document_type', 'unknown')}")

        if a.get("size_bytes"):
            size_mb = a["size_bytes"] / (1024 * 1024)
            lines.append(f"Size: {size_mb:.1f} MB")

        preview = a.get("preview_text", "")
        if preview:
            lines.append(f"Content Preview:\n{preview[:800]}")
        else:
            lines.append("Content Preview: (not available - score from filename/type)")

        if a.get("url"):
            lines.append(f"URL: {a['url']}")

    lines.append(
        "\n\nReturn results for each document in order. "
        "The response format is enforced by the schema."
    )

    user_prompt = "\n".join(lines)

    messages: list[dict[str, Any]] = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]

    batch, total_cost, _ = await acall_llm_structured(
        model=model,
        messages=messages,
        response_model=DocumentScoreBatch,
        max_tokens=16000,
        temperature=0.3,
        timeout=120,
        max_retries=5,
        retry_base_delay=4.0,
    )

    llm_results = batch.results

    # Convert to result dicts
    cost_per_document = total_cost / len(documents) if documents else 0.0
    results = []

    for r in llm_results[: len(documents)]:
        scores = {
            "score_data_documentation": r.score_data_documentation,
            "score_physics_content": r.score_physics_content,
            "score_code_documentation": r.score_code_documentation,
            "score_data_access": r.score_data_access,
            "score_calibration": r.score_calibration,
            "score_imas_relevance": r.score_imas_relevance,
        }

        combined_score = grounded_document_score(scores, r.document_purpose)

        # Find the matching document for filename
        matching = next((a for a in documents if a["id"] == r.id), {})

        results.append(
            {
                "id": r.id,
                "score_composite": combined_score,
                "document_purpose": r.document_purpose.value,
                "description": r.description,
                "reasoning": r.reasoning,
                "keywords": r.keywords[:5],
                "physics_domain": r.physics_domain.value if r.physics_domain else None,
                "should_ingest": r.should_ingest,
                "skip_reason": r.skip_reason or None,
                "score_data_documentation": r.score_data_documentation,
                "score_physics_content": r.score_physics_content,
                "score_code_documentation": r.score_code_documentation,
                "score_data_access": r.score_data_access,
                "score_calibration": r.score_calibration,
                "score_imas_relevance": r.score_imas_relevance,
                "score_cost": cost_per_document,
                # Pass through filename for display
                "filename": matching.get("filename", ""),
                "document_type": matching.get("document_type", ""),
            }
        )

    return results, total_cost


# =============================================================================


async def _score_images_batch(
    images: list[dict[str, Any]],
    model: str,
    focus: str | None = None,
    data_access_patterns: dict[str, Any] | None = None,
    facility_id: str | None = None,
) -> tuple[list[dict[str, Any]], float]:
    """Score a batch of images using VLM with structured output.

    Delegates to the shared implementation in discovery.base.image.
    """
    from imas_codex.discovery.base.image import score_images_batch

    return await score_images_batch(
        images, model, focus, data_access_patterns, facility_id=facility_id
    )


async def _fetch_html(
    url: str,
    ssh_host: str | None,
    auth_type: str | None = None,
    credential_service: str | None = None,
    async_wiki_client: Any = None,
    keycloak_client: Any = None,
    basic_auth_client: Any = None,
    confluence_client: Any = None,
) -> str:
    """Fetch HTML content from URL.

    Args:
        url: Page URL
        ssh_host: Optional SSH host for proxied fetching
        auth_type: Authentication type (tequila, session, keycloak, basic, etc.)
        credential_service: Keyring service for credentials
        async_wiki_client: Shared AsyncMediaWikiClient for Tequila auth
        keycloak_client: Shared httpx.AsyncClient for Keycloak auth
        basic_auth_client: Shared httpx.AsyncClient with HTTP Basic auth
        confluence_client: Shared ConfluenceClient for Confluence session auth

    Returns:
        HTML content string or empty string on error
    """

    def _ssh_fetch() -> str:
        """Fetch via SSH proxy."""
        cmd = f'curl -sk --noproxy "*" "{url}" 2>/dev/null'
        try:
            result = subprocess.run(
                ["ssh", ssh_host, cmd],
                capture_output=True,
                timeout=30,
            )
            if result.returncode == 0:
                # Wiki pages may be ISO-8859 encoded despite claiming UTF-8
                return result.stdout.decode("utf-8", errors="replace")
            return ""
        except Exception as e:
            logger.warning("SSH fetch failed for %s: %s", url, e)
            return ""

    def _ssh_file_read_and_convert(file_url: str) -> str:
        """Read a raw TWiki file via SSH and convert markup to HTML.

        Used for ssh:// URLs from the twiki_raw adapter. Reads the file
        directly from the filesystem (instant) and converts TWiki markup
        to HTML locally, avoiding the slow CGI path entirely.
        """
        from imas_codex.discovery.wiki.adapters import fetch_twiki_raw_content
        from imas_codex.discovery.wiki.pipeline import twiki_markup_to_html

        parts = file_url[len("ssh://") :]
        slash_idx = parts.index("/")
        raw_ssh_host = parts[:slash_idx]
        filepath = parts[slash_idx:]

        raw_markup = fetch_twiki_raw_content(raw_ssh_host, filepath)
        if not raw_markup or len(raw_markup) < 20:
            return ""
        return twiki_markup_to_html(raw_markup)

    async def _async_tequila_fetch() -> str:
        """Fetch with Tequila authentication using async client."""
        import urllib.parse as urlparse

        # Extract page name from URL
        page_name = url.split("/wiki/")[-1] if "/wiki/" in url else url.split("/")[-1]
        if "?" in page_name:
            parsed = urlparse.parse_qs(urlparse.urlparse(url).query)
            page_name = parsed.get("title", [page_name])[0]
        page_name = urlparse.unquote(page_name)

        # Use provided async client
        if async_wiki_client is not None:
            try:
                page = await async_wiki_client.get_page(page_name)
                if page:
                    return page.content_html
                return ""
            except Exception as e:
                logger.debug("Async client fetch failed for %s: %s", url, e)
                return ""

        # No client provided - create a new async client
        from imas_codex.discovery.wiki.mediawiki import AsyncMediaWikiClient

        base_url_local = (
            url.rsplit("/", 1)[0] if "/wiki/" in url else url.rsplit("/", 1)[0]
        )
        if "/wiki" in base_url_local:
            base_url_local = base_url_local.rsplit("/wiki", 1)[0] + "/wiki"

        async with AsyncMediaWikiClient(
            base_url=base_url_local,
            credential_service=credential_service or "tcv",
            verify_ssl=False,
        ) as client:
            if not await client.authenticate():
                logger.warning("Tequila auth failed for %s", url)
                return ""
            page = await client.get_page(page_name)
            return page.content_html if page else ""

    async def _async_keycloak_fetch() -> str:
        """Fetch with Keycloak auth using shared httpx.AsyncClient."""
        if keycloak_client is None:
            logger.warning("No Keycloak client available for %s", url)
            return ""
        try:
            response = await keycloak_client.get(url)
            if response.status_code == 200:
                return response.text
            logger.debug(
                "Keycloak fetch returned HTTP %d for %s", response.status_code, url
            )
            return ""
        except Exception as e:
            logger.debug("Keycloak fetch failed for %s: %s", url, e)
            return ""

    async def _async_basic_auth_fetch() -> str:
        """Fetch with HTTP Basic auth using shared httpx.AsyncClient."""
        if basic_auth_client is None:
            logger.warning("No HTTP Basic auth client available for %s", url)
            return ""
        try:
            response = await basic_auth_client.get(url)
            if response.status_code == 200:
                return response.text
            logger.debug(
                "Basic auth fetch returned HTTP %d for %s", response.status_code, url
            )
            return ""
        except Exception as e:
            logger.debug("Basic auth fetch failed for %s: %s", url, e)
            return ""

    async def _async_confluence_fetch() -> str:
        """Fetch via Confluence REST API using authenticated ConfluenceClient."""
        if confluence_client is None:
            logger.warning("No Confluence client available for %s", url)
            return ""
        try:
            import urllib.parse as urlparse

            # Extract pageId from viewpage.action URL
            parsed = urlparse.urlparse(url)
            params = urlparse.parse_qs(parsed.query)
            page_id = params.get("pageId", [None])[0]

            if not page_id:
                logger.debug("Cannot extract pageId from Confluence URL: %s", url)
                return ""

            page = await asyncio.to_thread(confluence_client.get_page_content, page_id)
            if page and page.content_html:
                return page.content_html
            return ""
        except Exception as e:
            logger.debug("Confluence REST API fetch failed for %s: %s", url, e)
            return ""

    # Handle ssh:// URLs (twiki_raw adapter pages) — read file via SSH
    if url.startswith("ssh://"):
        return await asyncio.to_thread(_ssh_file_read_and_convert, url)

    # Determine fetch strategy - prefer direct HTTP over SSH when credentials available
    if auth_type == "session" and confluence_client:
        return await _async_confluence_fetch()
    elif auth_type in ("tequila", "session"):
        return await _async_tequila_fetch()
    elif auth_type == "keycloak" and keycloak_client:
        return await _async_keycloak_fetch()
    elif auth_type == "basic" and basic_auth_client:
        return await _async_basic_auth_fetch()
    elif ssh_host:
        return await asyncio.to_thread(_ssh_fetch)
    else:
        # Direct HTTP fetch (no auth) - already async
        from imas_codex.discovery.wiki.prefetch import fetch_page_content

        html, error = await fetch_page_content(url)
        if html:
            return html
        if error:
            logger.debug("HTTP fetch failed for %s: %s", url, error)
        return ""


async def _fetch_and_summarize(
    url: str,
    ssh_host: str | None,
    auth_type: str | None = None,
    credential_service: str | None = None,
    max_chars: int = 2000,
    async_wiki_client: Any = None,
    keycloak_client: Any = None,
    basic_auth_client: Any = None,
    confluence_client: Any = None,
) -> str:
    """Fetch page content and extract text preview.

    No LLM is used here - prefetch extracts text deterministically.
    The summary is just cleaned text for the scorer to evaluate.

    Delegates HTML fetching to _fetch_html and applies text extraction.
    """
    from imas_codex.discovery.wiki.prefetch import extract_text_from_html

    html = await _fetch_html(
        url,
        ssh_host,
        auth_type=auth_type,
        credential_service=credential_service,
        async_wiki_client=async_wiki_client,
        keycloak_client=keycloak_client,
        basic_auth_client=basic_auth_client,
        confluence_client=confluence_client,
    )

    if html:
        return extract_text_from_html(html, max_chars=max_chars)
    return ""


async def _score_pages_batch(
    pages: list[dict[str, Any]],
    model: str,
    focus: str | None = None,
    data_access_patterns: dict[str, Any] | None = None,
    facility: str | None = None,
) -> tuple[list[dict[str, Any]], float]:
    """Score a batch of pages using LLM with structured output.

    Uses litellm.acompletion with WikiScoreBatch Pydantic model for
    structured output. Content-based scoring with per-dimension scores.

    The retry loop includes response parsing because JSON/validation
    errors from truncated responses are retryable — a fresh LLM call
    often returns valid output. Cost is accumulated across all attempts
    since API calls are billed regardless of parsing success.

    Args:
        pages: List of page dicts with id, title, summary, preview_text, etc.
        model: Model identifier from get_model()
        focus: Optional focus area for scoring
        data_access_patterns: Optional facility-specific data access patterns
            from facility config. When provided, injected into the prompt
            template so the LLM can boost pages matching facility tools/APIs.

    Returns:
        (results, cost) tuple where cost is actual LLM cost from OpenRouter.
    """
    from imas_codex.discovery.base.llm import acall_llm_structured
    from imas_codex.discovery.wiki.models import (
        WikiScoreBatch,
        grounded_wiki_score,
    )
    from imas_codex.llm.prompt_loader import render_prompt

    # Build system prompt using dynamic template with schema injection
    context: dict[str, Any] = {}
    if focus:
        context["focus"] = focus
    if data_access_patterns:
        context["data_access_patterns"] = data_access_patterns

    # Inject calibration examples from scored WikiPage nodes
    dimension_calibration = sample_wiki_page_calibration(facility=facility, per_level=2)
    has_calibration = any(
        any(examples for examples in levels.values())
        for levels in dimension_calibration.values()
    )
    if has_calibration:
        context["dimension_calibration"] = dimension_calibration

    system_prompt = render_prompt("wiki/scorer", context)

    # Build user prompt with page content (not graph metrics)
    lines = [
        f"Score these {len(pages)} wiki pages based on their content.",
        "(Use the preview text to infer value - graph metrics like in_degree are NOT indicators.)\n",
    ]

    for i, p in enumerate(pages, 1):
        lines.append(f"\n## Page {i}")
        lines.append(f"ID: {p['id']}")
        lines.append(f"Title: {p.get('title', 'Unknown')}")

        # Use preview_text for content-based scoring (preferred over summary)
        preview = p.get("preview_text") or p.get("summary") or ""
        if preview:
            lines.append(f"Preview: {preview[:800]}")

        # Include URL for context (Confluence vs MediaWiki structure hints)
        url = p.get("url")
        if url:
            lines.append(f"URL: {url}")

        # Include language hint for non-English content
        content_language = p.get("content_language")
        if content_language and content_language != "en":
            lines.append(
                f"Note: Content may be in {content_language}. "
                "Score based on technical content regardless of language."
            )

    lines.append(
        "\n\nReturn results for each page in order. "
        "The response format is enforced by the schema."
    )

    user_prompt = "\n".join(lines)

    # Shared retry+parse loop handles both API errors and JSON/validation
    # errors from truncated responses. Cost accumulated across retries.
    # Model-aware token limits + timeout applied automatically.
    batch, total_cost, _tokens = await acall_llm_structured(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        response_model=WikiScoreBatch,
        temperature=0.3,
    )

    llm_results = batch.results

    # Convert to result dicts, computing combined scores
    cost_per_page = total_cost / len(pages) if pages else 0.0
    results = []

    for r in llm_results[: len(pages)]:
        # Build per-dimension scores dict
        scores = {
            "score_data_documentation": r.score_data_documentation,
            "score_physics_content": r.score_physics_content,
            "score_code_documentation": r.score_code_documentation,
            "score_data_access": r.score_data_access,
            "score_calibration": r.score_calibration,
            "score_imas_relevance": r.score_imas_relevance,
        }

        # Compute combined score using grounded function
        combined_score = grounded_wiki_score(scores, r.page_purpose)

        results.append(
            {
                "id": r.id,
                "score_composite": combined_score,
                "purpose": r.page_purpose.value,
                "description": r.description,
                "reasoning": r.reasoning,
                "keywords": r.keywords[:5],
                "physics_domain": r.physics_domain.value if r.physics_domain else None,
                "should_ingest": r.should_ingest,
                "skip_reason": r.skip_reason or None,
                # Per-dimension scores
                "score_data_documentation": r.score_data_documentation,
                "score_physics_content": r.score_physics_content,
                "score_code_documentation": r.score_code_documentation,
                "score_data_access": r.score_data_access,
                "score_calibration": r.score_calibration,
                "score_imas_relevance": r.score_imas_relevance,
                # Legacy fields for compatibility
                "page_type": r.page_purpose.value,
                "is_physics": r.physics_domain is not None
                and r.physics_domain.value != "general",
                "score_cost": cost_per_page,
            }
        )

    return results, total_cost


def _score_pages_heuristic(
    pages: list[dict[str, Any]],
    data_access_patterns: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    """Heuristic fallback scoring based on keywords. Zero cost.

    When data_access_patterns is provided, facility-specific key_tools
    and code_import_patterns are added to the physics keywords for
    boosted matching.
    """
    # Build facility-specific keywords from data_access_patterns
    facility_keywords: list[str] = []
    if data_access_patterns:
        for tool in data_access_patterns.get("key_tools") or []:
            # Normalize: lowercase, strip parens/dots for matching
            kw = tool.lower().rstrip("(").split(".")[-1]
            if kw and kw not in facility_keywords:
                facility_keywords.append(kw)
        for pattern in data_access_patterns.get("code_import_patterns") or []:
            # Extract the key identifier from import patterns
            # e.g., "import ppf" -> "ppf", "ppfget(" -> "ppfget"
            kw = pattern.lower().rstrip("(").split()[-1].split(".")[-1]
            if kw and len(kw) > 2 and kw not in facility_keywords:
                facility_keywords.append(kw)

    results = []
    for page in pages:
        title = page.get("title", "").lower()
        summary = page.get("summary", "") or ""

        score = 0.5
        reasoning = "Default score"

        physics_keywords = [
            "thomson",
            "liuqe",
            "equilibrium",
            "mhd",
            "plasma",
            "diagnostic",
            "calibration",
            "signal",
            "node",
        ]
        low_value_keywords = [
            "meeting",
            "workshop",
            "todo",
            "draft",
            "notes",
            "personal",
            "test",
            "sandbox",
        ]

        for kw in physics_keywords:
            if kw in title or kw in summary.lower():
                score = min(score + 0.15, 1.0)
                reasoning = f"Contains physics keyword: {kw}"

        for kw in low_value_keywords:
            if kw in title:
                score = max(score - 0.2, 0.0)
                reasoning = f"Contains low-value keyword: {kw}"

        # Boost for facility-specific data access patterns
        for kw in facility_keywords:
            if kw in title or kw in summary.lower():
                score = min(score + 0.15, 1.0)
                reasoning = f"Contains facility data access keyword: {kw}"

        results.append(
            {
                "id": page["id"],
                "score_composite": score,
                "reasoning": reasoning,
                "page_type": "documentation",
                "is_physics": score >= 0.6,
            }
        )

    return results
