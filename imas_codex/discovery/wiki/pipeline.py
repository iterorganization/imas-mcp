"""Wiki ingestion pipeline using LlamaIndex.

Processes wiki pages and artifacts through a deterministic pipeline:
1. Scan: Discover pages/artifacts via link traversal (creates nodes with status='scanned')
2. Score: Agent evaluates interest_score (sets status='scored' or 'skipped')
3. Ingest: Fetch content, chunk, embed, and link to graph (sets status='ingested')

The pipeline is graph-driven and fully deterministic (no LLM calls during ingestion).
Entity extraction uses regex patterns for MDSplus paths, IMAS paths, units, conventions.

Example:
    # Step 1: Scan wiki (link extraction, no content fetch)
    # imas-codex wiki scan tcv

    # Step 2: Score pages with LLM agent
    # imas-codex wiki score tcv

    # Step 3: Ingest high-score pages (deterministic)
    pipeline = WikiIngestionPipeline(facility_id="tcv")
    stats = await pipeline.ingest_from_graph(min_interest_score=0.7)
    print(f"Created {stats['chunks']} chunks")
"""

from __future__ import annotations

import hashlib
import logging
import re
from collections.abc import Callable
from html.parser import HTMLParser
from typing import TYPE_CHECKING, TypedDict

from llama_index.core import Document
from llama_index.core.node_parser import SentenceSplitter

from imas_codex.graph import GraphClient
from imas_codex.settings import get_embedding_dimension

from .monitor import WikiProgressMonitor, set_current_monitor
from .scraper import WikiPage, fetch_wiki_page

if TYPE_CHECKING:
    from llama_index.core.embeddings import BaseEmbedding

logger = logging.getLogger(__name__)

# Batch size for UNWIND operations
BATCH_SIZE = 50

# Progress callback type: (current, total, message) -> None
ProgressCallback = Callable[[int, int, str], None]


class PageIngestionStats(TypedDict):
    """Statistics from ingesting a single wiki page."""

    chunks: int
    tree_nodes_linked: int
    imas_paths_linked: int
    conventions: int
    units: int
    content_preview: str
    mdsplus_paths: list[str]


# =============================================================================
# Graph Query Functions
# =============================================================================


def get_pending_wiki_pages(
    facility_id: str,
    limit: int | None = None,
    min_interest_score: float = 0.5,
) -> list[dict]:
    """Get wiki pages pending ingestion from the graph.

    Returns WikiPage nodes with status='scored' (passed agent evaluation),
    sorted by interest_score descending.

    Args:
        facility_id: Facility ID
        limit: Maximum pages to return (None for all)
        min_interest_score: Minimum interest score threshold

    Returns:
        List of page dicts with id, url, title, interest_score
    """
    # Build query with optional LIMIT clause
    query = """
        MATCH (wp:WikiPage {facility_id: $facility_id, status: 'scored'})
        WHERE wp.interest_score >= $min_score
        RETURN wp.id AS id, wp.url AS url, wp.title AS title,
               wp.interest_score AS interest_score
        ORDER BY wp.interest_score DESC, wp.discovered_at ASC
    """
    if limit is not None:
        query += f"LIMIT {limit}"

    with GraphClient() as gc:
        result = gc.query(
            query,
            facility_id=facility_id,
            min_score=min_interest_score,
        )
        return [dict(r) for r in result] if result else []


def get_wiki_queue_stats(facility_id: str) -> dict:
    """Get wiki page queue statistics by status.

    Args:
        facility_id: Facility ID

    Returns:
        Dict with status counts for pages and artifacts
    """
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility_id})
            RETURN wp.status AS status, count(*) AS count
            ORDER BY count DESC
            """,
            facility_id=facility_id,
        )

        # Initialize with zero counts for all expected statuses
        stats = {
            "scanned": 0,
            "scored": 0,
            "skipped": 0,
            "ingested": 0,
            "failed": 0,
            "stale": 0,
            "total": 0,
        }

        if result:
            for r in result:
                status = r["status"]
                count = r["count"]
                if status in stats:
                    stats[status] = count
                stats["total"] += count

        return stats


def get_pending_wiki_artifacts(
    facility_id: str,
    limit: int | None = None,
    min_interest_score: float = 0.5,
) -> list[dict]:
    """Get wiki artifacts pending ingestion from the graph.

    Returns WikiArtifact nodes with status='scored' (passed agent evaluation),
    sorted by interest_score descending.

    Args:
        facility_id: Facility ID
        limit: Maximum artifacts to return (None for all)
        min_interest_score: Minimum interest score threshold

    Returns:
        List of artifact dicts with id, url, filename, artifact_type, interest_score
    """
    # Build query with optional LIMIT clause
    query = """
        MATCH (wa:WikiArtifact {facility_id: $facility_id, status: 'scored'})
        WHERE wa.interest_score >= $min_score
        RETURN wa.id AS id, wa.url AS url, wa.filename AS filename,
               wa.artifact_type AS artifact_type,
               wa.interest_score AS interest_score
        ORDER BY wa.interest_score DESC
    """
    if limit is not None:
        query += f"LIMIT {limit}"

    with GraphClient() as gc:
        result = gc.query(
            query,
            facility_id=facility_id,
            min_score=min_interest_score,
        )
        return [dict(r) for r in result] if result else []


def mark_wiki_page_status(
    page_id: str,
    status: str,
    error: str | None = None,
) -> None:
    """Update the status of a wiki page.

    Args:
        page_id: WikiPage ID
        status: New status (scanned, scored, skipped, ingested, failed, stale)
        error: Error message if status is 'failed'
    """
    with GraphClient() as gc:
        gc.query(
            """
            MATCH (wp:WikiPage {id: $id})
            SET wp.status = $status,
                wp.error = $error
            """,
            id=page_id,
            status=status,
            error=error,
        )


def persist_chunks_batch(chunks: list[dict]) -> int:
    """Persist a batch of WikiChunk nodes using UNWIND for efficiency.

    Args:
        chunks: List of chunk dicts with required fields:
            - id, wiki_page_id, facility_id, content, embedding
            - mdsplus_paths, imas_paths, units, conventions (optional)

    Returns:
        Number of chunks persisted
    """
    if not chunks:
        return 0

    with GraphClient() as gc:
        gc.query(
            """
            UNWIND $chunks AS chunk
            MERGE (c:WikiChunk {id: chunk.id})
            SET c.wiki_page_id = chunk.wiki_page_id,
                c.facility_id = chunk.facility_id,
                c.content = chunk.content,
                c.embedding = chunk.embedding,
                c.mdsplus_paths_mentioned = chunk.mdsplus_paths,
                c.imas_paths_mentioned = chunk.imas_paths,
                c.units_mentioned = chunk.units,
                c.conventions_mentioned = chunk.conventions
            WITH c, chunk
            MATCH (p:WikiPage {id: chunk.wiki_page_id})
            MERGE (p)-[:HAS_CHUNK]->(c)
            """,
            chunks=chunks,
        )
        return len(chunks)


def link_chunks_to_entities(facility_id: str) -> dict[str, int]:
    """Create DOCUMENTS and MENTIONS relationships from chunk metadata.

    Uses batched queries to link WikiChunks to TreeNodes and IMASPaths
    based on the paths stored in chunk properties during ingestion.

    Args:
        facility_id: Facility to process

    Returns:
        Dict with counts: {tree_nodes_linked, imas_paths_linked}
    """
    stats = {"tree_nodes_linked": 0, "imas_paths_linked": 0}

    with GraphClient() as gc:
        # Link to TreeNodes via MDSplus paths
        result = gc.query(
            """
            MATCH (c:WikiChunk {facility_id: $facility_id})
            WHERE c.mdsplus_paths_mentioned IS NOT NULL
            UNWIND c.mdsplus_paths_mentioned AS mds_path
            MATCH (t:TreeNode)
            WHERE t.path = mds_path
               OR t.path ENDS WITH replace(mds_path, '\\\\', '')
               OR t.canonical_path = toUpper(mds_path)
            MERGE (c)-[:DOCUMENTS]->(t)
            RETURN count(*) AS linked
            """,
            facility_id=facility_id,
        )
        if result:
            stats["tree_nodes_linked"] = result[0]["linked"]

    return stats


class MediaWikiExtractor(HTMLParser):
    """MediaWiki-aware HTML to text converter.

    Targets the main content area (mw-parser-output) and skips
    navigation, sidebar, and footer elements.
    """

    # Tags that should be completely skipped
    SKIP_TAGS = {"script", "style", "nav", "footer", "header", "noscript"}

    # CSS classes that indicate navigation/sidebar content to skip
    SKIP_CLASSES = {
        "mw-navigation",
        "mw-sidebar",
        "mw-footer",
        "printfooter",
        "catlinks",
        "mw-editsection",
        "noprint",
        "toc",
        "navbox",
        "metadata",
        "mw-jump-link",
        "portal",
        "portlet",
        "p-personal",
        "p-navigation",
        "p-search",
        "p-tb",
        "p-coll-print_export",
    }

    # CSS classes that indicate main content
    CONTENT_CLASSES = {"mw-parser-output", "mw-body-content", "mw-content-text"}

    def __init__(self):
        super().__init__()
        self.text_parts: list[str] = []
        self.current_section: str = ""
        self.sections: dict[str, str] = {}
        self._skip_depth = 0
        self._content_depth = 0
        self._in_heading = False
        self._heading_text = ""
        self._in_table = False
        self._table_row: list[str] = []
        self._in_pre = False
        self._skipping_div = 0  # Track divs we're skipping

    def _get_class(self, attrs: list[tuple[str, str | None]]) -> set[str]:
        """Extract CSS classes from tag attributes."""
        for name, value in attrs:
            if name == "class" and value:
                return set(value.split())
        return set()

    def _get_id(self, attrs: list[tuple[str, str | None]]) -> str | None:
        """Extract id from tag attributes."""
        for name, value in attrs:
            if name == "id":
                return value
        return None

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        classes = self._get_class(attrs)
        element_id = self._get_id(attrs)

        # Check if this element should be skipped
        if tag in self.SKIP_TAGS:
            self._skip_depth += 1
            return

        if classes & self.SKIP_CLASSES:
            self._skip_depth += 1
            if tag == "div":
                self._skipping_div += 1
            return

        # Skip by ID patterns (common MediaWiki navigation)
        if element_id:
            # Skip navigation/sidebar divs by ID
            skip_ids = ("jump-to-nav", "siteSub", "contentSub")
            if element_id in skip_ids or element_id.startswith(
                ("mw-", "p-", "ca-", "n-")
            ):
                if element_id not in ("mw-content-text", "mw-body-content"):
                    self._skip_depth += 1
                    if tag == "div":
                        self._skipping_div += 1
                    return

        # Track when we enter content area
        if classes & self.CONTENT_CLASSES:
            self._content_depth += 1

        # Handle specific content tags
        if tag in ("h1", "h2", "h3", "h4", "h5"):
            self._in_heading = True
            self._heading_text = ""
        elif tag == "table":
            self._in_table = True
        elif tag == "tr":
            self._table_row = []
        elif tag == "pre":
            self._in_pre = True
            self.text_parts.append("\n```\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self.SKIP_TAGS:
            self._skip_depth = max(0, self._skip_depth - 1)
            return

        # Track end of skipped divs
        if tag == "div" and self._skipping_div > 0:
            self._skipping_div -= 1
            self._skip_depth = max(0, self._skip_depth - 1)
            return

        # Check for skip class divs - we don't track which specific div
        # so we just decrement if we're in a skip state
        if tag == "div" and self._skip_depth > 0:
            self._skip_depth = max(0, self._skip_depth - 1)
            return

        if tag in ("h1", "h2", "h3", "h4", "h5"):
            self._in_heading = False
            if self._heading_text.strip():
                # Save previous section
                if self.current_section and self.text_parts:
                    self.sections[self.current_section] = " ".join(self.text_parts)
                self.current_section = self._heading_text.strip()
                self.text_parts = []
                self.text_parts.append(f"\n## {self._heading_text.strip()}\n")
        elif tag == "table":
            self._in_table = False
            self.text_parts.append("\n")
        elif tag == "tr":
            if self._table_row:
                self.text_parts.append(" | ".join(self._table_row) + "\n")
            self._table_row = []
        elif tag in ("td", "th"):
            pass  # Cell content already added
        elif tag in ("p", "div", "li"):
            self.text_parts.append("\n")
        elif tag == "br":
            self.text_parts.append("\n")
        elif tag == "pre":
            self._in_pre = False
            self.text_parts.append("\n```\n")

    def handle_data(self, data: str) -> None:
        # Skip if we're in a skipped section
        if self._skip_depth > 0:
            return

        text = data.strip() if not self._in_pre else data

        if not text:
            return

        if self._in_heading:
            self._heading_text += text
        elif self._in_table and self._table_row is not None:
            # Collect table cell content
            self._table_row.append(text)
            self.text_parts.append(text + " ")
        else:
            self.text_parts.append(text)

    def get_result(self) -> tuple[str, dict[str, str]]:
        """Get extracted text and sections dict."""
        # Save final section
        if self.current_section and self.text_parts:
            self.sections[self.current_section] = " ".join(self.text_parts)

        full_text = " ".join(self.text_parts)
        # Clean up whitespace (but preserve code blocks)
        full_text = re.sub(r"[ \t]+", " ", full_text)
        full_text = re.sub(r"\n{3,}", "\n\n", full_text)
        return full_text.strip(), self.sections


def html_to_text(html: str) -> tuple[str, dict[str, str]]:
    """Convert wiki HTML to plain text with section extraction.

    Supports MediaWiki, TWiki, and static HTML page structures:
    - MediaWiki: bodyContent or mw-parser-output div
    - TWiki: patternTopic div (live TWiki rendered pages)
    - Fallback: full HTML with tag stripping

    Args:
        html: Raw HTML content from wiki page

    Returns:
        Tuple of (full_text, sections_dict)
    """
    content_html = html

    # Try TWiki structure first (class="patternTopic")
    twiki_match = re.search(
        r'class="patternTopic">(.*?)</div>\s*<!--\s*/patternTopic',
        html,
        re.DOTALL,
    )
    if twiki_match:
        content_html = twiki_match.group(1)
    else:
        # Try older MediaWiki structure first (id="bodyContent") - common on many wikis
        body_start = html.find('<div id="bodyContent">')
        if body_start >= 0:
            # End at printfooter (before footer/categories)
            footer_start = html.find('<div class="printfooter">', body_start)
            if footer_start > body_start:
                content_html = html[body_start:footer_start]
            else:
                # No printfooter, try to find end of bodyContent div
                # This is harder, so just take everything after bodyContent start
                content_html = html[body_start:]
        else:
            # Try modern MediaWiki structure (class="mw-parser-output")
            content_match = re.search(
                r'<div[^>]*class="[^"]*mw-parser-output[^"]*"[^>]*>(.*?)</div>\s*'
                r'(?:<div[^>]*class="[^"]*printfooter|$)',
                html,
                re.DOTALL | re.IGNORECASE,
            )
            if content_match:
                content_html = content_match.group(1)

    # Use simple tag stripping instead of complex parser
    # Remove script and style tags with their content
    content_html = re.sub(
        r"<script[^>]*>.*?</script>", " ", content_html, flags=re.DOTALL | re.IGNORECASE
    )
    content_html = re.sub(
        r"<style[^>]*>.*?</style>", " ", content_html, flags=re.DOTALL | re.IGNORECASE
    )

    # Remove HTML comments
    content_html = re.sub(r"<!--.*?-->", " ", content_html, flags=re.DOTALL)

    # Remove all HTML tags
    text = re.sub(r"<[^>]+>", " ", content_html)

    # Decode HTML entities
    import html as html_module

    text = html_module.unescape(text)

    # Clean up whitespace
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.strip()

    # Extract sections (simple heading detection)
    sections: dict[str, str] = {}
    lines = text.split("\n")
    current_section = ""
    current_text: list[str] = []

    for line in lines:
        # Detect headings (lines that look like section titles)
        if len(line) < 100 and line.strip() and not line.strip().startswith("("):
            # Could be a heading - save previous section
            if current_section and current_text:
                sections[current_section] = "\n".join(current_text).strip()
            current_section = line.strip()
            current_text = []
        else:
            current_text.append(line)

    # Save final section
    if current_section and current_text:
        sections[current_section] = "\n".join(current_text).strip()

    return text, sections


def twiki_markup_to_html(markup: str) -> str:
    """Convert raw TWiki markup to minimal HTML for the ingestion pipeline.

    Handles common TWiki markup patterns:
    - %META:...% lines (stripped)
    - ---+/---++/---+++ headings → <h1>/<h2>/<h3>
    - *bold*, _italic_ formatting
    - | table | cells | → <table> rows
    - Bullet lists (   * item) → <li>
    - <verbatim>...</verbatim> → <pre>
    - <literal>...</literal> → passed through (contains HTML)
    - %VARIABLE% → stripped (except %BR% → <br>)
    - Inline HTML tags preserved as-is

    Args:
        markup: Raw TWiki .txt file content

    Returns:
        Minimal HTML suitable for html_to_text() processing
    """
    lines = markup.split("\n")
    html_lines: list[str] = []
    title = ""
    in_verbatim = False
    in_literal = False

    for line in lines:
        stripped = line.strip()

        # Skip META lines
        if stripped.startswith("%META:"):
            # Extract title from first heading in TOPICPARENT or keep looking
            continue

        # Handle <verbatim> blocks → <pre>
        if "<verbatim>" in stripped and not in_verbatim:
            in_verbatim = True
            # Content may be on same line as tag
            content = stripped.split("<verbatim>", 1)[1]
            html_lines.append("<pre>")
            if content:
                html_lines.append(content)
            continue
        if "</verbatim>" in stripped and in_verbatim:
            content = stripped.split("</verbatim>", 1)[0]
            if content:
                html_lines.append(content)
            html_lines.append("</pre>")
            in_verbatim = False
            continue
        if in_verbatim:
            html_lines.append(line)
            continue

        # Handle <literal> blocks (pass through — they contain HTML)
        if "<literal>" in stripped:
            in_literal = True
            content = stripped.split("<literal>", 1)[1]
            if content:
                html_lines.append(content)
            continue
        if "</literal>" in stripped:
            content = stripped.split("</literal>", 1)[0]
            if content:
                html_lines.append(content)
            in_literal = False
            continue
        if in_literal:
            html_lines.append(line)
            continue

        # TWiki headings: ---+ to ---++++++ (h1 to h6)
        heading_match = re.match(r"^---(\++){1,6}\s*(.*)", stripped)
        if heading_match:
            level = len(heading_match.group(1))
            heading_text = heading_match.group(2).strip()
            # Strip TWiki formatting from heading text
            heading_text = _strip_twiki_formatting(heading_text)
            if not title and level <= 2:
                title = heading_text
            html_lines.append(f"<h{level}>{heading_text}</h{level}>")
            continue

        # TWiki tables: | cell | cell |
        if stripped.startswith("|") and stripped.endswith("|"):
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            row = "".join(f"<td>{_strip_twiki_formatting(c)}</td>" for c in cells)
            html_lines.append(f"<tr>{row}</tr>")
            continue

        # TWiki bullet lists: 3-space indent + *
        bullet_match = re.match(r"^(\s+)\*\s+(.*)", line)
        if bullet_match:
            content = _strip_twiki_formatting(bullet_match.group(2))
            html_lines.append(f"<li>{content}</li>")
            continue

        # TWiki numbered lists: 3-space indent + 1.
        num_match = re.match(r"^(\s+)\d+\.?\s+(.*)", line)
        if num_match:
            content = _strip_twiki_formatting(num_match.group(2))
            html_lines.append(f"<li>{content}</li>")
            continue

        # Regular text — apply TWiki variable and formatting cleanup
        processed = _strip_twiki_formatting(stripped)
        if processed:
            html_lines.append(f"<p>{processed}</p>")

    body = "\n".join(html_lines)
    if not title:
        title = "Untitled"

    return f"<html><head><title>{title}</title></head><body>{body}</body></html>"


def _strip_twiki_formatting(text: str) -> str:
    """Strip TWiki-specific formatting from text.

    Converts TWiki variables and formatting to plain text/HTML.
    """
    # %BR% → <br>
    text = text.replace("%BR%", "<br>")

    # Strip color formatting: %BLACK%, %RED%, %ENDCOLOR%, etc.
    text = re.sub(r"%[A-Z_]+%", "", text)

    # Strip %VARIABLE{...}% patterns (e.g., %USERSIG{...}%, %PARENTBC%)
    text = re.sub(r"%[A-Z_]+\{[^}]*\}%", "", text)

    # TWiki bold: *text* → <b>text</b> (but not ** or inside HTML)
    text = re.sub(r"(?<!\w)\*([^\*\n]+?)\*(?!\w)", r"<b>\1</b>", text)

    # TWiki italic: _text_ → <i>text</i> (but not __ or inside HTML)
    text = re.sub(r"(?<!\w)_([^_\n]+?)_(?!\w)", r"<i>\1</i>", text)

    # Strip remaining %PUBURLPATH% etc.
    text = re.sub(r"%[A-Z][A-Z0-9_]*%", "", text)

    # Clean up multiple spaces
    text = re.sub(r"  +", " ", text)

    return text.strip()


def get_embed_model() -> BaseEmbedding:
    """Get embedding model respecting embedding-backend config (local/remote).

    Delegates to the canonical cached singleton in imas_codex.embeddings.
    """
    from imas_codex.embeddings.llama_index import get_llama_embed_model

    return get_llama_embed_model()


class WikiIngestionPipeline:
    """Pipeline for ingesting wiki pages into the knowledge graph.

    Handles the full lifecycle: fetch → chunk → embed → link.
    Uses same embedding model as code examples for unified search.
    """

    def __init__(
        self,
        facility_id: str = "tcv",
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        use_rich: bool = True,
    ):
        """Initialize the pipeline.

        Args:
            facility_id: Facility ID (e.g., "tcv")
            chunk_size: Target chunk size in characters
            chunk_overlap: Overlap between chunks
            use_rich: Use Rich progress display
        """
        self.facility_id = facility_id
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.use_rich = use_rich

        # Initialize text splitter
        self.splitter = SentenceSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separator="\n",
        )

        # Embedding model is loaded lazily
        self._embed_model: BaseEmbedding | None = None

    @property
    def embed_model(self) -> BaseEmbedding:
        """Lazy-load embedding model respecting embedding-backend config."""
        if self._embed_model is None:
            self._embed_model = get_embed_model()
        return self._embed_model

    def _generate_page_id(self, page_name: str) -> str:
        """Generate unique ID for a WikiPage using canonical format."""
        from .scraper import canonical_page_id

        return canonical_page_id(page_name, self.facility_id)

    def _generate_chunk_id(self, page_id: str, chunk_idx: int) -> str:
        """Generate unique ID for a WikiChunk."""
        return f"{page_id}:chunk_{chunk_idx}"

    async def ingest_page(self, page: WikiPage) -> PageIngestionStats:
        """Ingest a single wiki page into the graph.

        Deterministic pipeline: fetch → extract → chunk → embed → persist.
        No LLM calls - all entity extraction uses regex patterns.

        IMPORTANT: Blocking I/O operations (embedding, Neo4j) are wrapped in
        asyncio.to_thread() to avoid blocking the event loop.

        Args:
            page: Scraped WikiPage instance

        Returns:
            Stats dict: {chunks, tree_nodes_linked, imas_paths_linked, conventions,
                        content_preview, mdsplus_paths}
        """
        import asyncio

        from .scraper import (
            extract_conventions,
            extract_mdsplus_paths,
            extract_units,
        )

        stats: PageIngestionStats = {
            "chunks": 0,
            "tree_nodes_linked": 0,
            "imas_paths_linked": 0,
            "conventions": 0,
            "units": 0,
            "content_preview": "",
            "mdsplus_paths": [],
        }

        page_id = self._generate_page_id(page.page_name)

        # Extract text content (CPU-bound, fast)
        text_content, sections = html_to_text(page.content_html)
        if not text_content or len(text_content) < 50:
            logger.warning("Skipping %s: insufficient content", page.page_name)
            return stats

        # Add content preview and MDSplus paths for progress display
        stats["content_preview"] = text_content[:300]
        stats["mdsplus_paths"] = page.mdsplus_paths[:10] if page.mdsplus_paths else []

        # Create LlamaIndex document for chunking
        doc = Document(
            text=text_content,
            metadata={
                "page_id": page_id,
                "title": page.title,
                "url": page.url,
                "facility_id": self.facility_id,
            },
        )

        # Split into chunks (CPU-bound, fast)
        nodes = self.splitter.get_nodes_from_documents([doc])
        stats["chunks"] = len(nodes)

        # Generate embeddings in batch - BLOCKING HTTP, run in thread pool
        chunk_texts = [node.text for node in nodes]  # type: ignore[attr-defined]
        embeddings = await asyncio.to_thread(
            self.embed_model.get_text_embedding_batch, chunk_texts
        )

        # Prepare batch data with pre-computed embeddings
        chunk_batch: list[dict] = []
        all_mdsplus: set[str] = set()

        for i, node in enumerate(nodes):
            chunk_text: str = node.text  # type: ignore[attr-defined]
            node.embedding = embeddings[i]

            # Extract entities from this chunk
            chunk_mdsplus = extract_mdsplus_paths(chunk_text)
            chunk_units = extract_units(chunk_text)
            chunk_conventions = extract_conventions(chunk_text)

            # Track for stats
            all_mdsplus.update(chunk_mdsplus)
            stats["units"] += len(chunk_units)
            stats["conventions"] += len(chunk_conventions)

            chunk_batch.append(
                {
                    "id": self._generate_chunk_id(page_id, i),
                    "wiki_page_id": page_id,
                    "facility_id": self.facility_id,
                    "chunk_index": i,
                    "content": chunk_text,
                    "embedding": node.embedding,
                    "mdsplus_paths": chunk_mdsplus,
                    "imas_paths": [],
                    "units": chunk_units,
                    "conventions": [c.get("name", "") for c in chunk_conventions],
                }
            )

        # Persist all chunks in batches
        with GraphClient() as gc:
            # Update WikiPage: MERGE on id to match discovery phase
            # Constraint is on (id, facility_id), so MERGE on id ensures we update
            # existing nodes rather than creating duplicates with title mismatches
            gc.query(
                """
                MERGE (p:WikiPage {id: $id})
                SET p.facility_id = $facility_id,
                    p.title = $title,
                    p.url = $url,
                    p.status = 'ingested',
                    p.content_hash = $hash,
                    p.last_scraped = datetime(),
                    p.ingested_at = datetime(),
                    p.claimed_at = null,
                    p.chunk_count = $chunk_count,
                    p.mdsplus_paths_found = $mdsplus_paths,
                    p.imas_paths_found = $imas_paths,
                    p.conventions_found = $conventions
                WITH p
                MATCH (f:Facility {id: $facility_id})
                MERGE (p)-[:FACILITY_ID]->(f)
                """,
                id=page_id,
                url=page.url,
                title=page.title,
                facility_id=self.facility_id,
                hash=page.content_hash,
                chunk_count=len(nodes),
                mdsplus_paths=page.mdsplus_paths,
                imas_paths=page.imas_paths,
                conventions=[c.get("name", "") for c in page.conventions],
            )

            # Batch persist chunks using UNWIND
            for i in range(0, len(chunk_batch), BATCH_SIZE):
                batch = chunk_batch[i : i + BATCH_SIZE]
                gc.query(
                    """
                    UNWIND $chunks AS chunk
                    MERGE (c:WikiChunk {id: chunk.id})
                    SET c.wiki_page_id = chunk.wiki_page_id,
                        c.facility_id = chunk.facility_id,
                        c.chunk_index = chunk.chunk_index,
                        c.content = chunk.content,
                        c.embedding = chunk.embedding,
                        c.mdsplus_paths_mentioned = chunk.mdsplus_paths,
                        c.imas_paths_mentioned = chunk.imas_paths,
                        c.units_mentioned = chunk.units,
                        c.conventions_mentioned = chunk.conventions
                    WITH c, chunk
                    MATCH (p:WikiPage {id: chunk.wiki_page_id})
                    MERGE (p)-[:HAS_CHUNK]->(c)
                    """,
                    chunks=batch,
                )

            # Create NEXT_CHUNK relationships for sequential navigation
            # Best practice from neo4j-graphrag: linked list pattern for chunk traversal
            if len(chunk_batch) > 1:
                gc.query(
                    """
                    MATCH (p:WikiPage {id: $page_id})-[:HAS_CHUNK]->(c:WikiChunk)
                    WITH c ORDER BY c.chunk_index
                    WITH collect(c) AS chunks
                    UNWIND range(0, size(chunks) - 2) AS idx
                    WITH chunks[idx] AS current, chunks[idx + 1] AS next
                    MERGE (current)-[:NEXT_CHUNK]->(next)
                    """,
                    page_id=page_id,
                )

            # Batch link chunks to TreeNodes
            if all_mdsplus:
                result = gc.query(
                    """
                    MATCH (c:WikiChunk {wiki_page_id: $page_id})
                    WHERE c.mdsplus_paths_mentioned IS NOT NULL
                    UNWIND c.mdsplus_paths_mentioned AS mds_path
                    MATCH (t:TreeNode)
                    WHERE t.path = mds_path
                       OR t.path ENDS WITH replace(mds_path, '\\\\', '')
                       OR t.canonical_path = toUpper(mds_path)
                    MERGE (c)-[:DOCUMENTS]->(t)
                    RETURN count(*) AS linked
                    """,
                    page_id=page_id,
                )
                if result and result[0]["linked"]:
                    stats["tree_nodes_linked"] = result[0]["linked"]

            # Create SignConvention nodes from page conventions (batch)
            if page.conventions:
                conv_batch = []
                for conv in page.conventions:
                    conv_id = f"{self.facility_id}:{conv.get('type', 'unknown')}:{hashlib.md5(conv.get('name', '').encode()).hexdigest()[:8]}"
                    conv_batch.append(
                        {
                            "id": conv_id,
                            "facility_id": self.facility_id,
                            "type": conv.get("type", "sign"),
                            "name": conv.get("name", ""),
                            "context": conv.get("context", ""),
                            "url": page.url,
                            "cocos_index": conv.get("cocos_index"),
                        }
                    )

                gc.query(
                    """
                    UNWIND $convs AS conv
                    MERGE (sc:SignConvention {id: conv.id})
                    SET sc.facility_id = conv.facility_id,
                        sc.convention_type = conv.type,
                        sc.name = conv.name,
                        sc.description = conv.context,
                        sc.wiki_source = conv.url,
                        sc.cocos_index = conv.cocos_index
                    WITH sc, conv
                    MATCH (f:Facility {id: conv.facility_id})
                    MERGE (sc)-[:FACILITY_ID]->(f)
                    """,
                    convs=conv_batch,
                )

            # Update link counts on WikiPage
            gc.query(
                """
                MATCH (p:WikiPage {id: $page_id})
                SET p.link_count = $links
                """,
                page_id=page_id,
                links=stats["tree_nodes_linked"] + stats["imas_paths_linked"],
            )

        return stats

    async def ingest_pages(
        self,
        page_names: list[str],
        progress_callback: ProgressCallback | None = None,
        rate_limit: float = 0.5,
    ) -> dict[str, int]:
        """Ingest multiple wiki pages by name (legacy method).

        For graph-driven workflow, use ingest_from_graph() instead.

        Args:
            page_names: List of wiki page names to ingest
            progress_callback: Optional callback for progress updates
            rate_limit: Minimum seconds between requests

        Returns:
            Aggregate stats dict
        """
        import asyncio

        total_stats = {
            "pages": 0,
            "pages_failed": 0,
            "chunks": 0,
            "tree_nodes_linked": 0,
            "imas_paths_linked": 0,
            "conventions": 0,
            "units": 0,
        }

        # Set up progress monitoring
        monitor = WikiProgressMonitor(use_rich=self.use_rich)
        set_current_monitor(monitor)
        monitor.start(total_pages=len(page_names))

        def report(current: int, total: int, message: str) -> None:
            if progress_callback:
                progress_callback(current, total, message)

        try:
            for i, page_name in enumerate(page_names):
                report(i, len(page_names), f"Processing {page_name}")
                page_id = self._generate_page_id(page_name)

                try:
                    # Determine site type from facility config
                    from imas_codex.discovery.wiki.config import WikiConfig

                    wiki_config = WikiConfig.from_facility(self.facility_id)
                    site_type = wiki_config.site_type

                    # Fetch the page
                    page = fetch_wiki_page(
                        page_name,
                        facility=self.facility_id,
                        site_type=site_type,
                        auth_type=wiki_config.auth_type,
                        credential_service=wiki_config.credential_service,
                    )

                    # Check if fetch succeeded
                    if page is None:
                        raise ValueError(f"Failed to fetch page {page_name}")

                    # Ingest it
                    page_stats = await self.ingest_page(page)

                    # Aggregate stats
                    total_stats["pages"] += 1
                    total_stats["chunks"] += page_stats["chunks"]
                    total_stats["tree_nodes_linked"] += page_stats["tree_nodes_linked"]
                    total_stats["imas_paths_linked"] += page_stats["imas_paths_linked"]
                    total_stats["conventions"] += page_stats["conventions"]
                    total_stats["units"] += page_stats["units"]

                    # Update progress monitor with content preview
                    monitor.update_scrape(
                        page_name,
                        chunks=page_stats["chunks"],
                        tree_nodes=page_stats["tree_nodes_linked"],
                        imas_paths=page_stats["imas_paths_linked"],
                        conventions=page_stats["conventions"],
                        units=page_stats["units"],
                        content_preview=str(page_stats.get("content_preview", "")),
                        mdsplus_paths=page_stats.get("mdsplus_paths"),
                    )

                except Exception as e:
                    logger.error("Failed to ingest %s: %s", page_name, e)
                    total_stats["pages_failed"] += 1
                    monitor.update_scrape(page_name, failed=True)
                    # Mark as failed in graph
                    mark_wiki_page_status(page_id, "failed", str(e))

                # Rate limiting
                await asyncio.sleep(rate_limit)

        finally:
            monitor.finish()
            set_current_monitor(None)

        report(len(page_names), len(page_names), "Ingestion complete")
        return total_stats

    async def ingest_from_graph(
        self,
        limit: int | None = None,
        min_interest_score: float = 0.5,
        progress_callback: ProgressCallback | None = None,
        rate_limit: float = 0.5,
    ) -> dict[str, int]:
        """Ingest wiki pages from the graph queue (graph-driven workflow).

        Reads WikiPage nodes with status='scored', fetches content,
        generates embeddings, and creates chunks with graph links.

        This is the preferred method - discover creates the queue,
        ingest processes it.

        Args:
            limit: Maximum pages to process (None for all)
            min_interest_score: Minimum interest score threshold
            progress_callback: Optional callback for progress updates
            rate_limit: Minimum seconds between requests

        Returns:
            Aggregate stats dict
        """
        # Get pending pages from graph
        pending = get_pending_wiki_pages(
            self.facility_id,
            limit=limit,
            min_interest_score=min_interest_score,
        )

        if not pending:
            logger.info("No pending wiki pages for %s", self.facility_id)
            return {
                "pages": 0,
                "pages_failed": 0,
                "chunks": 0,
                "tree_nodes_linked": 0,
                "imas_paths_linked": 0,
                "conventions": 0,
                "units": 0,
            }

        # Extract page identifiers from the graph results
        # For Confluence: use page ID (numeric)
        # For MediaWiki: use page title
        from imas_codex.discovery.wiki.config import WikiConfig

        wiki_config = WikiConfig.from_facility(self.facility_id)

        if wiki_config.site_type == "confluence":
            # Extract page ID from the full ID (format: "facility:page_id")
            page_names = [
                p["id"].split(":", 1)[1] if ":" in p["id"] else p["id"] for p in pending
            ]
        else:
            # Use title for MediaWiki
            page_names = [p["title"] for p in pending]

        logger.info(
            "Processing %d wiki pages from graph queue for %s",
            len(page_names),
            self.facility_id,
        )

        # Use the existing ingest_pages method
        return await self.ingest_pages(
            page_names,
            progress_callback=progress_callback,
            rate_limit=rate_limit,
        )


def ensure_wiki_vector_index() -> bool:
    """Ensure Neo4j vector index exists for WikiChunk embeddings.

    Creates the index if it doesn't exist. Safe to call multiple times.
    Dimension is determined by the configured embedding model.

    Returns:
        True if index was created, False if it already existed.
    """
    with GraphClient() as gc:
        # Check if index already exists
        existing = gc.query(
            "SHOW INDEXES YIELD name WHERE name = 'wiki_chunk_embedding' RETURN name"
        )
        if existing:
            logger.debug("wiki_chunk_embedding index already exists")
            return False

        # Get dimension from configured embedding model
        dim = get_embedding_dimension()
        gc.query(
            f"""
            CREATE VECTOR INDEX wiki_chunk_embedding IF NOT EXISTS
            FOR (c:WikiChunk) ON c.embedding
            OPTIONS {{
                indexConfig: {{
                    `vector.dimensions`: {dim},
                    `vector.similarity_function`: 'cosine'
                }}
            }}
            """
        )
        logger.info(f"Created wiki_chunk_embedding vector index ({dim} dims)")
        return True


# Backward compatibility alias
def create_wiki_vector_index() -> None:
    """Deprecated: Use ensure_wiki_vector_index() instead."""
    ensure_wiki_vector_index()


def get_wiki_stats(facility_id: str) -> dict:
    """Get wiki ingestion statistics for a facility.

    Args:
        facility_id: Facility ID

    Returns:
        Stats dict with page/chunk counts and link statistics
    """
    with GraphClient() as gc:
        result = gc.query(
            """
            MATCH (wp:WikiPage {facility_id: $facility_id})
            OPTIONAL MATCH (wp)-[:HAS_CHUNK]->(wc:WikiChunk)
            OPTIONAL MATCH (wc)-[:DOCUMENTS]->(t:TreeNode)
            OPTIONAL MATCH (wc)-[:MENTIONS_IMAS]->(ip:IMASPath)
            WITH wp, count(DISTINCT wc) AS chunks,
                 count(DISTINCT t) AS tree_nodes,
                 count(DISTINCT ip) AS imas_paths
            RETURN count(wp) AS pages,
                   sum(chunks) AS total_chunks,
                   sum(tree_nodes) AS tree_nodes_linked,
                   sum(imas_paths) AS imas_paths_linked
            """,
            facility_id=facility_id,
        )

        if result:
            return {
                "pages": result[0]["pages"],
                "chunks": result[0]["total_chunks"],
                "tree_nodes_linked": result[0]["tree_nodes_linked"],
                "imas_paths_linked": result[0]["imas_paths_linked"],
            }
        return {"pages": 0, "chunks": 0, "tree_nodes_linked": 0, "imas_paths_linked": 0}


def clear_facility_wiki(facility: str, batch_size: int = 1000) -> dict:
    """Delete all wiki discovery nodes for a facility in batches.

    Always cascades: deletes WikiChunks and WikiArtifacts along with WikiPages.
    Wiki chunks and artifacts are always dependent on their parent pages/facility,
    so cascade is the only sensible behaviour.

    Deletion order follows referential integrity:
    1. WikiChunks linked via HAS_CHUNK from WikiPages
    2. WikiChunks linked via HAS_CHUNK from WikiArtifacts
    3. Any remaining orphaned WikiChunks by facility_id
    4. WikiArtifacts by facility_id
    5. WikiPages by facility_id

    Args:
        facility: Facility ID
        batch_size: Nodes to delete per batch (default 1000)

    Returns:
        Dict with deletion counts: pages_deleted, chunks_deleted, artifacts_deleted
    """
    chunks_deleted = 0
    artifacts_deleted = 0
    pages_deleted = 0

    with GraphClient() as gc:
        # First delete WikiChunks linked to WikiPages
        while True:
            result = gc.query(
                """
                MATCH (wp:WikiPage {facility_id: $facility})
                      -[:HAS_CHUNK]->(wc:WikiChunk)
                WITH wc LIMIT $batch_size
                DETACH DELETE wc
                RETURN count(wc) AS deleted
                """,
                facility=facility,
                batch_size=batch_size,
            )
            deleted = result[0]["deleted"] if result else 0
            chunks_deleted += deleted
            if deleted < batch_size:
                break

        # Delete WikiChunks linked to WikiArtifacts
        while True:
            result = gc.query(
                """
                MATCH (wa:WikiArtifact {facility_id: $facility})
                      -[:HAS_CHUNK]->(wc:WikiChunk)
                WITH wc LIMIT $batch_size
                DETACH DELETE wc
                RETURN count(wc) AS deleted
                """,
                facility=facility,
                batch_size=batch_size,
            )
            deleted = result[0]["deleted"] if result else 0
            chunks_deleted += deleted
            if deleted < batch_size:
                break

        # Delete any remaining orphaned WikiChunks by facility_id
        while True:
            result = gc.query(
                """
                MATCH (wc:WikiChunk {facility_id: $facility})
                WITH wc LIMIT $batch_size
                DETACH DELETE wc
                RETURN count(wc) AS deleted
                """,
                facility=facility,
                batch_size=batch_size,
            )
            deleted = result[0]["deleted"] if result else 0
            chunks_deleted += deleted
            if deleted < batch_size:
                break

        # Delete WikiArtifacts (uses facility_id to catch both linked and orphaned)
        while True:
            result = gc.query(
                """
                MATCH (wa:WikiArtifact {facility_id: $facility})
                WITH wa LIMIT $batch_size
                DETACH DELETE wa
                RETURN count(wa) AS deleted
                """,
                facility=facility,
                batch_size=batch_size,
            )
            deleted = result[0]["deleted"] if result else 0
            artifacts_deleted += deleted
            if deleted < batch_size:
                break

        # Finally delete WikiPages
        while True:
            result = gc.query(
                """
                MATCH (wp:WikiPage {facility_id: $facility})
                WITH wp LIMIT $batch_size
                DETACH DELETE wp
                RETURN count(wp) AS deleted
                """,
                facility=facility,
                batch_size=batch_size,
            )
            deleted = result[0]["deleted"] if result else 0
            pages_deleted += deleted
            if deleted < batch_size:
                break

    return {
        "pages_deleted": pages_deleted,
        "chunks_deleted": chunks_deleted,
        "artifacts_deleted": artifacts_deleted,
    }


# =============================================================================
# Artifact Ingestion
# =============================================================================


class ArtifactIngestionStats(TypedDict):
    """Statistics from ingesting a single artifact."""

    chunks: int
    content_preview: str
    artifact_type: str


# Default size limits for artifact ingestion
DEFAULT_MAX_ARTIFACT_SIZE_MB = 5.0  # 5 MB default
MAX_ARTIFACT_SIZE_BYTES = int(DEFAULT_MAX_ARTIFACT_SIZE_MB * 1024 * 1024)


def _decode_url(url: str) -> str:
    """Decode URL-encoded path components.

    Wiki artifact URLs may have encoded slashes (%2F) that need decoding.
    """
    from urllib.parse import unquote

    return unquote(url)


def fetch_artifact_size(
    url: str,
    facility: str = "tcv",
    timeout: int = 30,
) -> int | None:
    """Fetch artifact file size via HTTP HEAD request.

    Uses the transfer module for SSH-proxied or direct HTTP access.

    Args:
        url: Full URL to artifact
        facility: SSH host alias (None for direct HTTP)
        timeout: Timeout in seconds

    Returns:
        File size in bytes, or None if size cannot be determined
    """
    import asyncio

    from imas_codex.discovery.base.transfer import TransferClient

    async def _get_size():
        async with TransferClient(ssh_host=facility) as client:
            return await client.get_size(url, timeout=timeout)

    # Run async in sync context
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            # Already in async context - create a task
            import concurrent.futures

            with concurrent.futures.ThreadPoolExecutor() as pool:
                future = pool.submit(asyncio.run, _get_size())
                return future.result(timeout=timeout + 5)
        else:
            return asyncio.run(_get_size())
    except Exception as e:
        logger.debug("Error getting artifact size: %s", e)
        return None


async def fetch_artifact_content(
    url: str,
    facility: str = "tcv",
    timeout: int = 120,
) -> tuple[str, bytes]:
    """Fetch artifact content with validation.

    Uses the transfer module for SSH-proxied or direct HTTP access.
    Validates content type matches expected file extension.

    Args:
        url: Full URL to artifact
        facility: SSH host alias (None for direct HTTP)
        timeout: Timeout in seconds

    Returns:
        Tuple of (content_type, raw_bytes)

    Raises:
        RuntimeError: If download fails
        ValueError: If content doesn't match expected type
    """
    from imas_codex.discovery.base.transfer import TransferClient

    # Get expected type from URL extension
    ext = url.rsplit(".", 1)[-1].lower()

    async with TransferClient(ssh_host=facility) as client:
        result = await client.download(url, timeout=timeout, expected_type=ext)

    if not result.success:
        raise RuntimeError(f"Failed to fetch {url}: {result.error}")

    return result.content_type or "application/octet-stream", result.content


class WikiArtifactPipeline:
    """Pipeline for ingesting wiki artifacts (PDFs, presentations, etc.).

    Uses LlamaIndex readers for PDF extraction. Other artifact types
    are currently marked as deferred (require OCR or specialized parsers).

    Supported types:
    - PDF: Full text extraction via pypdf (up to max_size_mb limit)
    - Others: Deferred for future implementation

    Size limits:
    - Artifacts larger than max_size_mb are marked as 'deferred' with size_bytes stored
    - This prevents flooding the graph with oversized content
    - Deferred artifacts can still be searched via metadata (filename, linked pages)
    """

    def __init__(
        self,
        facility_id: str = "tcv",
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        use_rich: bool = True,
        max_size_mb: float = DEFAULT_MAX_ARTIFACT_SIZE_MB,
    ):
        """Initialize the artifact pipeline.

        Args:
            facility_id: Facility ID
            chunk_size: Target chunk size in characters
            chunk_overlap: Overlap between chunks
            use_rich: Use Rich progress display
            max_size_mb: Maximum artifact size in MB (default 5.0)
        """
        self.facility_id = facility_id
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.use_rich = use_rich
        self.max_size_bytes = int(max_size_mb * 1024 * 1024)

        self.splitter = SentenceSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separator="\n",
        )

        self._embed_model: BaseEmbedding | None = None

    @property
    def embed_model(self) -> BaseEmbedding:
        """Lazy-load embedding model respecting embedding-backend config."""
        if self._embed_model is None:
            self._embed_model = get_embed_model()
        return self._embed_model

    async def ingest_pdf(
        self,
        artifact_id: str,
        pdf_bytes: bytes,
    ) -> ArtifactIngestionStats:
        """Ingest a PDF artifact.

        Args:
            artifact_id: WikiArtifact node ID
            pdf_bytes: Raw PDF content

        Returns:
            Ingestion stats

        Raises:
            ValueError: If content is not a valid PDF (wrong magic bytes)
        """
        import tempfile
        from pathlib import Path

        from llama_index.readers.file import PDFReader

        stats: ArtifactIngestionStats = {
            "chunks": 0,
            "content_preview": "",
            "artifact_type": "pdf",
        }

        # Validate PDF magic bytes before attempting to parse
        # Valid PDFs start with %PDF (possibly with leading whitespace/BOM)
        header = pdf_bytes[:1024] if len(pdf_bytes) > 1024 else pdf_bytes
        if b"%PDF" not in header:
            # Identify what the file actually is
            preview = pdf_bytes[:5].decode("latin-1", errors="replace")
            if pdf_bytes[:2] == b"PK":
                raise ValueError(
                    f"Content is a ZIP/Office file, not PDF (header: {preview!r})"
                )
            elif pdf_bytes[:5] == b"<!DOC" or pdf_bytes[:5] == b"<html":
                raise ValueError(f"Content is HTML, not PDF (header: {preview!r})")
            else:
                raise ValueError(f"Invalid PDF: missing %PDF header (got: {preview!r})")

        # Write to temp file for PDFReader
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(pdf_bytes)
            temp_path = Path(f.name)

        try:
            # Suppress pypdf's verbose internal warnings about corrupt PDF objects
            # These are benign warnings about non-critical PDF structure issues
            pypdf_logger = logging.getLogger("pypdf")
            original_level = pypdf_logger.level
            pypdf_logger.setLevel(logging.ERROR)

            try:
                reader = PDFReader()
                documents = reader.load_data(temp_path)
            finally:
                pypdf_logger.setLevel(original_level)

            if not documents:
                logger.warning("No content extracted from PDF: %s", artifact_id)
                return stats

            # Combine all pages
            full_text = "\n\n".join(doc.text for doc in documents if doc.text)

            if not full_text.strip():
                logger.warning("No text content in PDF: %s", artifact_id)
                return stats

            stats["content_preview"] = full_text[:300]
            stats = await self._create_artifact_chunks(
                artifact_id, full_text, "pdf", stats
            )
            return stats

        finally:
            temp_path.unlink(missing_ok=True)

    async def ingest_docx(
        self,
        artifact_id: str,
        content_bytes: bytes,
    ) -> ArtifactIngestionStats:
        """Ingest a Word document artifact.

        Args:
            artifact_id: WikiArtifact node ID
            content_bytes: Raw document content

        Returns:
            Ingestion stats
        """
        import tempfile
        from pathlib import Path

        from docx import Document as DocxDocument

        stats: ArtifactIngestionStats = {
            "chunks": 0,
            "content_preview": "",
            "artifact_type": "docx",
        }

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            f.write(content_bytes)
            temp_path = Path(f.name)

        try:
            doc = DocxDocument(temp_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)

            if not full_text.strip():
                logger.warning("No content extracted from DOCX: %s", artifact_id)
                return stats

            stats["content_preview"] = full_text[:300]
            stats = await self._create_artifact_chunks(
                artifact_id, full_text, "docx", stats
            )
            return stats
        finally:
            temp_path.unlink(missing_ok=True)

    async def ingest_pptx(
        self,
        artifact_id: str,
        content_bytes: bytes,
    ) -> ArtifactIngestionStats:
        """Ingest a PowerPoint artifact.

        Args:
            artifact_id: WikiArtifact node ID
            content_bytes: Raw presentation content

        Returns:
            Ingestion stats
        """
        import tempfile
        from pathlib import Path

        from pptx import Presentation

        stats: ArtifactIngestionStats = {
            "chunks": 0,
            "content_preview": "",
            "artifact_type": "pptx",
        }

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(content_bytes)
            temp_path = Path(f.name)

        try:
            prs = Presentation(temp_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

            full_text = "\n\n".join(text_parts)

            if not full_text.strip():
                logger.warning("No content extracted from PPTX: %s", artifact_id)
                return stats

            stats["content_preview"] = full_text[:300]
            stats = await self._create_artifact_chunks(
                artifact_id, full_text, "pptx", stats
            )
            return stats
        finally:
            temp_path.unlink(missing_ok=True)

    async def ingest_xlsx(
        self,
        artifact_id: str,
        content_bytes: bytes,
    ) -> ArtifactIngestionStats:
        """Ingest an Excel artifact.

        Args:
            artifact_id: WikiArtifact node ID
            content_bytes: Raw spreadsheet content

        Returns:
            Ingestion stats
        """
        stats: ArtifactIngestionStats = {
            "chunks": 0,
            "content_preview": "",
            "artifact_type": "xlsx",
        }

        from imas_codex.discovery.wiki.excel import extract_excel_full

        full_text = extract_excel_full(content_bytes)

        if not full_text.strip():
            logger.warning("No content extracted from XLSX: %s", artifact_id)
            return stats

        stats["content_preview"] = full_text[:300]
        stats = await self._create_artifact_chunks(
            artifact_id, full_text, "xlsx", stats
        )
        return stats

    async def ingest_notebook(
        self,
        artifact_id: str,
        content_bytes: bytes,
    ) -> ArtifactIngestionStats:
        """Ingest a Jupyter notebook artifact.

        Args:
            artifact_id: WikiArtifact node ID
            content_bytes: Raw notebook content

        Returns:
            Ingestion stats
        """
        import json

        import nbformat

        stats: ArtifactIngestionStats = {
            "chunks": 0,
            "content_preview": "",
            "artifact_type": "ipynb",
        }

        try:
            nb = nbformat.reads(content_bytes.decode("utf-8"), as_version=4)
        except (json.JSONDecodeError, nbformat.reader.NotJSONError) as e:
            logger.warning("Failed to parse notebook %s: %s", artifact_id, e)
            return stats

        text_parts = []
        for cell_num, cell in enumerate(nb.cells, 1):
            if cell.cell_type == "markdown":
                text_parts.append(f"[Cell {cell_num} - Markdown]\n{cell.source}")
            elif cell.cell_type == "code":
                text_parts.append(
                    f"[Cell {cell_num} - Code]\n```python\n{cell.source}\n```"
                )

        full_text = "\n\n".join(text_parts)

        if not full_text.strip():
            logger.warning("No content extracted from notebook: %s", artifact_id)
            return stats

        stats["content_preview"] = full_text[:300]
        stats = await self._create_artifact_chunks(
            artifact_id, full_text, "ipynb", stats
        )
        return stats

    async def ingest_artifact(
        self,
        artifact_id: str,
        content_bytes: bytes,
        artifact_type: str,
    ) -> ArtifactIngestionStats:
        """Dispatch to appropriate extraction method based on artifact type.

        Args:
            artifact_id: WikiArtifact node ID
            content_bytes: Raw artifact content
            artifact_type: Type of artifact (pdf, docx, pptx, xlsx, ipynb)

        Returns:
            Ingestion stats as dict

        Raises:
            ValueError: If artifact type is not supported
        """
        handlers = {
            "pdf": self.ingest_pdf,
            "docx": self.ingest_docx,
            "doc": self.ingest_docx,  # Try docx parser for .doc files
            "pptx": self.ingest_pptx,
            "ppt": self.ingest_pptx,  # Try pptx parser for .ppt files
            "xlsx": self.ingest_xlsx,
            "xls": self.ingest_xlsx,  # Try xlsx parser for .xls files
            "ipynb": self.ingest_notebook,
        }

        handler = handlers.get(artifact_type.lower())
        if handler is None:
            raise ValueError(f"Unsupported artifact type: {artifact_type}")

        return await handler(artifact_id, content_bytes)

    async def _create_artifact_chunks(
        self,
        artifact_id: str,
        full_text: str,
        artifact_type: str,
        stats: ArtifactIngestionStats,
    ) -> ArtifactIngestionStats:
        """Create chunks from extracted text and persist to graph.

        Common implementation used by all artifact type extractors.

        IMPORTANT: Blocking I/O operations (embedding) are wrapped in
        asyncio.to_thread() to avoid blocking the event loop.
        """
        import asyncio

        from .scraper import (
            extract_conventions,
            extract_mdsplus_paths,
            extract_units,
        )

        # Split into chunks
        combined_doc = Document(
            text=full_text,
            metadata={
                "artifact_id": artifact_id,
                "facility_id": self.facility_id,
            },
        )
        nodes = self.splitter.get_nodes_from_documents([combined_doc])
        stats["chunks"] = len(nodes)

        if not nodes:
            return stats

        # Generate embeddings in batch - BLOCKING HTTP, run in thread pool
        chunk_texts = [node.text for node in nodes]  # type: ignore[attr-defined]
        embeddings = await asyncio.to_thread(
            self.embed_model.get_text_embedding_batch, chunk_texts
        )

        # Prepare batch with pre-computed embeddings
        chunk_batch: list[dict] = []
        for i, node in enumerate(nodes):
            chunk_text: str = node.text  # type: ignore[attr-defined]
            node.embedding = embeddings[i]

            chunk_mdsplus = extract_mdsplus_paths(chunk_text)
            chunk_units = extract_units(chunk_text)
            chunk_conventions = extract_conventions(chunk_text)

            chunk_batch.append(
                {
                    "id": f"{artifact_id}:chunk_{i}",
                    "artifact_id": artifact_id,
                    "facility_id": self.facility_id,
                    "content": chunk_text,
                    "embedding": node.embedding,
                    "mdsplus_paths": chunk_mdsplus,
                    "imas_paths": [],
                    "units": chunk_units,
                    "conventions": [c.get("name", "") for c in chunk_conventions],
                }
            )

        # Persist chunks
        with GraphClient() as gc:
            # Update artifact status
            gc.query(
                """
                MATCH (wa:WikiArtifact {id: $id})
                SET wa.status = 'ingested',
                    wa.chunk_count = $chunks,
                    wa.ingested_at = datetime(),
                    wa.claimed_at = null
                """,
                id=artifact_id,
                chunks=len(nodes),
            )

            # Batch persist chunks
            for i in range(0, len(chunk_batch), BATCH_SIZE):
                batch = chunk_batch[i : i + BATCH_SIZE]
                gc.query(
                    """
                    UNWIND $chunks AS chunk
                    MERGE (c:WikiChunk {id: chunk.id})
                    SET c.artifact_id = chunk.artifact_id,
                        c.facility_id = chunk.facility_id,
                        c.content = chunk.content,
                        c.embedding = chunk.embedding,
                        c.mdsplus_paths_mentioned = chunk.mdsplus_paths,
                        c.imas_paths_mentioned = chunk.imas_paths,
                        c.units_mentioned = chunk.units,
                        c.conventions_mentioned = chunk.conventions
                    WITH c, chunk
                    MATCH (wa:WikiArtifact {id: chunk.artifact_id})
                    MERGE (wa)-[:HAS_CHUNK]->(c)
                    """,
                    chunks=batch,
                )

        return stats

    async def ingest_from_graph(
        self,
        limit: int | None = None,
        min_interest_score: float = 0.5,
    ) -> dict[str, int]:
        """Ingest artifacts from the graph queue.

        Pre-checks file size via HTTP HEAD before downloading. Artifacts
        exceeding max_size_bytes are marked as 'deferred' with size stored.
        Currently only processes PDFs. Other types are marked as deferred.

        Args:
            limit: Maximum artifacts to process (None for all)
            min_interest_score: Minimum score threshold

        Returns:
            Stats dict
        """
        total_stats = {
            "artifacts": 0,
            "artifacts_failed": 0,
            "artifacts_deferred": 0,
            "artifacts_oversized": 0,
            "chunks": 0,
        }

        pending = get_pending_wiki_artifacts(
            self.facility_id,
            limit=limit,
            min_interest_score=min_interest_score,
        )

        if not pending:
            logger.info("No pending artifacts for %s", self.facility_id)
            return total_stats

        for artifact in pending:
            artifact_id = artifact["id"]
            artifact_type = artifact.get("artifact_type", "unknown")
            url = artifact["url"]
            filename = artifact.get("filename", "unknown")

            try:
                # Check size before downloading
                size_bytes = fetch_artifact_size(url, facility=self.facility_id)

                if size_bytes is not None:
                    # Update size in graph
                    with GraphClient() as gc:
                        gc.query(
                            """
                            MATCH (wa:WikiArtifact {id: $id})
                            SET wa.size_bytes = $size
                            """,
                            id=artifact_id,
                            size=size_bytes,
                        )

                    # Check against limit
                    if size_bytes > self.max_size_bytes:
                        size_mb = size_bytes / (1024 * 1024)
                        max_mb = self.max_size_bytes / (1024 * 1024)
                        defer_reason = (
                            f"File size {size_mb:.1f} MB exceeds limit {max_mb:.1f} MB"
                        )
                        logger.info(
                            "Deferring oversized artifact %s: %s",
                            filename,
                            defer_reason,
                        )
                        with GraphClient() as gc:
                            gc.query(
                                """
                                MATCH (wa:WikiArtifact {id: $id})
                                SET wa.status = 'deferred',
                                    wa.defer_reason = $reason
                                """,
                                id=artifact_id,
                                reason=defer_reason,
                            )
                        total_stats["artifacts_deferred"] += 1
                        total_stats["artifacts_oversized"] += 1
                        continue

                if artifact_type == "pdf":
                    _, content = await fetch_artifact_content(
                        url, facility=self.facility_id
                    )
                    stats = await self.ingest_pdf(artifact_id, content)
                    total_stats["artifacts"] += 1
                    total_stats["chunks"] += stats["chunks"]
                else:
                    # Mark non-PDF as deferred
                    with GraphClient() as gc:
                        gc.query(
                            """
                            MATCH (wa:WikiArtifact {id: $id})
                            SET wa.status = 'deferred',
                                wa.defer_reason = $reason
                            """,
                            id=artifact_id,
                            reason=f"Artifact type '{artifact_type}' not yet supported",
                        )
                    total_stats["artifacts_deferred"] += 1

            except Exception as e:
                logger.error("Failed to ingest artifact %s: %s", artifact_id, e)
                with GraphClient() as gc:
                    gc.query(
                        """
                        MATCH (wa:WikiArtifact {id: $id})
                        SET wa.status = 'failed', wa.error = $error
                        """,
                        id=artifact_id,
                        error=str(e),
                    )
                total_stats["artifacts_failed"] += 1

        return total_stats


__all__ = [
    "ArtifactIngestionStats",
    "DEFAULT_MAX_ARTIFACT_SIZE_MB",
    "ProgressCallback",
    "WikiArtifactPipeline",
    "WikiIngestionPipeline",
    "clear_facility_wiki",
    "create_wiki_vector_index",
    "fetch_artifact_content",
    "fetch_artifact_size",
    "get_embed_model",
    "get_pending_wiki_artifacts",
    "get_pending_wiki_pages",
    "get_wiki_queue_stats",
    "get_wiki_stats",
    "html_to_text",
    "link_chunks_to_entities",
    "mark_wiki_page_status",
    "persist_chunks_batch",
]
